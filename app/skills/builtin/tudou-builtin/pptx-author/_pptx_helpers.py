"""pptx-author 共享 helpers — 用这个模块代替在每个脚本里复制 200 行样板。

## 典型用法（agent 脚本模板 —— 就这 5 行）

```python
from _pptx_helpers import *   # bash 工具自动注入 PYTHONPATH，直接 import

prs = new_deck()
def slide_cover(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide, THEME["bg"])
    add_text(slide, Inches(1), Inches(3), Inches(11), Inches(1.5),
             "2026 云服务市场洞察", size=40, bold=True, color=THEME["fg"])
slide_cover(prs)
prs.save("out.pptx")
verify_slides("out.pptx")     # 自动跑 shape-count 质量门
```

所有 helper 的第一个形参都叫 `slide`（不是 `s` / `sl`）——保持一致避免 NameError。
Inches/Pt/RGBColor/PP_ALIGN/MSO_ANCHOR/MSO_SHAPE 已通过 * 导出，**不要**再写
`from pptx.util import ...` —— 会 shadow 掉这里导出的版本。
"""
from __future__ import annotations

import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


__all__ = [
    # re-exported from python-pptx so the caller only needs one import line
    "Presentation", "Inches", "Pt", "Emu", "RGBColor",
    "MSO_SHAPE", "PP_ALIGN", "MSO_ANCHOR",
    "CategoryChartData", "XL_CHART_TYPE",
    # constants
    "THEME", "FONT", "SW", "SH", "EMU_PER_INCH",
    # primitives
    "hex_color", "new_deck", "blank_layout",
    # shape builders
    "set_bg", "add_text", "add_rect", "add_card",
    "add_styled_table", "add_bullets", "add_bar_chart",
    "add_line_chart", "add_image",
    # slide helpers
    "header_bar", "slide_full_chart",
    # quality gates
    "check_bounds", "check_safe_margins", "check_one_chart_per_slide",
    "verify_slides",
    # markdown helpers (optional, for md-to-deck scripts)
    "strip_md", "parse_md_outline",
]


EMU_PER_INCH = 914400

# 16:9 canvas
SW = Inches(13.333)
SH = Inches(7.5)

FONT = "Microsoft YaHei"   # CN/EN 通用。英文 deck 可临时覆盖为 "Inter" / "Arial"


def hex_color(s: str) -> RGBColor:
    """'#2563EB' or '2563EB' -> RGBColor."""
    s = s.lstrip("#")
    if len(s) != 6:
        raise ValueError(f"hex_color expects a 6-digit hex, got {s!r}")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


THEME = {
    "bg":        hex_color("#0F172A"),
    "fg":        hex_color("#F8FAFC"),
    "accent":    hex_color("#22D3EE"),
    "accent2":   hex_color("#F59E0B"),
    "muted":     hex_color("#94A3B8"),
    "card_bg":   hex_color("#1E293B"),
    "card_alt":  hex_color("#273449"),
    "row_alt":   hex_color("#273449"),
    "divider":   hex_color("#1E40AF"),
    "ok":        hex_color("#22C55E"),
    "warn":      hex_color("#F59E0B"),
    "bad":       hex_color("#EF4444"),
}


def new_deck() -> Presentation:
    """Create a blank 16:9 Presentation with SW/SH already set."""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank_layout(prs: Presentation):
    """Return the blank (layout index 6) — no placeholders."""
    return prs.slide_layouts[6]


# ─── shape builders ───────────────────────────────────────────────


def set_bg(slide, color: RGBColor) -> None:
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *,
             size: int = 18, bold: bool = False,
             color: RGBColor | None = None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    # Sanity-warn on suspicious font sizes. Typos like Pt(4) or Pt(150)
    # usually mean the LLM meant 14 / 15. Body text < 10pt is unreadable,
    # > 50pt almost never fits outside cover / big-number divider slides.
    try:
        if int(size) < 10 or int(size) > 50:
            print(
                f"[pptx_helpers WARNING] add_text size={size}pt is outside "
                f"normal range (10-50). text={str(text)[:40]!r}",
                file=sys.stderr,
            )
    except Exception:
        pass
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color if color is not None else THEME["fg"]
    return tb


def add_rect(slide, x, y, w, h, fill: RGBColor, *,
             rounded: bool = False,
             line_color: RGBColor | None = None,
             line_width_pt: float = 0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_width_pt:
            sh.line.width = Pt(line_width_pt)
    return sh


def add_card(slide, x, y, w, h, title: str, body: str, *,
             fill: RGBColor | None = None,
             accent: RGBColor | None = None,
             title_size: int = 16, body_size: int = 12):
    """Rounded card with a bold title + body text. Returns the card shape."""
    card_fill = fill if fill is not None else THEME["card_bg"]
    card = add_rect(slide, x, y, w, h, card_fill, rounded=True)
    if accent is not None:
        # thin accent strip on the left
        add_rect(slide, x, y, Inches(0.08), h, accent, rounded=True)
    # title
    pad = Inches(0.2)
    add_text(slide, x + pad, y + pad, w - 2 * pad, Inches(0.5),
             title, size=title_size, bold=True, color=THEME["fg"])
    # body
    add_text(slide, x + pad, y + Inches(0.7), w - 2 * pad, h - Inches(0.9),
             body, size=body_size, color=THEME["muted"])
    return card


def add_styled_table(slide, x, y, w, h, headers, rows):
    n_cols = max(len(headers), 1)
    n_rows = max(len(rows) + 1, 2)
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = shape.table
    # header
    for c in range(n_cols):
        cell = tbl.cell(0, c)
        cell.text = headers[c] if c < len(headers) else ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME["accent"]
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.name = FONT
                r.font.size = Pt(12)
                r.font.color.rgb = THEME["bg"]
    # data rows (alternating fills)
    for ri, row in enumerate(rows, start=1):
        for c in range(n_cols):
            cell = tbl.cell(ri, c)
            cell.text = row[c] if c < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                THEME["card_bg"] if ri % 2 == 0 else THEME["row_alt"])
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(10)
                    r.font.color.rgb = THEME["fg"]
    return shape


def add_bullets(slide, x, y, w, h, items, *,
                size: int = 14, color: RGBColor | None = None,
                bullet: str = "• "):
    """Bulleted list in a textbox. Items can be strings or (title, body) tuples."""
    # Readability guard — > 7 bullets on a single slide is a wall of text.
    # Not a hard cap (comparison pages sometimes need 8), just a nudge.
    if len(items) > 7:
        print(
            f"[pptx_helpers WARNING] add_bullets: {len(items)} items on one "
            f"slide — consider splitting (≤7 per page keeps it readable).",
            file=sys.stderr,
        )
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    clr = color if color is not None else THEME["fg"]
    for idx, item in enumerate(items):
        text = item if isinstance(item, str) else \
            f"{item[0]}: {item[1]}" if len(item) == 2 else str(item)
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = bullet + text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = clr
    return tb


def add_bar_chart(slide, x, y, w, h, categories, series_name, values):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, data)
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False
    return chart


def add_line_chart(slide, x, y, w, h, categories, series: dict):
    """series: {series_name: [v1, v2, ...]} — supports multi-series."""
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series.items():
        data.add_series(name, vals)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, w, h, data)
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    return chart


def add_image(slide, x, y, w, h, image_path: str):
    """Insert image. Missing-file → placeholder rect, never crashes."""
    if not os.path.isfile(image_path):
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        ph.fill.solid()
        ph.fill.fore_color.rgb = THEME["muted"]
        add_text(slide, x, y + h / 2 - Inches(0.2), w, Inches(0.4),
                 f"[image missing: {os.path.basename(image_path)}]",
                 size=12, color=THEME["bg"], align=PP_ALIGN.CENTER)
        return ph
    return slide.shapes.add_picture(image_path, x, y, w, h)


# ─── common slide helpers ─────────────────────────────────────────


def header_bar(prs_or_slide, title: str = ""):
    """Top title bar (card_bg strip + title text). Accepts a prs to create
    a new blank slide OR an existing slide to decorate. Returns the slide."""
    if hasattr(prs_or_slide, "slide_layouts"):
        # It's a Presentation — create a new blank slide.
        slide = prs_or_slide.slides.add_slide(blank_layout(prs_or_slide))
        set_bg(slide, THEME["bg"])
    else:
        slide = prs_or_slide
    add_rect(slide, 0, 0, SW, Inches(1.0), THEME["card_bg"])
    if title:
        add_text(slide, Inches(0.6), Inches(0.15),
                 Inches(12.1), Inches(0.7),
                 title, size=24, bold=True, color=THEME["fg"],
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def slide_full_chart(prs, title: str, chart_builder):
    """Standard layout D — title bar + one full-width chart below it.

    ``chart_builder(slide, x, y, w, h)`` should invoke one of
    ``add_bar_chart`` / ``add_line_chart`` / ``add_styled_table`` on the
    slide with the given box coords. Keeps the "one chart per slide"
    rule obvious and the callsite short:

        slide_full_chart(prs, "2026 市场份额",
                         lambda s, x, y, w, h: add_bar_chart(
                             s, x, y, w, h,
                             ["AWS", "Azure", "GCP"], "份额%", [32, 25, 11]))
    """
    slide = header_bar(prs, title)
    chart_builder(slide,
                  Inches(0.6), Inches(1.3),
                  Inches(12.13), Inches(5.9))
    return slide


# ─── quality gates ────────────────────────────────────────────────


def check_bounds(pptx_path: str, tol_inch: float = 0.02) -> list[str]:
    """Walk all shapes, report any that extend past slide edges. Empty list = pass."""
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    tol = int(tol_inch * EMU_PER_INCH)
    issues: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            x = shp.left or 0
            y = shp.top or 0
            w = shp.width or 0
            h = shp.height or 0
            name = getattr(shp, "name", "") or str(shp.shape_type)
            if x < -tol or y < -tol:
                issues.append(
                    f"slide {i}: '{name}' 左上角越界 "
                    f"({x / EMU_PER_INCH:.2f}, {y / EMU_PER_INCH:.2f})"
                )
            if x + w > sw + tol:
                issues.append(
                    f"slide {i}: '{name}' 右边越界: "
                    f"right={((x + w) / EMU_PER_INCH):.2f}\" > "
                    f"slide_width={(sw / EMU_PER_INCH):.2f}\""
                )
            if y + h > sh + tol:
                issues.append(
                    f"slide {i}: '{name}' 下边越界: "
                    f"bottom={((y + h) / EMU_PER_INCH):.2f}\" > "
                    f"slide_height={(sh / EMU_PER_INCH):.2f}\""
                )
    return issues


def check_safe_margins(pptx_path: str, min_margin_inch: float = 0.3) -> list[str]:
    """Warn when shapes sit inside the safe margins (default 0.3").

    Exempted (design, not content):
      - Full-width top/bottom decorative bars (x=0, w≈SW) — e.g. header bar
      - Full-height side strips (y=0, h≈SH) — e.g. cover accent strip
      - Any shape whose y-band is entirely inside the top/bottom 1" shelf
        (the header / footer decorative zone, where title text legitimately
        sits at y≈0.15" on top of the header bar).
    """
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    m = int(min_margin_inch * EMU_PER_INCH)
    shelf = EMU_PER_INCH  # 1" top-shelf / bottom-shelf decorative zone
    issues: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        for shp in slide.shapes:
            x = shp.left or 0
            y = shp.top or 0
            w = shp.width or 0
            h = shp.height or 0
            name = getattr(shp, "name", "") or str(shp.shape_type)
            # Decorative-edge skip: full-width top bar / full-height side strip
            is_full_width_bar = (x == 0 and w >= sw - EMU_PER_INCH // 4)
            is_full_height_strip = (y == 0 and h >= sh - EMU_PER_INCH // 4)
            if is_full_width_bar or is_full_height_strip:
                continue
            # Header / footer shelf skip: shape fully inside y∈[0, 1"]
            # or y∈[SH-1", SH] is an overlay on the decorative zone.
            if (y + h) <= shelf or y >= (sh - shelf):
                continue
            if x < m:
                issues.append(
                    f"slide {i}: '{name}' 左边距过近 "
                    f"(x={x / EMU_PER_INCH:.2f}\" < {min_margin_inch}\")"
                )
            if y < m:
                issues.append(
                    f"slide {i}: '{name}' 上边距过近 "
                    f"(y={y / EMU_PER_INCH:.2f}\" < {min_margin_inch}\")"
                )
            if x + w > sw - m:
                issues.append(
                    f"slide {i}: '{name}' 右边距过近 "
                    f"(right={(x + w) / EMU_PER_INCH:.2f}\", "
                    f"limit={(sw - m) / EMU_PER_INCH:.2f}\")"
                )
            if y + h > sh - m:
                issues.append(
                    f"slide {i}: '{name}' 下边距过近 "
                    f"(bottom={(y + h) / EMU_PER_INCH:.2f}\", "
                    f"limit={(sh - m) / EMU_PER_INCH:.2f}\")"
                )
    return issues


def check_one_chart_per_slide(pptx_path: str) -> list[str]:
    """Report slides with more than one chart. Two charts on one slide
    almost always means cramped / illegible — split into two slides.
    Tables don't count (they behave differently from charts visually).
    """
    prs = Presentation(pptx_path)
    issues: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        n = 0
        for shp in slide.shapes:
            # has_chart / has_table are safe attrs across python-pptx
            # shape types (BaseShape returns False for non-matching kinds).
            try:
                if getattr(shp, "has_chart", False):
                    n += 1
            except Exception:
                pass
        if n > 1:
            issues.append(f"slide {i}: {n} 个图表 — 一页最多 1 个")
    return issues


def _is_title_only_body(slide, sw_emu: int) -> bool:
    """True when a slide has nothing but a top header bar + its title text.

    Detects the "只有一行字" failure mode where an md-to-deck builder emitted
    a page for a H3 / H4 subsection but never populated its body. We flag
    these as FAIL (not just THIN) because users can see them and they
    obviously look broken. Cover slides (no full-width top bar) and the
    closing "谢谢" slide (no top bar) don't trip this.
    """
    shapes = list(slide.shapes)
    if len(shapes) == 0 or len(shapes) > 2:
        return False
    # Look for a rectangle at (x=0, y=0, w≈SW, h<1.5") — that's the header bar.
    ONE_INCH_EMU = 914400
    for sh in shapes:
        x = sh.left or 0
        y = sh.top or 0
        w = sh.width or 0
        h = sh.height or 0
        # tolerance 0.25" on width, max height 1.5"
        if (x == 0 and y == 0
                and abs(w - sw_emu) < (ONE_INCH_EMU // 4)
                and h < (ONE_INCH_EMU + ONE_INCH_EMU // 2)
                and h > 0):
            return True
    return False


def verify_slides(pptx_path: str, *, strict: bool = True) -> list[dict]:
    """Print per-slide shape counts + flags. Returns a list of dicts.

    Flag values:
      OK          — ≥3 shapes
      THIN        — 1-2 shapes that don't match the broken title-only pattern
                    (e.g. cover / closing slide) — printed but not a failure
      TITLE_ONLY  — only a top header bar + title text, no body content
                    (an md-to-deck subsection slide that forgot its body)
                    → FAIL
      BLANK       — 0 shapes → FAIL

    strict=True (default) → raises SystemExit(2) when any BLANK or
    TITLE_ONLY slide is found, which makes the bash step fail loudly so
    the agent retries that slide's function instead of delivering a
    deck where half the pages are just titles.
    """
    prs = Presentation(pptx_path)
    sw = prs.slide_width
    report = []
    fail = False
    print(f"—— {pptx_path} ——")
    for i, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        texts = [sh.text_frame.text[:40] for sh in shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()]
        n = len(shapes)
        if n == 0:
            flag = "BLANK"
            fail = True
        elif _is_title_only_body(slide, sw):
            flag = "TITLE_ONLY"
            fail = True
        elif n < 3:
            flag = "THIN"
        else:
            flag = "OK"
        print(f"  {i:2d}: {n:2d} shapes [{flag:10s}]  {texts[:2]}")
        report.append({"index": i, "shapes": n, "flag": flag, "texts": texts})
    # bounds check — hard fail (shapes off the canvas)
    bounds = check_bounds(pptx_path)
    if bounds:
        print("\n⚠️ bounds issues:")
        for b in bounds:
            print("  " + b)
    # safe-margin check — hard fail (shapes inside the 0.3" bleed zone)
    margins = check_safe_margins(pptx_path)
    if margins:
        print("\n⚠️ safe-margin issues:")
        for m in margins:
            print("  " + m)
    # one-chart-per-slide — hard fail (two charts in one slide = cramped)
    charts = check_one_chart_per_slide(pptx_path)
    if charts:
        print("\n⚠️ chart-density issues:")
        for c in charts:
            print("  " + c)
    if strict and (fail or bounds or margins or charts):
        raise SystemExit(2)
    return report


# ─── markdown helpers (optional, for md-to-deck builders) ──────────

_BOLD = re.compile(r"\*\*(.+?)\*\*")
_ITAL = re.compile(r"(?<!\*)\*(?!\s)([^*\n]+?)\*(?!\*)")
_CODE = re.compile(r"`([^`]+)`")
_LINK = re.compile(r"\[([^\]]+)\]\([^)]+\)")
_REFS = re.compile(r"\[\d+\]")


def strip_md(s: str) -> str:
    """Strip inline markdown markers: **bold** / *ital* / `code` / [link](url) / [1]."""
    s = _BOLD.sub(r"\1", s)
    s = _ITAL.sub(r"\1", s)
    s = _CODE.sub(r"\1", s)
    s = _LINK.sub(r"\1", s)
    s = _REFS.sub("", s)
    return s.strip()


def parse_md_outline(text: str) -> dict:
    """Light-weight md → outline. Returns {title, sections: [{title, bullets: [...]}]}.

    For quick deck generation from a markdown report. For the richer
    ### / #### / table parser, see the full v2 reference script in SKILL.md.
    """
    lines = text.split("\n")
    title = ""
    sections: list[dict] = []
    cur: dict | None = None
    for raw in lines:
        line = raw.rstrip()
        if not title and line.startswith("# "):
            title = strip_md(line[2:].strip())
            continue
        if line.startswith("## "):
            cur = {"title": strip_md(line[3:].strip()), "bullets": []}
            sections.append(cur)
            continue
        if cur is None:
            continue
        m = re.match(r"^\s*[-*]\s+(.+)$", line)
        if m:
            cur["bullets"].append(strip_md(m.group(1)))
    return {"title": title, "sections": sections}


# ─── quick self-test ──────────────────────────────────────────────

if __name__ == "__main__":
    out = Path("/tmp/pptx_helpers_selftest.pptx")
    prs = new_deck()
    # cover
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide, THEME["bg"])
    add_text(slide, Inches(1), Inches(3), Inches(11), Inches(1.5),
             "pptx_helpers self-test", size=40, bold=True)
    # cards
    slide = header_bar(prs, "Cards")
    add_card(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(2),
             "A", "one", accent=THEME["accent"])
    add_card(slide, Inches(4.8), Inches(1.5), Inches(4), Inches(2),
             "B", "two", accent=THEME["accent2"])
    add_card(slide, Inches(9.1), Inches(1.5), Inches(3.7), Inches(2),
             "C", "three", accent=THEME["ok"])
    # table
    slide = header_bar(prs, "Table")
    add_styled_table(slide, Inches(0.5), Inches(1.5),
                     Inches(12), Inches(4),
                     headers=["厂商", "份额", "同比"],
                     rows=[["AWS", "32%", "+3%"],
                           ["Azure", "25%", "+4%"],
                           ["GCP", "11%", "+2%"]])
    prs.save(str(out))
    print(f"✅ selftest saved → {out}")
    verify_slides(str(out))
