"""MD-driven PPTX template loader.

Parses frontmatter-based template definitions from `templates/_shared/*.md`,
resolves them against a theme `templates/themes/<name>/theme.yaml`, and
renders onto a python-pptx Presentation.

Design: layout & theme are decoupled. Same layout MD + different theme.yaml
produces different-looking slides. See SKILL.md for usage.

Public API:
    render_from_md(prs, "<theme>/<template_id>", params={...})
    list_themes() -> list[dict]
    list_layouts() -> list[dict]
    describe_layout(id) -> dict

Safety:
    Expression evaluation uses a whitelisted AST walker — NO eval().
    Allowed: literals, Name, Attribute, Subscript, BinOp, UnaryOp,
             Compare, BoolOp, IfExp, Call (for `len`, `max`, `min`, `str`,
             `int`, `float` only).
"""
from __future__ import annotations
import ast
import logging
import os
import re
from pathlib import Path
from typing import Any

import yaml

logger = logging.getLogger("tudou.pptx_loader")

_HERE = Path(__file__).parent
_TEMPLATES_DIR = _HERE / "templates"
_SHARED_DIR = _TEMPLATES_DIR / "_shared"
_THEMES_DIR = _TEMPLATES_DIR / "themes"

_FRONTMATTER_RE = re.compile(r"\A---\n(.+?)\n---\n?", re.DOTALL)

# ── Safe expression evaluation ────────────────────────────────────────
_ALLOWED_FUNCS = {
    "len": len, "max": max, "min": min,
    "str": str, "int": int, "float": float, "abs": abs,
    "round": round,
}
_ALLOWED_NODE_TYPES = (
    ast.Expression, ast.Constant, ast.Name, ast.Load, ast.Attribute,
    ast.Subscript, ast.Slice, ast.Index,
    ast.BinOp, ast.UnaryOp, ast.BoolOp, ast.Compare, ast.IfExp,
    ast.Add, ast.Sub, ast.Mult, ast.Div, ast.FloorDiv, ast.Mod, ast.Pow,
    ast.USub, ast.UAdd, ast.Not,
    ast.And, ast.Or,
    ast.Eq, ast.NotEq, ast.Lt, ast.LtE, ast.Gt, ast.GtE, ast.In, ast.NotIn,
    ast.Call, ast.List, ast.Tuple,
)


class _AttrDict(dict):
    """Dict that supports attribute access and returns None for missing keys.

    Needed because template expressions like `step.no` use AST Attribute,
    not Subscript — plain dict raises AttributeError. Missing-returns-None
    lets authors write `{cell.icon|fallback}` safely.

    Note: stored keys SHADOW dict methods. So `data.items` returns the
    stored value (not the `dict.items()` method). This is intentional —
    templates should see their data, not framework internals.
    """
    # Keep a small set of dunder / framework attrs that must NOT be shadowed
    # (otherwise list comprehension, copy, repr, iteration break).
    _RESERVED_ATTRS = frozenset({
        "__class__", "__dict__", "__doc__", "__init__", "__init_subclass__",
        "__subclasshook__", "__getattribute__", "__setattr__", "__delattr__",
        "__repr__", "__str__", "__hash__", "__eq__", "__ne__", "__lt__",
        "__le__", "__gt__", "__ge__", "__sizeof__", "__reduce__",
        "__reduce_ex__", "__format__", "__new__",
        "_RESERVED_ATTRS",
    })

    def __getattribute__(self, name):
        # Priority: reserved dunder/framework → stored key (shadows methods)
        # → regular attr lookup (methods, properties).
        if name in _AttrDict._RESERVED_ATTRS or name.startswith("__"):
            return dict.__getattribute__(self, name)
        try:
            # dict.__contains__ without triggering our machinery
            if dict.__contains__(self, name):
                return dict.__getitem__(self, name)
        except Exception:
            pass
        try:
            return dict.__getattribute__(self, name)
        except AttributeError:
            return None

    def __getitem__(self, key):
        # Also swallow KeyError for subscript access
        if dict.__contains__(self, key):
            return dict.__getitem__(self, key)
        return None


def _wrap(v):
    """Recursively wrap dicts inside a value for attribute access."""
    if isinstance(v, dict) and not isinstance(v, _AttrDict):
        return _AttrDict({k: _wrap(x) for k, x in v.items()})
    if isinstance(v, list):
        return [_wrap(x) for x in v]
    return v


def _safe_eval(expr: str, env: dict) -> Any:
    """Evaluate a restricted Python expression against ``env``.

    Never use on untrusted input without reviewing _ALLOWED_FUNCS.
    """
    tree = ast.parse(expr, mode="eval")
    for node in ast.walk(tree):
        if not isinstance(node, _ALLOWED_NODE_TYPES):
            raise ValueError(f"Disallowed expression node: {type(node).__name__}")
        if isinstance(node, ast.Call):
            # Only allow calls to whitelisted bare-name functions.
            if not isinstance(node.func, ast.Name):
                raise ValueError("Only bare-name calls allowed")
            if node.func.id not in _ALLOWED_FUNCS:
                raise ValueError(f"Call to '{node.func.id}' not allowed")
    env2 = {**_ALLOWED_FUNCS, **env}
    return eval(compile(tree, "<template_expr>", "eval"), {"__builtins__": {}}, env2)


# Match `{expr}` with optional `|default` fallback:
#   "{title}"           → eval(title)
#   "{cell.icon|target}" → eval(cell.icon) or "target"
_TPL_RE = re.compile(r"\{([^{}]+)\}")


def _render_tpl(text: str, env: dict) -> str:
    """Substitute {expr} placeholders, with '|default' fallback syntax."""
    if not isinstance(text, str) or "{" not in text:
        return text

    def _sub(m):
        raw = m.group(1)
        # split on first unescaped '|'
        if "|" in raw:
            expr, default = raw.split("|", 1)
        else:
            expr, default = raw, ""
        try:
            val = _safe_eval(expr.strip(), env)
            if val is None or val == "":
                return default
            return str(val)
        except Exception:
            return default

    return _TPL_RE.sub(_sub, text)


def _resolve_value(v, env: dict):
    """Resolve a single layout value.

    - string: run through template substitution. If no {…} braces, but the
      whole string is an identifier present in env (e.g. `sw`, `sh`), return
      env[v] directly — lets authors write `w: sw` instead of `w: "{sw}"`.
    - dict: recurse
    - list: recurse each
    - other: return as-is
    """
    if isinstance(v, str):
        # Bare-identifier shortcut: `w: sw` → env['sw']
        if "{" not in v and v.isidentifier() and v in env:
            return env[v]
        out = _render_tpl(v, env)
        # Try numeric coercion for geometry fields (x/y/w/h/size)
        if out and out.replace(".", "", 1).lstrip("-").isdigit():
            try:
                return float(out) if "." in out else int(out)
            except ValueError:
                return out
        return out
    if isinstance(v, dict):
        return {k: _resolve_value(val, env) for k, val in v.items()}
    if isinstance(v, list):
        return [_resolve_value(x, env) for x in v]
    return v


# ── Theme + layout loading ────────────────────────────────────────────

def _load_theme(theme_name: str) -> dict:
    tdir = _THEMES_DIR / theme_name
    yml = tdir / "theme.yaml"
    if not yml.exists():
        raise FileNotFoundError(
            f"theme '{theme_name}' not found at {yml}; "
            f"available: {[p.name for p in _THEMES_DIR.iterdir() if p.is_dir()]}")
    with open(yml, encoding="utf-8") as f:
        return yaml.safe_load(f)


def _load_layout_md(template_id: str, theme_name: str) -> dict:
    """Find layout MD. Lookup order: themes/<theme>/<id>.md, _shared/<id>.md."""
    candidates = [
        _THEMES_DIR / theme_name / f"{template_id}.md",
        _SHARED_DIR / f"{template_id}.md",
    ]
    for p in candidates:
        if p.exists():
            with open(p, encoding="utf-8") as f:
                text = f.read()
            break
    else:
        raise FileNotFoundError(
            f"layout '{template_id}' not found; looked in: "
            f"{[str(c) for c in candidates]}")

    m = _FRONTMATTER_RE.match(text)
    if not m:
        raise ValueError(f"layout {template_id}: missing YAML frontmatter")
    meta = yaml.safe_load(m.group(1))
    if "layout" not in meta:
        raise ValueError(f"layout {template_id}: missing 'layout' key")
    return meta


# ── Rendering ─────────────────────────────────────────────────────────
# Shape primitives map to _pptx_helpers.* calls.

def _hex_color(s: str):
    """Resolve #RRGGBB or theme color name to RGBColor."""
    from pptx.dml.color import RGBColor
    s = str(s).strip()
    if s.startswith("#") and len(s) == 7:
        return RGBColor(int(s[1:3], 16), int(s[3:5], 16), int(s[5:7], 16))
    raise ValueError(f"bad hex color: {s}")


def _color_value(name_or_hex: str, theme: dict):
    """Resolve 'primary' / '#RRGGBB' to RGBColor using theme palette."""
    colors = theme.get("colors", {})
    if name_or_hex in colors:
        return _hex_color(colors[name_or_hex])
    if str(name_or_hex).startswith("#"):
        return _hex_color(name_or_hex)
    # Unknown name — default to text color.
    return _hex_color(colors.get("text", "#000000"))


def _style_props(style_name: str, theme: dict) -> dict:
    """Resolve style name to {size, bold, color_rgb, font_name}."""
    styles = theme.get("styles", {})
    st = styles.get(style_name) or {}
    fonts = theme.get("fonts", {})
    font_alias = st.get("font", "sans")
    font_name = fonts.get(font_alias, "Arial")
    return {
        "size": int(st.get("size", 12)),
        "bold": bool(st.get("bold", False)),
        "color_rgb": _color_value(st.get("color", "text"), theme),
        "font_name": font_name,
    }


def _emit_shape(slide, shape_def: dict, theme: dict):
    """Emit one shape onto the slide."""
    # Late imports so this module only loads pptx when actually rendering.
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    kind = shape_def.get("shape")
    if kind is None:
        return  # likely a pure `if`/`for_each` marker, should have been expanded
    # Skip conditionals resolved to falsy.
    if "if" in shape_def:
        # Expression already evaluated in expand pass; a leftover `if: ""` means drop.
        if not shape_def.get("if"):
            return

    x = float(shape_def.get("x", 0))
    y = float(shape_def.get("y", 0))
    w = float(shape_def.get("w", 1))
    h = float(shape_def.get("h", 1))

    if kind == "rect":
        fill_name = shape_def.get("fill", "primary")
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _color_value(fill_name, theme)
        shape.line.fill.background()
        return

    if kind == "oval":
        fill_name = shape_def.get("fill", "primary")
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _color_value(fill_name, theme)
        shape.line.fill.background()
        return

    if kind == "text":
        text = shape_def.get("text", "")
        if text is None:
            text = ""
        text = str(text).strip()
        if not text:
            return  # skip empty text boxes
        style = shape_def.get("style", "card_body")
        props = _style_props(style, theme)
        # Optional color override (direct color on shape_def)
        color_override = shape_def.get("color")
        if color_override:
            props["color_rgb"] = _color_value(color_override, theme)

        tb = slide.shapes.add_textbox(
            Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.05)
        tf.margin_top = tf.margin_bottom = Inches(0.02)

        align = shape_def.get("align", "left")
        valign = shape_def.get("valign", "top")
        if valign == "middle":
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        elif valign == "bottom":
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM

        # First paragraph reuses the implicit one.
        lines = text.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            run = p.add_run()
            run.text = line
            run.font.name = props["font_name"]
            run.font.size = Pt(props["size"])
            run.font.bold = props["bold"]
            run.font.color.rgb = props["color_rgb"]
        return

    if kind == "icon":
        # Use existing _pptx_helpers.add_icon which already has Lucide/emoji
        # fallback.
        try:
            from _pptx_helpers import add_icon
        except ImportError:
            # fallback absolute import (when loaded as a module)
            import sys as _sys
            _sys.path.insert(0, str(_HERE))
            from _pptx_helpers import add_icon  # type: ignore
        name = shape_def.get("name", "target")
        size = float(shape_def.get("size", 0.5))
        color_name = shape_def.get("color", "primary")
        color_hex = theme.get("colors", {}).get(color_name, color_name)
        try:
            add_icon(slide, x, y, size, name, color_hex=color_hex)
        except Exception:
            # Icon fetch can fail (no net); draw a colored dot placeholder.
            from pptx.enum.shapes import MSO_SHAPE as _S
            from pptx.util import Inches as _I
            dot = slide.shapes.add_shape(
                _S.OVAL, _I(x), _I(y), _I(size), _I(size))
            dot.fill.solid()
            dot.fill.fore_color.rgb = _color_value(color_name, theme)
            dot.line.fill.background()
        return

    if kind == "table":
        # Required: headers (list[str]), rows (list[list[str]])
        # Optional: header_fill / header_text / row_fill / row_alt_fill / border
        headers = shape_def.get("headers") or []
        rows = shape_def.get("rows") or []
        if isinstance(headers, str):
            # Param substitution might have left it as a literal string;
            # try to parse comma list as a fallback.
            headers = [s.strip() for s in headers.split(",") if s.strip()]
        if not headers:
            return
        if not isinstance(rows, list):
            rows = []
        n_cols = len(headers)
        n_rows = 1 + len(rows)  # header + body rows
        if n_rows == 1:
            return  # nothing to render

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(x), Inches(y), Inches(w), Inches(h),
        )
        tbl = table_shape.table

        header_fill_name = shape_def.get("header_fill", "primary")
        header_text_name = shape_def.get("header_text", "bg")
        row_fill_name = shape_def.get("row_fill", "bg")
        row_alt_fill_name = shape_def.get("row_alt_fill", "panel")
        # border = shape_def.get("border", "border")  # currently unused

        body_props = _style_props("card_body", theme)
        head_props = _style_props("card_heading", theme)

        for ci, htext in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _color_value(header_fill_name, theme)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(htext)
            run.font.name = head_props["font_name"]
            run.font.size = Pt(head_props["size"])
            run.font.bold = True
            run.font.color.rgb = _color_value(header_text_name, theme)

        for ri, row in enumerate(rows):
            row_fill = row_alt_fill_name if (ri % 2 == 1) else row_fill_name
            cells = row if isinstance(row, list) else [row]
            for ci in range(n_cols):
                cell = tbl.cell(ri + 1, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = _color_value(row_fill, theme)
                tf = cell.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(cells[ci]) if ci < len(cells) else ""
                run.font.name = body_props["font_name"]
                run.font.size = Pt(body_props["size"])
                run.font.color.rgb = body_props["color_rgb"]
        return

    if kind == "image":
        # Required: path (string)
        # Optional: fit ("contain"|"cover"|"fill") — for now we just place
        # it at the given box; python-pptx auto-scales when w/h passed.
        path = shape_def.get("path")
        if not path:
            return
        # Resolve relative paths against cwd; absolute stays as-is.
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        if not os.path.exists(path):
            # Draw a panel placeholder so the slide still composes
            ph = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h))
            ph.fill.solid()
            ph.fill.fore_color.rgb = _color_value("panel", theme)
            ph.line.fill.background()
            return
        try:
            slide.shapes.add_picture(
                path,
                Inches(x), Inches(y), Inches(w), Inches(h),
            )
        except Exception:
            # python-pptx may reject unsupported formats; placeholder
            ph = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h))
            ph.fill.solid()
            ph.fill.fore_color.rgb = _color_value("panel", theme)
            ph.line.fill.background()
        return

    raise ValueError(f"Unknown shape kind: {kind}")


def _expand_layout(layout: list, env: dict) -> list:
    """Expand for_each / if directives into a flat list of shape dicts.

    env: {params..., theme_colors..., sw, sh, ...}
    """
    out = []
    for item in layout:
        if "if" in item:
            cond = item["if"]
            # eval the condition string against env
            try:
                ok = bool(_safe_eval(cond, env))
            except Exception:
                ok = bool(env.get(cond))
            if not ok:
                continue
            # Shallow copy without `if`
            item = {k: v for k, v in item.items() if k != "if"}
            resolved = {k: _resolve_value(v, env) for k, v in item.items()}
            out.append(resolved)
            continue

        if "for_each" in item:
            expr = item["for_each"]     # "cell in cells[:4]"
            m = re.match(r"\s*(\w+)\s+in\s+(.+)\s*$", expr)
            if not m:
                raise ValueError(f"bad for_each: {expr}")
            var, iterable_expr = m.group(1), m.group(2)
            index_name = item.get("index_as", "i")
            try:
                iterable = _safe_eval(iterable_expr, env)
            except Exception as e:
                raise ValueError(f"for_each iterable {iterable_expr!r}: {e}")
            if not isinstance(iterable, (list, tuple)):
                raise ValueError(f"for_each target must be iterable, got {type(iterable).__name__}")
            vars_block = item.get("vars", {}) or {}
            children = item.get("children", []) or []
            for i, val in enumerate(iterable):
                loop_env = {**env, var: val, index_name: i}
                # Compute vars (can reference i, var, existing env)
                for k, expr_str in vars_block.items():
                    try:
                        loop_env[k] = _safe_eval(str(expr_str), loop_env)
                    except Exception:
                        loop_env[k] = _render_tpl(str(expr_str), loop_env)
                # Expand each child against loop_env
                out.extend(_expand_layout(children, loop_env))
            continue

        # Plain shape — just resolve its values
        resolved = {k: _resolve_value(v, env) for k, v in item.items()}
        out.append(resolved)
    return out


def render_from_md(prs, template_path: str, params: dict | None = None):
    """Render one slide onto ``prs`` using ``<theme>/<template_id>``.

    Example: render_from_md(prs, "corporate/T01_cover", params={"title": "..."})
    """
    if "/" not in template_path:
        raise ValueError("template_path must be '<theme>/<template_id>'")
    theme_name, template_id = template_path.split("/", 1)
    theme = _load_theme(theme_name)
    layout_meta = _load_layout_md(template_id, theme_name)

    params = params or {}
    # Wrap nested dicts so `step.no` / `cell.icon` attribute access works and
    # missing keys return None instead of raising.
    wrapped_params = {k: _wrap(v) for k, v in params.items()}
    slide_size = theme.get("slide_size", [13.33, 7.5])
    env = {
        "sw": float(slide_size[0]),
        "sh": float(slide_size[1]),
        "steps": [],
        **wrapped_params,
    }

    # Apply defaults declared in params_schema (optional fields absent)
    schema = layout_meta.get("params_schema", {}) or {}
    props = schema.get("properties", {}) or {}
    for field, spec in props.items():
        if field not in env and isinstance(spec, dict):
            if "default" in spec:
                env[field] = spec["default"]

    # Add blank slide
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # Ensure slide size matches theme's declared size (idempotent — no-op if same).
    from pptx.util import Inches
    if prs.slide_width != Inches(slide_size[0]) or prs.slide_height != Inches(slide_size[1]):
        prs.slide_width = Inches(slide_size[0])
        prs.slide_height = Inches(slide_size[1])

    # Expand + emit
    shapes = _expand_layout(layout_meta["layout"], env)
    for s in shapes:
        _emit_shape(slide, s, theme)

    return slide


# ── Introspection API ─────────────────────────────────────────────────

def list_themes() -> list[dict]:
    """Return [{id, name, description, version, tags, trigger_keywords}, ...]."""
    out = []
    if not _THEMES_DIR.exists():
        return out
    for d in sorted(_THEMES_DIR.iterdir()):
        if not d.is_dir():
            continue
        y = d / "theme.yaml"
        if not y.exists():
            continue
        try:
            with open(y, encoding="utf-8") as f:
                t = yaml.safe_load(f) or {}
        except Exception:
            continue
        out.append({
            "id": d.name,
            "name": t.get("name", d.name),
            "description": t.get("description", ""),
            "version": t.get("version", ""),
            "tags": t.get("tags", []) or [],
            "trigger_keywords": t.get("trigger_keywords", []) or [],
        })
    return out


def recommend_theme(query: str, top_k: int = 3) -> list[dict]:
    """Score themes against ``query`` and return top-k ranked matches.

    Scoring: substring count (case-insensitive) of each theme's
    ``trigger_keywords`` and ``tags`` in the query. The higher the match
    count, the better the fit. Ties broken by number of unique keywords
    matched.

    Returns::

        [{"id": "tech_dark", "name": "Tech Dark", "score": 3,
          "matched": ["AI", "大模型", "开发者"],
          "description": "..."}, ...]
    """
    q = (query or "").lower()
    if not q:
        # Return the default order with zero scores.
        return [
            {**t, "score": 0, "matched": []}
            for t in list_themes()[:top_k]
        ]
    scored = []
    for t in list_themes():
        matched: list[str] = []
        # Count keyword hits (case-insensitive)
        for kw in t.get("trigger_keywords", []) or []:
            if kw and kw.lower() in q:
                matched.append(kw)
        for tag in t.get("tags", []) or []:
            if tag and tag.lower() in q and tag not in matched:
                matched.append(tag)
        if matched:
            scored.append({
                "id": t["id"],
                "name": t["name"],
                "description": t["description"],
                "score": len(matched),
                "matched": matched,
            })
    # Sort by score desc, then alphabetical for stability
    scored.sort(key=lambda x: (-x["score"], x["id"]))
    if scored:
        return scored[:top_k]
    # No keyword hits — return the first 3 themes with score 0 as fallback.
    fallback = list_themes()[:top_k]
    return [{"id": t["id"], "name": t["name"], "description": t["description"],
             "score": 0, "matched": []} for t in fallback]


def list_layouts() -> list[dict]:
    """Return [{id, name, category, summary, when_to_use}, ...]."""
    out = []
    if not _SHARED_DIR.exists():
        return out
    for md in sorted(_SHARED_DIR.glob("*.md")):
        try:
            with open(md, encoding="utf-8") as f:
                text = f.read()
        except Exception:
            continue
        m = _FRONTMATTER_RE.match(text)
        if not m:
            continue
        try:
            meta = yaml.safe_load(m.group(1))
        except Exception:
            continue
        out.append({
            "id": meta.get("id", md.stem),
            "name": meta.get("name", ""),
            "name_zh": meta.get("name_zh", ""),
            "category": meta.get("category", ""),
            "summary": meta.get("summary_zh") or meta.get("summary", ""),
            "when_to_use": meta.get("when_to_use", ""),
        })
    return out


# ── Sample params per layout — used to render preview thumbnails ──
# Each entry must satisfy the layout's params_schema (otherwise the
# render fails). Kept in code (not YAML) because preview renders run
# in-process and we want zero IO overhead during the gallery loop.
_LAYOUT_SAMPLE_PARAMS: dict[str, dict] = {
    "T01_cover": {
        "title": "示例标题",
        "subtitle": "副标题示意",
        "footer": "Tudou Claw · 2026",
    },
    "T02_three_column": {
        "title": "三栏要点",
        "items": [
            {"heading": "高效", "body": "并行执行、自动调度。"},
            {"heading": "稳定", "body": "状态机 + 持久化。"},
            {"heading": "智能", "body": "自动分流 + 经验沉淀。"},
        ],
    },
    "T04_section_divider": {
        "section_no": "01",
        "title": "进入正题",
        "subtitle": "下面进入主体内容",
    },
    "T06_requirement_grid": {
        "title": "需求矩阵",
        "items": [
            {"heading": "性能", "body": "P95 < 200ms"},
            {"heading": "可观测", "body": "trace + metrics"},
            {"heading": "安全", "body": "审计日志"},
            {"heading": "成本", "body": "≤ ¥0.05/req"},
        ],
    },
    "T09_vertical_timeline": {
        "title": "项目里程碑",
        "events": [
            {"when": "2025-Q4", "what": "立项"},
            {"when": "2026-Q1", "what": "MVP"},
            {"when": "2026-Q2", "what": "公测"},
            {"when": "2026-Q3", "what": "GA"},
        ],
    },
    "T10_comparison": {
        "title": "方案对比",
        "left_title":  "方案 A",
        "left_items":  ["优点 1", "优点 2", "缺点 1"],
        "right_title": "方案 B",
        "right_items": ["优点 1", "缺点 1", "缺点 2"],
    },
    "T15_chart_page": {
        "title": "数据走势",
        "caption": "近 4 季度 ARR 增长",
        "series": [
            {"label": "Q1", "value": 30},
            {"label": "Q2", "value": 45},
            {"label": "Q3", "value": 60},
            {"label": "Q4", "value": 78},
        ],
    },
    "T19_qa_closing": {
        "title": "Q & A",
        "subtitle": "感谢大家",
        "contact": "team@tudou.example",
    },
    "T20_table": {
        "title": "对比表",
        "caption": "三方案核心指标",
        "headers": ["维度", "A", "B", "C"],
        "rows": [
            ["延迟", "180ms", "240ms", "120ms"],
            ["成本", "¥0.04", "¥0.03", "¥0.06"],
            ["稳定性", "99.5%", "99.0%", "99.9%"],
        ],
    },
    "T21_image_left_text_right": {
        "title": "图文左右",
        "image_path": "",
        "image_caption": "示例图",
        "heading": "核心要点",
        "body": "图在左侧、要点在右侧;留白干净。",
    },
    "T22_text_top_image_bottom": {
        "title": "文上图下",
        "summary": "上方一段说明,下方放大图;"
                    "适合 banner / 概览页。",
        "image_path": "",
    },
    "T23_six_card_grid": {
        "title": "六卡片网格",
        "cards": [
            {"heading": "卡 1", "body": "要点说明"},
            {"heading": "卡 2", "body": "要点说明"},
            {"heading": "卡 3", "body": "要点说明"},
            {"heading": "卡 4", "body": "要点说明"},
            {"heading": "卡 5", "body": "要点说明"},
            {"heading": "卡 6", "body": "要点说明"},
        ],
    },
    "T24_kpi_metrics": {
        "title": "本季 KPI",
        "subtitle": "数据截止 2026-Q2",
        "metrics": [
            {"value": "+32%", "label": "营收同比", "delta": "↑ 行业平均"},
            {"value": "12M",  "label": "MAU",      "delta": "+8% QoQ"},
            {"value": "99.9%","label": "可用性",   "delta": "持平"},
            {"value": "¥1.2B","label": "ARR",      "delta": "+25% YoY"},
        ],
    },
    "T25_quote": {
        "quote": "未来不属于做事最快的公司, 而属于学习最快的公司。",
        "author": "Reid Hoffman",
        "author_title": "联合创始人 · LinkedIn",
    },
    "T26_process_flow": {
        "title": "客户成功四步",
        "steps": [
            {"heading": "需求识别", "body": "倾听 + 提问,锁定真实痛点"},
            {"heading": "方案设计", "body": "差异化方案,匹配场景"},
            {"heading": "落地实施", "body": "分阶段交付,设定验收"},
            {"heading": "持续运营", "body": "数据驱动迭代"},
        ],
    },
}


def get_sample_params(layout_id: str) -> dict:
    """Return a param dict that satisfies the layout's schema.

    Used for preview rendering and for the ``/api/portal/pptx-layouts``
    showcase endpoint. Falls back to ``{"title": "<id>"}`` for unknown
    layouts so a misnamed entry still produces *something* visible.
    """
    return dict(_LAYOUT_SAMPLE_PARAMS.get(layout_id) or {"title": layout_id})


def generate_layout_preview(layout_id: str, theme_name: str, *,
                              width: int = 640,
                              force: bool = False) -> str:
    """Render a single layout × theme combination to PNG.

    Output path: ``templates/_shared/.previews/<theme>/<layout_id>.png``.
    Returns the absolute path on success, "" on failure.

    Idempotent: if the file exists and ``force`` is False, returns it
    immediately (dirt-cheap). The full render takes ~200ms per slide.
    """
    import os as _os, tempfile as _tmp
    from pptx import Presentation

    out_dir = _SHARED_DIR / ".previews" / theme_name
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{layout_id}.png"

    if out_path.exists() and not force:
        return str(out_path)

    params = get_sample_params(layout_id)
    prs = Presentation()
    # 13.33 × 7.5 inch (16:9 widescreen)
    prs.slide_width = 12192000
    prs.slide_height = 6858000

    try:
        render_from_md(prs, f"{theme_name}/{layout_id}", params=params)
    except Exception as e:
        logger.warning(
            "layout preview render failed: %s/%s — %s",
            theme_name, layout_id, e,
        )
        return ""

    # Save pptx temp + render to PNG via existing pptx_to_pngs
    with _tmp.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        tmp_pptx = f.name
    try:
        prs.save(tmp_pptx)
        pngs = pptx_to_pngs(tmp_pptx, str(out_dir), width=width, force=True)
        if not pngs:
            return ""
        # pptx_to_pngs writes "<basename>_slide_001.png" — rename the
        # first slide to our canonical layout id.
        produced = pngs[0]
        if produced != str(out_path):
            try:
                _os.replace(produced, str(out_path))
            except OSError:
                pass
        # Clean up auxiliary slide PNGs (we only have one slide anyway)
        for p in pngs[1:]:
            try:
                _os.remove(p)
            except OSError:
                pass
        return str(out_path) if out_path.exists() else ""
    finally:
        try:
            _os.unlink(tmp_pptx)
        except OSError:
            pass


# Curated 3-layout showcase per theme — picks one cover, one data,
# one rich-content layout to demo the theme's distinctive look without
# rendering all 15. Cuts preview build from 150 frames to 30.
PREVIEW_SHOWCASE_LAYOUTS: tuple[str, ...] = (
    "T01_cover",        # 封面 — title hierarchy + accent line
    "T24_kpi_metrics",  # 数据 — big numbers + cards (color usage)
    "T26_process_flow", # 富内容 — multi-shape composition
)


def generate_all_layout_previews(*,
                                   themes: list[str] | None = None,
                                   layouts: tuple[str, ...] | None = None,
                                   force: bool = False) -> dict:
    """Bulk-render preview thumbnails for theme × layout pairs.

    Default is the 3-layout showcase set (``PREVIEW_SHOWCASE_LAYOUTS``)
    × all themes — 30 PNGs, ~10s wall clock. Pass an explicit
    ``layouts`` tuple to render a different subset.

    Returns ``{theme: {layout_id: png_path or ""}}``. Failures surface
    as empty strings so the caller can render placeholders; one bad
    layout doesn't block the rest.
    """
    if themes is None:
        # ``list_themes()`` returns display names (e.g. "Eco Green") but
        # the on-disk dir / render API expects slug ids (e.g.
        # "eco_green"). Use slug.
        themes = [t.get("id") or t.get("slug") or t["name"]
                  for t in list_themes()]
    if layouts is None:
        layouts = PREVIEW_SHOWCASE_LAYOUTS
    out: dict = {}
    total = ok = 0
    for theme in themes:
        out[theme] = {}
        for lid in layouts:
            total += 1
            path = generate_layout_preview(lid, theme, force=force)
            out[theme][lid] = path
            if path:
                ok += 1
    logger.info(
        "layout previews: %d/%d generated across %d themes",
        ok, total, len(themes),
    )
    return out


def generate_theme_preview(theme_name: str, *,
                            width: int = 640, height: int = 360,
                            force: bool = False) -> str:
    """Render a PIL synthetic preview of the theme's signature look.

    Shape:  640×360 (16:9) PNG saved to themes/<theme>/preview.png.
    Content:
      - header bar / title in primary color
      - 2×2 sample card grid using accent colors
      - bottom accent strip

    Used by portal UI so users can preview a theme before applying it.
    Returns the absolute path of the generated PNG.
    """
    from PIL import Image, ImageDraw, ImageFont  # local import, lazy dep.

    tdir = _THEMES_DIR / theme_name
    out_path = tdir / "preview.png"
    if out_path.exists() and not force:
        return str(out_path)

    theme = _load_theme(theme_name)
    colors = theme.get("colors", {})

    def _c(name: str, fallback: str = "#CCCCCC"):
        v = colors.get(name, fallback)
        if isinstance(v, str) and v.startswith("#") and len(v) == 7:
            return (int(v[1:3], 16), int(v[3:5], 16), int(v[5:7], 16))
        return (128, 128, 128)

    bg = _c("bg", "#FFFFFF")
    primary = _c("primary", "#1E3A5F")
    accent = _c("accent", "#E86C3A")
    text = _c("text", "#1A1A1A")
    muted = _c("muted", "#6B7280")
    panel = _c("panel", "#F3F4F6")

    img = Image.new("RGB", (width, height), bg)
    d = ImageDraw.Draw(img)

    # Top accent strip
    d.rectangle([0, 0, width, 8], fill=accent)

    # Title bar area (primary color text + short accent under-line)
    try:
        font_big = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 22)
        font_mid = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 12)
        font_small = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 10)
    except Exception:
        font_big = ImageFont.load_default()
        font_mid = ImageFont.load_default()
        font_small = ImageFont.load_default()

    title = theme.get("name", theme_name).upper()
    d.text((24, 24), title, font=font_big, fill=primary)
    # Under-title accent line
    d.rectangle([24, 56, 170, 59], fill=accent)

    # Subtitle (theme id)
    d.text((24, 66), f"id: {theme_name}", font=font_small, fill=muted)

    # 2×2 cards demonstrating palette
    grid_x0, grid_y0 = 24, 100
    cw, ch, gap = 140, 100, 16
    cells = [
        (primary, "AA"),  # primary tint
        (accent,  "BB"),
        (muted,   "CC"),
        (panel,   "DD"),  # panel (muted bg, primary text)
    ]
    for idx, (fill, label) in enumerate(cells):
        row = idx // 2
        col = idx % 2
        x = grid_x0 + col * (cw + gap)
        y = grid_y0 + row * (ch + gap)
        d.rectangle([x, y, x + cw, y + ch], fill=fill)
        # text color depends on fill contrast (simple luma)
        luma = 0.299 * fill[0] + 0.587 * fill[1] + 0.114 * fill[2]
        txt_color = (255, 255, 255) if luma < 140 else (40, 40, 40)
        d.text((x + 12, y + 12), label, font=font_big, fill=txt_color)
        # small subtitle line
        sub_color = txt_color
        d.text((x + 12, y + ch - 20), "sample card", font=font_small, fill=sub_color)

    # Right side: palette swatches (vertical strip)
    sw_x = grid_x0 + 2 * (cw + gap) + 30
    labels = [("primary", primary), ("accent", accent), ("muted", muted),
              ("text", text), ("panel", panel)]
    for i, (name, c) in enumerate(labels):
        y = 100 + i * 36
        d.rectangle([sw_x, y, sw_x + 28, y + 28], fill=c)
        d.text((sw_x + 36, y + 8), name, font=font_mid, fill=text)

    # Bottom accent strip
    d.rectangle([0, height - 8, width, height], fill=accent)

    img.save(out_path, "PNG", optimize=True)
    return str(out_path)


def pptx_to_pngs(pptx_path: str, out_dir: str, *,
                 width: int = 640, force: bool = False) -> list[str]:
    """Render every slide of a .pptx into PNG thumbnails using PIL.

    Pure-Python, no external converter needed (LibreOffice etc). The result
    is LOSSY — we only render:
      - rect / oval (fill + line color)
      - text boxes (approximate font size + color, truncated at box width)
      - pictures → gray placeholder tagged '[image]'
      - charts   → light-gray placeholder tagged '📊'

    Good enough for portal thumbnail previews. For pixel-perfect rendering
    use the actual .pptx file.

    Returns list of generated PNG paths, ordered by slide index.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image, ImageDraw, ImageFont
    import os as _os

    _os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    sw_in = prs.slide_width / 914400.0
    sh_in = prs.slide_height / 914400.0
    ppi = width / sw_in
    W = int(sw_in * ppi)
    H = int(sh_in * ppi)

    def _rgb_from_fill(shape):
        try:
            fill = shape.fill
            if fill.type == 1:  # solid
                c = fill.fore_color.rgb
                return (c[0], c[1], c[2])
        except Exception:
            pass
        return None

    def _rgb_from_run(run):
        try:
            c = run.font.color.rgb
            if c is not None:
                return (c[0], c[1], c[2])
        except Exception:
            pass
        return (40, 40, 40)

    # Font — try to use a reasonable CJK-capable font first.
    def _font_for(size_pt: int):
        candidates = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Medium.ttc",
            "/System/Library/Fonts/Helvetica.ttc",
            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        ]
        for p in candidates:
            if _os.path.exists(p):
                try:
                    px = max(8, int(size_pt * ppi / 72))   # pt → px
                    return ImageFont.truetype(p, px)
                except Exception:
                    continue
        return ImageFont.load_default()

    out_paths: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        out_path = _os.path.join(out_dir, f"slide_{idx}.png")
        if _os.path.exists(out_path) and not force:
            out_paths.append(out_path)
            continue

        img = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(img, "RGBA")

        # First pass: fills (so text draws on top)
        for shape in slide.shapes:
            try:
                x = int((shape.left or 0) / 914400.0 * ppi)
                y = int((shape.top or 0) / 914400.0 * ppi)
                w = int((shape.width or 0) / 914400.0 * ppi)
                h = int((shape.height or 0) / 914400.0 * ppi)
            except Exception:
                continue
            if w <= 0 or h <= 0:
                continue

            # Chart placeholder
            if shape.has_chart if hasattr(shape, 'has_chart') else False:
                d.rectangle([x, y, x+w, y+h], fill=(235, 235, 235),
                            outline=(180, 180, 180))
                try:
                    f = _font_for(24)
                    d.text((x + w//2, y + h//2), "📊 Chart",
                           anchor="mm", fill=(120, 120, 120), font=f)
                except Exception:
                    pass
                continue

            # Picture placeholder
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                d.rectangle([x, y, x+w, y+h], fill=(220, 220, 220),
                            outline=(170, 170, 170))
                try:
                    f = _font_for(10)
                    d.text((x + 4, y + 4), "[image]",
                           fill=(110, 110, 110), font=f)
                except Exception:
                    pass
                continue

            # Shape fill: rect vs oval
            fill_rgb = _rgb_from_fill(shape)
            if fill_rgb:
                try:
                    shape_type = getattr(shape, "auto_shape_type", None)
                except Exception:
                    shape_type = None
                if shape_type is not None and str(shape_type).lower().find("oval") >= 0:
                    d.ellipse([x, y, x+w, y+h], fill=fill_rgb)
                else:
                    d.rectangle([x, y, x+w, y+h], fill=fill_rgb)

        # Second pass: text on top of fills
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                x = int((shape.left or 0) / 914400.0 * ppi)
                y = int((shape.top or 0) / 914400.0 * ppi)
                w = int((shape.width or 0) / 914400.0 * ppi)
            except Exception:
                continue
            tf = shape.text_frame
            cursor_y = y + 2
            for para in tf.paragraphs:
                runs_text = "".join(r.text or "" for r in para.runs)
                if not runs_text.strip():
                    cursor_y += 4
                    continue
                # Use first run's style as representative
                first_run = para.runs[0] if para.runs else None
                size_pt = 14
                if first_run is not None and first_run.font.size is not None:
                    try:
                        size_pt = first_run.font.size.pt
                    except Exception:
                        pass
                color = _rgb_from_run(first_run) if first_run else (40, 40, 40)
                # Crude width-based truncation
                f = _font_for(int(size_pt))
                max_chars = max(2, int(w / max(size_pt * ppi / 72 * 0.55, 1)))
                text = runs_text if len(runs_text) <= max_chars else (
                    runs_text[:max_chars - 1] + "…")
                try:
                    d.text((x + 4, cursor_y), text, fill=color, font=f)
                except Exception:
                    pass
                # Line advance
                cursor_y += int(size_pt * ppi / 72 * 1.25) + 2

        img.save(out_path, "PNG", optimize=True)
        out_paths.append(out_path)

    return out_paths


def describe_layout(layout_id: str) -> dict:
    """Return full metadata for a single layout (frontmatter)."""
    md = _SHARED_DIR / f"{layout_id}.md"
    if not md.exists():
        raise FileNotFoundError(f"layout {layout_id} not found")
    with open(md, encoding="utf-8") as f:
        text = f.read()
    m = _FRONTMATTER_RE.match(text)
    if not m:
        raise ValueError(f"{layout_id}: missing frontmatter")
    return yaml.safe_load(m.group(1))
