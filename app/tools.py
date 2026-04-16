"""
Tool system — Claude Code style tools with JSON schema definitions.
"""
import fnmatch
import json
import logging
import os
import re
import subprocess
import threading
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable, Optional

from .defaults import (
    MAX_PARALLEL_WORKERS as _DEF_MAX_WORKERS,
    MAX_HTTP_RESPONSE_CHARS, MAX_JSON_RESULT_CHARS,
)

from . import sandbox as _sandbox
from . import knowledge as _knowledge

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# ToolRegistry pattern (singleton, inspired by Hermes Agent)
# ---------------------------------------------------------------------------

@dataclass
class ToolEntry:
    """Registry entry for a single tool."""
    name: str
    toolset: str  # e.g. "core", "web", "system", "coordination"
    schema: dict  # JSON schema definition (the function dict)
    handler: Callable  # The actual function to call
    check_fn: Optional[Callable] = None  # Optional availability check (returns bool)
    requires_env: list[str] = field(default_factory=list)  # Required environment variables
    is_async: bool = False  # Whether the tool is async
    description: str = ""  # Tool description
    risk_level: str = "safe"  # "safe", "moderate", or "dangerous"


class ToolRegistry:
    """Singleton registry for managing tools."""
    _instance: Optional["ToolRegistry"] = None
    _lock = threading.Lock()

    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._initialized = False
        return cls._instance

    def __init__(self):
        if self._initialized:
            return
        self._tools: dict[str, ToolEntry] = {}
        self._aliases: dict[str, str] = {}  # alias → canonical name
        self._initialized = True

    def register(
        self,
        name: str,
        toolset: str,
        schema: dict,
        handler: Callable,
        check_fn: Optional[Callable] = None,
        requires_env: Optional[list[str]] = None,
        is_async: bool = False,
        description: str = "",
        risk_level: str = "safe",
    ) -> None:
        """Register a new tool in the registry."""
        if name in self._tools:
            logger.warning(f"Tool '{name}' already registered, overwriting")

        entry = ToolEntry(
            name=name,
            toolset=toolset,
            schema=schema,
            handler=handler,
            check_fn=check_fn,
            requires_env=requires_env or [],
            is_async=is_async,
            description=description,
            risk_level=risk_level,
        )
        self._tools[name] = entry

    def unregister(self, name: str) -> bool:
        """Remove a tool from the registry. Returns True if removed, False if not found."""
        if name in self._tools:
            del self._tools[name]
            # Also remove any aliases pointing to this tool
            aliases_to_remove = [alias for alias, target in self._aliases.items() if target == name]
            for alias in aliases_to_remove:
                del self._aliases[alias]
            return True
        return False

    def add_alias(self, alias: str, canonical_name: str) -> None:
        """Add an alias for a tool."""
        if canonical_name not in self._tools:
            raise ValueError(f"Cannot alias '{alias}' to unknown tool '{canonical_name}'")
        self._aliases[alias] = canonical_name

    def dispatch(self, name: str, arguments: dict) -> str:
        """
        Dispatch a tool call by name.
        - Resolves aliases
        - Checks availability (check_fn)
        - Calls handler with arguments
        Returns a string result.
        """
        # Resolve alias
        canonical_name = self._aliases.get(name, name)

        if canonical_name not in self._tools:
            available = list(self._tools.keys())
            return (f"Error: Unknown tool '{name}'. "
                    f"Available: {available}. "
                    f"For shell commands use 'bash'.")

        entry = self._tools[canonical_name]

        # Check availability
        if entry.check_fn and not entry.check_fn():
            return f"Error: Tool '{canonical_name}' is not available in this context"

        # Check required environment variables
        missing_env = [var for var in entry.requires_env if var not in os.environ]
        if missing_env:
            return f"Error: Tool '{canonical_name}' requires environment variables: {missing_env}"

        # Call handler
        try:
            return entry.handler(**arguments)
        except TypeError as e:
            # Special handling for bash tool (argument name mismatch)
            if canonical_name == "bash" and arguments and "command" not in arguments:
                cmd = (arguments.get("cmd") or arguments.get("script") or
                       arguments.get("code") or next(iter(arguments.values()), ""))
                if isinstance(cmd, str) and cmd:
                    try:
                        return entry.handler(command=cmd)
                    except Exception as e2:
                        return f"Error executing tool '{canonical_name}': {e2}"
            return f"Error executing tool '{canonical_name}': {e}"
        except Exception as e:
            return f"Error executing tool '{canonical_name}': {e}"

    def get_definitions(self) -> list[dict]:
        """Return JSON schema definitions for all available tools.

        Returns tools in OpenAI function-calling format:
        {"type": "function", "function": {"name": ..., "description": ..., "parameters": ...}}
        """
        definitions = []
        for entry in self._tools.values():
            if entry.check_fn is None or entry.check_fn():
                schema = entry.schema
                # Ensure OpenAI function-calling wrapper is present
                if schema.get("type") == "function" and "function" in schema:
                    # Already wrapped correctly
                    definitions.append(schema)
                elif "name" in schema:
                    # Bare schema (name, description, parameters) — wrap it
                    definitions.append({
                        "type": "function",
                        "function": schema,
                    })
                else:
                    definitions.append(schema)
        return definitions

    def get_available_tools(self) -> list[str]:
        """Return list of tool names that pass their check_fn (or have no check_fn)."""
        return [
            name for name, entry in self._tools.items()
            if entry.check_fn is None or entry.check_fn()
        ]

    def is_parallel_safe(self, name: str) -> bool:
        """Check if a tool is safe for parallel execution."""
        canonical_name = self._aliases.get(name, name)
        return canonical_name in PARALLEL_SAFE_TOOLS

    def get_tool_entry(self, name: str) -> Optional[ToolEntry]:
        """Get the ToolEntry for a tool (resolving aliases)."""
        canonical_name = self._aliases.get(name, name)
        return self._tools.get(canonical_name)

    def list_tools(self) -> list[str]:
        """Return sorted list of all registered tool names."""
        return sorted(self._tools.keys())


def tool_result(result: Any, tool_name: str = "") -> str:
    """Standardized JSON tool result response."""
    if isinstance(result, str):
        return result
    return json.dumps({"status": "success", "result": result, "tool": tool_name})


def tool_error(message: str, tool_name: str = "", details: Optional[dict] = None) -> str:
    """Standardized JSON tool error response."""
    error_obj = {"status": "error", "message": message, "tool": tool_name}
    if details:
        error_obj["details"] = details
    return json.dumps(error_obj)


# ---------------------------------------------------------------------------
# Parallel execution configuration
# ---------------------------------------------------------------------------

# Tools that are safe to execute in parallel (read-only, no side effects)
PARALLEL_SAFE_TOOLS = frozenset({
    "read_file", "search_files", "glob_files",
    "web_search", "web_fetch", "web_screenshot",
    "datetime_calc", "json_process", "text_process",
    "get_skill_guide",
})

# Max parallel workers
MAX_PARALLEL_WORKERS = _DEF_MAX_WORKERS


# ---------------------------------------------------------------------------
# Tool definitions (JSON schema for function calling)
# ---------------------------------------------------------------------------

TOOL_DEFINITIONS: list[dict] = [
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": (
                "Read the contents of a file. Returns the file text. "
                "Supports optional offset (start line, 0-based) and limit (number of lines)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Absolute or relative file path"},
                    "offset": {
                        "type": "integer",
                        "description": "Start reading from this line number (0-based). Default 0.",
                    },
                    "limit": {
                        "type": "integer",
                        "description": "Maximum number of lines to read. Default: read all.",
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "Write content to a file. Creates parent directories if needed. Overwrites existing content.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "File path to write to"},
                    "content": {"type": "string", "description": "Content to write"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "edit_file",
            "description": (
                "Edit a file by performing an exact string replacement. "
                "Finds old_string in the file and replaces it with new_string. "
                "The old_string must appear exactly once for a unique match."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "File path to edit"},
                    "old_string": {"type": "string", "description": "Exact string to find"},
                    "new_string": {"type": "string", "description": "Replacement string"},
                },
                "required": ["path", "old_string", "new_string"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "bash",
            "description": "Execute a shell command (also aliased as 'exec', 'shell', 'run_command'). Returns stdout/stderr. Dangerous commands require human approval. Configurable timeout (default 30s).",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "Shell command to execute"},
                    "timeout": {
                        "type": "integer",
                        "description": "Timeout in seconds (default 30)",
                    },
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_files",
            "description": (
                "Search file contents using a regular expression (like grep -rn). "
                "Returns matching lines with file path and line number."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {"type": "string", "description": "Regular expression pattern to search for"},
                    "path": {
                        "type": "string",
                        "description": "Directory or file to search in (default: current directory)",
                    },
                    "include": {
                        "type": "string",
                        "description": "Glob pattern to filter files, e.g. '*.py'",
                    },
                },
                "required": ["pattern"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "glob_files",
            "description": "Find files matching a glob pattern. Returns a list of matching file paths.",
            "parameters": {
                "type": "object",
                "properties": {
                    "pattern": {
                        "type": "string",
                        "description": "Glob pattern, e.g. '**/*.py' or 'src/**/*.js'",
                    },
                    "path": {
                        "type": "string",
                        "description": "Base directory for the search (default: current directory)",
                    },
                },
                "required": ["pattern"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": (
                "Search the internet using DuckDuckGo. Returns search results with titles, "
                "URLs, and snippets. Use this to find up-to-date information from the web."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "The search query"},
                    "max_results": {
                        "type": "integer",
                        "description": "Maximum number of results to return (default: 8)",
                    },
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_fetch",
            "description": (
                "Fetch the text content of a web page URL. Returns the page content as plain text. "
                "Useful for reading articles, documentation, or any web page after finding it via web_search."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "The URL to fetch"},
                    "max_length": {
                        "type": "integer",
                        "description": "Maximum number of characters to return (default: 10000)",
                    },
                },
                "required": ["url"],
            },
        },
    },
    # ---- MCP bridge ----
    {
        "type": "function",
        "function": {
            "name": "mcp_call",
            "description": (
                "Invoke a tool on an external MCP server bound to this agent. "
                "Use this to call email (send_email), slack (send_message), github "
                "(create_pr), browser, postgres, or any other bound MCP. "
                "First call with list_mcps=true to see what MCPs are available and "
                "what tools each one provides."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "mcp_id": {
                        "type": "string",
                        "description": "The bound MCP id or name (e.g. 'email', 'slack', 'github')",
                    },
                    "tool": {
                        "type": "string",
                        "description": "The MCP tool name to invoke (e.g. 'send_email', 'send_message')",
                    },
                    "arguments": {
                        "type": "object",
                        "description": "Arguments object to pass to the MCP tool",
                    },
                    "list_mcps": {
                        "type": "boolean",
                        "description": "If true, list bound MCPs instead of calling one",
                    },
                },
            },
        },
    },
    # ---- Coordination tools (Claude Code architecture: TeamCreate / SendMessage / TaskList) ----
    {
        "type": "function",
        "function": {
            "name": "team_create",
            "description": (
                "Spawn a sub-agent for parallel task execution. The sub-agent runs "
                "independently with its own context window and tools. Use this when a "
                "task can be decomposed into independent sub-tasks that run in parallel "
                "(e.g., 3 sub-agents running in parallel = ~1 min vs 3 min serial)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "Name for the sub-agent"},
                    "role": {
                        "type": "string",
                        "description": "Role preset: coder, reviewer, researcher, tester, devops, writer",
                    },
                    "task": {"type": "string", "description": "Task description for the sub-agent to execute"},
                    "working_dir": {
                        "type": "string",
                        "description": "Working directory for the sub-agent (default: current dir)",
                    },
                },
                "required": ["name", "task"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_message",
            "description": (
                "Send a message to another agent for inter-agent communication. "
                "Use this to coordinate with other agents, share results, or request help. "
                "The target agent will receive the message and may respond."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "to_agent": {
                        "type": "string",
                        "description": "Agent ID or name to send the message to",
                    },
                    "content": {"type": "string", "description": "Message content"},
                    "msg_type": {
                        "type": "string",
                        "description": "Message type: task | info | result | question (default: task)",
                    },
                },
                "required": ["to_agent", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "task_update",
            "description": (
                "Manage the shared task list. Create one-time OR RECURRING tasks, update "
                "status, or complete tasks.\n"
                "• RECURRING: When the user asks for periodic work (\"每天9点\", \"每周一\", "
                "\"daily\", \"weekly\"), use action=create with recurrence + recurrence_spec.\n"
                "• DELAYED ONE-TIME: When the user asks \"5分钟后\", \"in 10 mins\", "
                "\"下午3点做X\", use action=create with run_at (e.g. run_at='+5m').\n"
                "The scheduler fires these tasks automatically at the configured time. "
                "Do NOT reply that you cannot run scheduled tasks."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": "Action: create | update | complete | list",
                    },
                    "task_id": {
                        "type": "string",
                        "description": "Task ID (required for update/complete)",
                    },
                    "title": {"type": "string", "description": "Task title (for create)"},
                    "description": {"type": "string", "description": "Task description"},
                    "status": {
                        "type": "string",
                        "description": "New status: todo | in_progress | done | blocked",
                    },
                    "result": {"type": "string", "description": "Result summary (for complete)"},
                    "recurrence": {
                        "type": "string",
                        "description": (
                            "Recurrence type: once (default, one-time) | daily | weekly | "
                            "monthly | cron. Use 'daily' for 每天, 'weekly' for 每周, "
                            "'monthly' for 每月."
                        ),
                    },
                    "recurrence_spec": {
                        "type": "string",
                        "description": (
                            "Schedule spec: daily='HH:MM' (e.g. '09:00'), "
                            "weekly='DOW HH:MM' (DOW=SUN|MON|TUE|WED|THU|FRI|SAT, e.g. 'MON 09:00'), "
                            "monthly='D HH:MM' (e.g. '1 09:00'), cron='m h dom mon dow'."
                        ),
                    },
                    "run_at": {
                        "type": "string",
                        "description": (
                            "For delayed one-time tasks: when to execute. "
                            "Accepts '+Nm' (N minutes from now, e.g. '+5m'), "
                            "'+Nh' (N hours from now, e.g. '+2h'), "
                            "or 'HH:MM' (today at specific time, e.g. '18:30'). "
                            "When set, the scheduler will auto-trigger this task at "
                            "the specified time. Use this for '5分钟后', 'in 10 mins', "
                            "'下午3点' etc."
                        ),
                    },
                },
                "required": ["action"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "plan_update",
            "description": (
                "Manage your execution plan — a visible step-by-step checklist "
                "that shows the user your progress. Use this at the START of a task "
                "to decompose it into steps, then update each step as you work through them. "
                "Actions: create_plan (set task_summary + steps array), "
                "start_step (mark a step as in_progress), "
                "complete_step (mark done with result_summary), "
                "add_step (add new step during execution), "
                "fail_step (mark as failed), "
                "replan (keep completed steps, replace pending with new steps)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": "Action: create_plan | start_step | complete_step | add_step | fail_step | replan",
                    },
                    "task_summary": {
                        "type": "string",
                        "description": "Brief summary of the task (for create_plan)",
                    },
                    "steps": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "detail": {"type": "string"},
                                "depends_on": {"type": "array", "items": {"type": "string"}, "description": "Step IDs this step depends on"},
                            },
                        },
                        "description": "List of step objects with title and optional detail (for create_plan)",
                    },
                    "step_id": {
                        "type": "string",
                        "description": "Step ID to update (for start_step/complete_step/fail_step)",
                    },
                    "title": {
                        "type": "string",
                        "description": "Step title (for add_step)",
                    },
                    "result_summary": {
                        "type": "string",
                        "description": "Brief result description (for complete_step/fail_step)",
                    },
                },
                "required": ["action"],
            },
        },
    },
    # ---- Screenshot tool ----
    {
        "type": "function",
        "function": {
            "name": "web_screenshot",
            "description": (
                "Take a screenshot of a web page. Returns the screenshot as a base64-encoded "
                "PNG image saved to a file. Requires Playwright or falls back to a simple "
                "HTML-to-image approach. Useful for capturing visual state of web pages, "
                "generating thumbnails, or documenting UI."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "The URL to screenshot"},
                    "output_path": {
                        "type": "string",
                        "description": "File path to save the screenshot (default: auto-generated in workspace)",
                    },
                    "full_page": {
                        "type": "boolean",
                        "description": "Capture the full scrollable page (default: false, viewport only)",
                    },
                    "width": {
                        "type": "integer",
                        "description": "Viewport width in pixels (default: 1280)",
                    },
                    "height": {
                        "type": "integer",
                        "description": "Viewport height in pixels (default: 720)",
                    },
                },
                "required": ["url"],
            },
        },
    },
    # ---- HTTP request tool ----
    {
        "type": "function",
        "function": {
            "name": "http_request",
            "description": (
                "Make an HTTP request (GET, POST, PUT, DELETE, PATCH) to any URL. "
                "Supports custom headers, JSON body, and form data. Useful for calling "
                "REST APIs, webhooks, or testing endpoints. Returns status code, headers, and body."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "The URL to request"},
                    "method": {
                        "type": "string",
                        "description": "HTTP method: GET, POST, PUT, DELETE, PATCH (default: GET)",
                    },
                    "headers": {
                        "type": "object",
                        "description": "Request headers as key-value pairs",
                    },
                    "body": {
                        "type": "string",
                        "description": "Request body (string or JSON string)",
                    },
                    "json_body": {
                        "type": "object",
                        "description": "Request body as JSON object (auto-sets Content-Type)",
                    },
                    "timeout": {
                        "type": "integer",
                        "description": "Request timeout in seconds (default: 30)",
                    },
                },
                "required": ["url"],
            },
        },
    },
    # ---- DateTime calculation tool ----
    {
        "type": "function",
        "function": {
            "name": "datetime_calc",
            "description": (
                "Perform date/time calculations. Get current time in any timezone, "
                "calculate date differences, add/subtract durations, format dates, "
                "convert between timezones. Use this instead of bash for date operations."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": (
                            "Action: 'now' (current time), 'diff' (difference between dates), "
                            "'add' (add duration to date), 'format' (reformat a date), "
                            "'convert' (convert timezone)"
                        ),
                    },
                    "date": {
                        "type": "string",
                        "description": "Date string (ISO format preferred, e.g. '2024-03-15T10:30:00')",
                    },
                    "date2": {
                        "type": "string",
                        "description": "Second date for 'diff' action",
                    },
                    "days": {"type": "integer", "description": "Days to add (for 'add' action)"},
                    "hours": {"type": "integer", "description": "Hours to add (for 'add' action)"},
                    "minutes": {"type": "integer", "description": "Minutes to add (for 'add' action)"},
                    "timezone": {
                        "type": "string",
                        "description": "Timezone name (e.g. 'Asia/Shanghai', 'US/Eastern', 'UTC')",
                    },
                    "format": {
                        "type": "string",
                        "description": "Output format string (Python strftime, e.g. '%%Y-%%m-%%d %%H:%%M')",
                    },
                },
                "required": ["action"],
            },
        },
    },
    # ---- JSON process tool ----
    {
        "type": "function",
        "function": {
            "name": "json_process",
            "description": (
                "Process JSON data: parse, format/pretty-print, extract fields using "
                "JSONPath-like expressions, transform, validate, or convert between "
                "JSON/CSV/YAML. Useful for data manipulation tasks."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": (
                            "Action: 'parse' (validate & pretty-print), 'extract' (extract field), "
                            "'keys' (list top-level keys), 'flatten' (flatten nested), "
                            "'to_csv' (JSON array to CSV), 'from_csv' (CSV to JSON), "
                            "'merge' (merge two JSON objects), 'count' (count items)"
                        ),
                    },
                    "data": {
                        "type": "string",
                        "description": "JSON string or file path to process",
                    },
                    "path": {
                        "type": "string",
                        "description": "JSONPath-like expression for 'extract' (e.g. 'users[0].name', 'data.items')",
                    },
                    "data2": {
                        "type": "string",
                        "description": "Second JSON string for 'merge' action",
                    },
                },
                "required": ["action", "data"],
            },
        },
    },
    # ---- Text process tool ----
    {
        "type": "function",
        "function": {
            "name": "text_process",
            "description": (
                "Process and transform text: count words/lines/chars, find & replace with regex, "
                "extract patterns, sort lines, deduplicate, base64 encode/decode, URL encode/decode, "
                "generate hash (md5/sha256), convert encoding. Batch text operations without bash."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "description": (
                            "Action: 'count' (word/line/char count), 'replace' (find & replace), "
                            "'extract' (extract regex matches), 'sort' (sort lines), "
                            "'dedup' (remove duplicates), 'base64_encode', 'base64_decode', "
                            "'url_encode', 'url_decode', 'hash' (md5/sha256), 'head' (first N lines), "
                            "'tail' (last N lines), 'split' (split by delimiter)"
                        ),
                    },
                    "text": {"type": "string", "description": "Input text to process"},
                    "pattern": {
                        "type": "string",
                        "description": "Regex pattern (for replace/extract)",
                    },
                    "replacement": {
                        "type": "string",
                        "description": "Replacement string (for replace)",
                    },
                    "n": {
                        "type": "integer",
                        "description": "Number of lines (for head/tail, default: 10)",
                    },
                    "algorithm": {
                        "type": "string",
                        "description": "Hash algorithm: md5, sha256, sha1 (for hash, default: sha256)",
                    },
                    "delimiter": {
                        "type": "string",
                        "description": "Delimiter (for split, default: newline)",
                    },
                },
                "required": ["action", "text"],
            },
        },
    },
    # ---- Experience persistence ----
    # NOTE: 经验条目(experience) 写入 experience_library 对应角色分桶。
    # 当经验积累到一定程度, agent 可通过 propose_skill 工具提议将经验
    # 锻造为技能(skill), 提交管理员审批后正式导入技能商店。
    {
        "type": "function",
        "function": {
            "name": "save_experience",
            "description": (
                "Persist a retrospective / active-learning finding as a reusable experience entry "
                "in the calling agent's role-based experience library. Use this for lessons learned, "
                "do/don't rules, and scene-specific action playbooks. "
                "After accumulating enough experiences on a topic, use propose_skill to "
                "crystallize them into a reusable skill package for admin approval."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "scene": {
                        "type": "string",
                        "description": "Trigger scenario / when this experience applies",
                    },
                    "core_knowledge": {
                        "type": "string",
                        "description": "Core insight / knowledge point",
                    },
                    "action_rules": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "1-3 positive action rules (do-this)",
                    },
                    "taboo_rules": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "1-2 taboo rules (avoid-this)",
                    },
                    "priority": {
                        "type": "string",
                        "enum": ["high", "medium", "low"],
                        "description": "Importance; default medium",
                    },
                    "tags": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional classification tags",
                    },
                    "exp_type": {
                        "type": "string",
                        "enum": ["retrospective", "active_learning"],
                        "description": "retrospective = 复盘产出; active_learning = 主动学习产出",
                    },
                    "source": {
                        "type": "string",
                        "description": "Human-readable origin (e.g. 'POC 贪吃蛇 产品复盘')",
                    },
                    "role": {
                        "type": "string",
                        "description": "Override the role bucket; defaults to the calling agent's role",
                    },
                },
                "required": ["scene", "core_knowledge"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "knowledge_lookup",
            "description": (
                "Look up shared knowledge base entries. Use this when you need reference information "
                "like design guidelines, tech stack standards, website lists, coding conventions, etc. "
                "Pass a search query to find relevant entries, or pass an entry_id to read a specific entry."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Search keyword or entry title",
                    },
                    "entry_id": {
                        "type": "string",
                        "description": "Specific entry ID to read (from a previous search result)",
                    },
                },
            },
        },
    },
    # ---- Cross-agent knowledge sharing ----
    {
        "type": "function",
        "function": {
            "name": "share_knowledge",
            "description": (
                "Share knowledge, best practices, or experience insights with all agents "
                "via the shared Knowledge Base. Use this when you've learned something valuable "
                "that other agents could benefit from — e.g., 'How to create professional PPTXs', "
                "'Best practices for API error handling', etc. All agents can read shared knowledge."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Concise title for the knowledge entry",
                    },
                    "content": {
                        "type": "string",
                        "description": "Detailed knowledge content — include steps, tips, examples, templates as needed",
                    },
                    "tags": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Tags for categorization, e.g. ['pptx', 'design', 'template']",
                    },
                },
                "required": ["title", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "learn_from_peers",
            "description": (
                "Learn from other agents' experiences. Browse and import high-quality experiences "
                "from other roles to expand your own capabilities. For example, a PM agent can learn "
                "design skills from a designer agent's experience pool."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "source_role": {
                        "type": "string",
                        "description": "The role to learn from, e.g. 'designer', 'coder', 'analyst'",
                    },
                    "topic": {
                        "type": "string",
                        "description": "Specific topic to search for, e.g. 'PPTX creation', 'API design'",
                    },
                    "limit": {
                        "type": "integer",
                        "description": "Max number of experiences to import (default 5)",
                    },
                },
                "required": ["source_role"],
            },
        },
    },
    # ---- Web login request (human-in-the-loop) ----
    {
        "type": "function",
        "function": {
            "name": "request_web_login",
            "description": (
                "Explicitly request the user to log into a website. "
                "Shows an interactive login card (iframe, credential form, cookie/token). "
                "Note: login walls encountered during browser navigation are handled automatically — "
                "you only need this tool to proactively request login before navigating."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {
                        "type": "string",
                        "description": "The URL that requires login",
                    },
                    "site_name": {
                        "type": "string",
                        "description": "Human-readable site name, e.g. 'GitHub', 'Jira', '企业微信'",
                    },
                    "reason": {
                        "type": "string",
                        "description": "Why you need access — what task requires this login",
                    },
                    "login_url": {
                        "type": "string",
                        "description": "Optional: the specific login page URL if different from the target URL",
                    },
                },
                "required": ["url", "site_name", "reason"],
            },
        },
    },
    # ---- Package management tool ----
    {
        "type": "function",
        "function": {
            "name": "pip_install",
            "description": "Install or upgrade Python packages using pip. Supports space-separated package names.",
            "parameters": {
                "type": "object",
                "properties": {
                    "packages": {
                        "type": "string",
                        "description": "Space-separated package names to install (e.g., 'requests numpy pandas')",
                    },
                    "upgrade": {
                        "type": "boolean",
                        "description": "Whether to upgrade packages to the latest version (default: false)",
                    },
                },
                "required": ["packages"],
            },
        },
    },
    # ---- PowerPoint creation tool ----
    {
        "type": "function",
        "function": {
            "name": "create_pptx",
            "description": "Create a PowerPoint presentation (.pptx) file with custom slides. Auto-installs python-pptx if needed.",
            "parameters": {
                "type": "object",
                "properties": {
                    "output_path": {
                        "type": "string",
                        "description": "Path where the .pptx file will be saved",
                    },
                    "title": {
                        "type": "string",
                        "description": "Optional title for the presentation deck",
                    },
                    "slides": {
                        "type": "array",
                        "description": "Array of slide objects, each with title, content, optional layout, and optional images",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {
                                    "type": "string",
                                    "description": "Slide title",
                                },
                                "content": {
                                    "type": "string",
                                    "description": "Slide content (bullet text or paragraphs)",
                                },
                                "layout": {
                                    "type": "string",
                                    "description": "Layout type: 'title', 'content', 'title_content', 'blank' (default: 'title_content')",
                                },
                                "images": {
                                    "type": "array",
                                    "description": "Optional list of images to place on the slide",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "path": {"type": "string", "description": "Image file path"},
                                            "left": {"type": "number", "description": "Left position in inches (default 1)"},
                                            "top": {"type": "number", "description": "Top position in inches (default 2)"},
                                            "width": {"type": "number", "description": "Width in inches (0=auto)"},
                                            "height": {"type": "number", "description": "Height in inches (0=auto)"},
                                        },
                                        "required": ["path"],
                                    },
                                },
                            },
                            "required": ["title", "content"],
                        },
                    },
                },
                "required": ["output_path", "slides"],
            },
        },
    },
    # ---- Advanced PPTX tool ----
    {
        "type": "function",
        "function": {
            "name": "create_pptx_advanced",
            "description": "创建高级精美PowerPoint演示文稿。支持形状、图表、表格、多栏布局、信息图表等专业元素。每页通过elements数组精确控制所有元素的位置和样式。",
            "parameters": {
                "type": "object",
                "properties": {
                    "output_path": {
                        "type": "string",
                        "description": "输出 .pptx 文件路径",
                    },
                    "theme": {
                        "type": "object",
                        "description": "全局配色主题",
                        "properties": {
                            "primary": {"type": "string", "description": "主色 hex (如 'E8590C')"},
                            "secondary": {"type": "string", "description": "辅色 hex (如 '2B2B2B')"},
                            "accent": {"type": "string", "description": "强调色 hex (如 'F4A261')"},
                            "background": {"type": "string", "description": "默认背景色 hex (如 'FFFFFF')"},
                            "title_font": {"type": "string", "description": "标题字体 (如 'Microsoft YaHei')"},
                            "body_font": {"type": "string", "description": "正文字体 (如 'Microsoft YaHei')"},
                        },
                    },
                    "slides": {
                        "type": "array",
                        "description": "页面数组。推荐用layout自动排版，也可用elements手动控制，或两者结合。",
                        "items": {
                            "type": "object",
                            "properties": {
                                "layout": {
                                    "type": "object",
                                    "description": "智能布局（推荐）。设置type和items，工具自动计算坐标。type: cover|toc|section|cards|process|kpi|comparison|timeline|chart|table|closing。示例: {\"type\":\"process\",\"title\":\"流程\",\"page_num\":3,\"items\":[{\"title\":\"步骤1\",\"detail\":\"说明\"}]}",
                                    "properties": {
                                        "type": {"type": "string", "description": "布局类型: cover|toc|section|cards|process|kpi|comparison|timeline|chart|table|closing"},
                                        "title": {"type": "string", "description": "页面标题"},
                                        "page_num": {"type": "integer", "description": "页码编号"},
                                        "items": {"type": "array", "description": "内容项数组，结构因布局类型而异"},
                                        "subtitle": {"type": "string", "description": "[cover/closing] 副标题"},
                                        "date": {"type": "string", "description": "[cover] 日期"},
                                        "author": {"type": "string", "description": "[cover] 作者"},
                                        "left": {"type": "object", "description": "[comparison] 左侧 {title, items:[]}"},
                                        "right": {"type": "object", "description": "[comparison] 右侧 {title, items:[]}"},
                                        "headers": {"type": "array", "description": "[table] 表头"},
                                        "rows": {"type": "array", "description": "[table] 数据行"},
                                        "summary": {"type": "string", "description": "底部说明文字"},
                                    },
                                },
                                "background": {
                                    "type": "string",
                                    "description": "页面背景色 hex，覆盖主题默认值",
                                },
                                "elements": {
                                    "type": "array",
                                    "description": "手动元素数组（可与layout组合使用，手动元素追加在layout自动元素之后）。每个元素须有type和x,y,w,h(英寸)。",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "type": {"type": "string", "description": "元素类型: text|shape|chart|table|image|icon_circle|line"},
                                            "x": {"type": "number", "description": "左边距(英寸)"},
                                            "y": {"type": "number", "description": "上边距(英寸)"},
                                            "w": {"type": "number", "description": "宽度(英寸)"},
                                            "h": {"type": "number", "description": "高度(英寸)"},
                                            "content": {"type": "string", "description": "[text] 文本内容，支持\\n换行"},
                                            "font_size": {"type": "number", "description": "[text/icon_circle] 字号(pt)"},
                                            "font_name": {"type": "string", "description": "[text] 字体名"},
                                            "bold": {"type": "boolean", "description": "[text] 是否粗体"},
                                            "italic": {"type": "boolean", "description": "[text] 是否斜体"},
                                            "color": {"type": "string", "description": "[text/icon_circle] 字体颜色 hex"},
                                            "bg_color": {"type": "string", "description": "[text] 文本框背景色 hex"},
                                            "align": {"type": "string", "description": "[text] 对齐: left|center|right"},
                                            "valign": {"type": "string", "description": "[text] 垂直对齐: top|middle|bottom"},
                                            "line_spacing": {"type": "number", "description": "[text] 行间距倍数(如1.5)"},
                                            "shape_type": {"type": "string", "description": "[shape] 形状: rectangle|rounded_rect|oval|triangle|arrow_right|arrow_left|chevron|diamond|pentagon|hexagon|star"},
                                            "fill_color": {"type": "string", "description": "[shape/icon_circle] 填充色 hex"},
                                            "line_color": {"type": "string", "description": "[shape/line] 线条颜色 hex"},
                                            "line_width": {"type": "number", "description": "[shape/line] 线宽(pt)"},
                                            "rotation": {"type": "number", "description": "[shape] 旋转角度(度)"},
                                            "chart_type": {"type": "string", "description": "[chart] 图表类型: bar|column|line|pie|doughnut|radar|area"},
                                            "categories": {"type": "array", "items": {"type": "string"}, "description": "[chart] 分类标签"},
                                            "series": {"type": "array", "description": "[chart] 数据系列 [{name,values}]"},
                                            "colors": {"type": "array", "items": {"type": "string"}, "description": "[chart] 系列颜色数组"},
                                            "show_labels": {"type": "boolean", "description": "[chart] 显示数据标签"},
                                            "show_percent": {"type": "boolean", "description": "[chart] 显示百分比(饼图)"},
                                            "show_legend": {"type": "boolean", "description": "[chart] 显示图例"},
                                            "headers": {"type": "array", "items": {"type": "string"}, "description": "[table] 表头"},
                                            "rows": {"type": "array", "description": "[table] 数据行 [[cell,...],...]"},
                                            "header_color": {"type": "string", "description": "[table] 表头背景色"},
                                            "header_font_color": {"type": "string", "description": "[table] 表头字色"},
                                            "stripe_color": {"type": "string", "description": "[table] 斑马纹颜色"},
                                            "path": {"type": "string", "description": "[image] 图片文件路径"},
                                            "text": {"type": "string", "description": "[icon_circle] 圆内文字"},
                                            "font_color": {"type": "string", "description": "[icon_circle] 文字颜色"},
                                        },
                                        "required": ["type"],
                                    },
                                },
                            },
                            "required": [],
                        },
                    },
                },
                "required": ["output_path", "slides"],
            },
        },
    },
    # ---- Desktop screenshot tool ----
    {
        "type": "function",
        "function": {
            "name": "desktop_screenshot",
            "description": "Take a screenshot of the local desktop. Can specify a region to crop. Returns PNG image path.",
            "parameters": {
                "type": "object",
                "properties": {
                    "output_path": {
                        "type": "string",
                        "description": "Optional path where the PNG will be saved (defaults to auto-generated path in working directory)",
                    },
                    "region": {
                        "type": "object",
                        "description": "Optional region to crop (x, y, w, h coordinates)",
                        "properties": {
                            "x": {"type": "integer", "description": "Top-left X coordinate"},
                            "y": {"type": "integer", "description": "Top-left Y coordinate"},
                            "w": {"type": "integer", "description": "Width in pixels"},
                            "h": {"type": "integer", "description": "Height in pixels"},
                        },
                    },
                },
            },
        },
    },
    # ---- Video creation tool ----
    {
        "type": "function",
        "function": {
            "name": "create_video",
            "description": "Create a video file from image frames. Auto-installs moviepy if needed. Can add audio track.",
            "parameters": {
                "type": "object",
                "properties": {
                    "output_path": {
                        "type": "string",
                        "description": "Path where the .mp4 video file will be saved",
                    },
                    "frames": {
                        "type": "array",
                        "description": "Array of frame objects with image_path and optional duration",
                        "items": {
                            "type": "object",
                            "properties": {
                                "image_path": {
                                    "type": "string",
                                    "description": "Path to the image file",
                                },
                                "duration": {
                                    "type": "number",
                                    "description": "Duration in seconds to display this frame (default: 3)",
                                },
                            },
                            "required": ["image_path"],
                        },
                    },
                    "fps": {
                        "type": "integer",
                        "description": "Frames per second for the video (default: 24)",
                    },
                    "audio_path": {
                        "type": "string",
                        "description": "Optional path to audio file to add as soundtrack",
                    },
                },
                "required": ["output_path", "frames"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_skill_guide",
            "description": (
                "Load the full operating guide for a granted skill. "
                "Returns the complete SKILL.md instructions, install_dir path, "
                "and ancillary file list. Use this when you need detailed "
                "step-by-step instructions or scripts for a specific skill "
                "(e.g. pdf, docx, xlsx, pptx). Run scripts from skill_dir."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {
                        "type": "string",
                        "description": "Skill name (e.g. 'pdf', 'docx', 'xlsx')",
                    },
                    "agent_id": {
                        "type": "string",
                        "description": "Optional agent ID to resolve agent-local skill path",
                    },
                },
                "required": ["name"],
            },
        },
    },
    # ---- Skill generation (propose a new skill from accumulated experiences) ----
    {
        "type": "function",
        "function": {
            "name": "propose_skill",
            "description": (
                "Scan the experience library for recurring patterns and propose a new skill draft. "
                "When enough similar, high-success-rate experiences exist (≥3, ≥75% success), "
                "SkillForge clusters them and generates a skill package (SKILL.md + manifest.yaml). "
                "The draft enters a pending-approval queue visible to admin in the portal. "
                "Use this after accumulating experiences on a topic (e.g. PPTX best practices, "
                "code review checklists) to crystallize them into a reusable skill."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "role": {
                        "type": "string",
                        "description": "Limit scan to experiences of this role (empty = all roles)",
                    },
                    "topic": {
                        "type": "string",
                        "description": "Optional topic hint to guide which experience cluster to target",
                    },
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "submit_skill",
            "description": (
                "Submit a skill package directory you created in your workspace for admin approval. "
                "The directory MUST contain a valid manifest.yaml with ALL required fields, and a SKILL.md.\n\n"
                "manifest.yaml REQUIRED fields:\n"
                "  name: string          # skill name, kebab-case (e.g. pptx-maker)\n"
                "  version: string       # semver (e.g. 1.0.0)\n"
                "  description: string   # one-line description\n"
                "  runtime: string       # MUST be one of: python, shell, markdown\n"
                "  author: string        # your name\n"
                "  entry: string         # entry file (e.g. main.py for python, SKILL.md for markdown)\n\n"
                "manifest.yaml OPTIONAL fields:\n"
                "  display_name: string  # human-friendly name\n"
                "  tags: [string]        # classification tags\n"
                "  triggers: [string]    # activation keywords\n"
                "  depends_on: object    # MCP or other dependencies\n"
                "  inputs: [{name, type, required, description}]  # input parameters\n"
                "  outputs: [{name, type}]  # output fields\n"
                "  hint: {sensitive: bool}  # security hint\n\n"
                "SKILL.md should document: what the skill does, how to use it, examples, and caveats.\n\n"
                "IMPORTANT for python runtime skills:\n"
                "  - The entry file (e.g. main.py) MUST define: def run(ctx, **kwargs)\n"
                "  - ctx provides: ctx.env('KEY') for environment variables, ctx.log(...), ctx.output(...)\n"
                "  - Do NOT use open(), exec(), eval(), __import__() — these are forbidden by the sandbox\n"
                "  - Use ctx.output(name, value) to return results instead of print()\n\n"
                "After submission the draft appears in the SkillForge review queue for admin approval."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "dir_name": {
                        "type": "string",
                        "description": "Name of the skill directory in your workspace (e.g. 'pptx_skill')",
                    },
                },
                "required": ["dir_name"],
            },
        },
    },
]


# ---------------------------------------------------------------------------
# Tool implementations
# ---------------------------------------------------------------------------

def _tool_read_file(path: str, offset: int = 0, limit: int | None = None, **_: Any) -> str:
    pol = _sandbox.get_current_policy()
    try:
        p = pol.safe_path(path)
    except _sandbox.SandboxViolation as e:
        return f"Error: {e}"
    if not p.exists():
        return f"Error: File not found: {path}"
    if not p.is_file():
        return f"Error: Not a file: {path}"
    try:
        lines = p.read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
    except Exception as e:
        return f"Error reading file: {e}"

    total = len(lines)
    start = max(0, offset)
    end = total if limit is None else min(total, start + limit)
    selected = lines[start:end]

    # Format with line numbers
    numbered = []
    for i, line in enumerate(selected, start=start + 1):
        numbered.append(f"{i:>6}\t{line.rstrip()}")
    header = f"[{p} — lines {start + 1}-{end} of {total}]"
    return header + "\n" + "\n".join(numbered)


def _tool_write_file(path: str, content: str, **_: Any) -> str:
    pol = _sandbox.get_current_policy()
    try:
        p = pol.safe_path(path)
    except _sandbox.SandboxViolation as e:
        return f"Error: {e}"
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding="utf-8")
        # Return the resolved absolute path so the artifact system can
        # locate the file reliably (relative paths break file card downloads).
        return f"Successfully wrote {len(content)} bytes to {p}"
    except Exception as e:
        return f"Error writing file: {e}"


def _tool_edit_file(path: str, old_string: str, new_string: str, **_: Any) -> str:
    pol = _sandbox.get_current_policy()
    try:
        p = pol.safe_path(path)
    except _sandbox.SandboxViolation as e:
        return f"Error: {e}"
    if not p.exists():
        return f"Error: File not found: {path}"
    try:
        text = p.read_text(encoding="utf-8")
    except Exception as e:
        return f"Error reading file: {e}"

    count = text.count(old_string)
    if count == 0:
        return f"Error: old_string not found in {path}"
    if count > 1:
        return f"Error: old_string found {count} times in {path}. Must be unique. Provide more context."

    new_text = text.replace(old_string, new_string, 1)
    p.write_text(new_text, encoding="utf-8")
    return f"Successfully edited {path} (replaced 1 occurrence)"


def _tool_bash(command: str, timeout: int = 30, **_: Any) -> str:
    pol = _sandbox.get_current_policy()
    ok, err = pol.check_command(command)
    if not ok:
        return f"Error: {err}"
    # Clamp timeout to protect against stuck processes
    try:
        timeout = max(1, min(int(timeout), 600))
    except Exception:
        timeout = 30
    try:
        jailed = pol.mode in ("restricted", "strict")
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
            # Always prefer the sandbox policy root (the agent's
            # working_dir). Falling back to os.getcwd() would run the
            # command in the server-process CWD (the code package
            # directory), causing runtime artefacts to leak into the
            # source tree.
            cwd=str(pol.root) if getattr(pol, "root", None) else os.getcwd(),
            env=pol.scrub_env() if jailed else None,
        )
        output_parts = []
        if result.stdout:
            output_parts.append(result.stdout)
        if result.stderr:
            output_parts.append(f"[stderr]\n{result.stderr}")
        output_parts.append(f"[exit code: {result.returncode}]")
        return "\n".join(output_parts)
    except subprocess.TimeoutExpired:
        return f"Error: Command timed out after {timeout}s"
    except Exception as e:
        return f"Error executing command: {e}"


def _tool_search_files(pattern: str, path: str = ".", include: str = "", **_: Any) -> str:
    pol = _sandbox.get_current_policy()
    try:
        base = pol.safe_path(path)
    except _sandbox.SandboxViolation as e:
        return f"Error: {e}"
    if not base.exists():
        return f"Error: Path not found: {path}"

    try:
        regex = re.compile(pattern)
    except re.error as e:
        return f"Error: Invalid regex pattern: {e}"

    matches = []
    max_matches = 200

    def _search_file(fpath: Path):
        try:
            with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                for lineno, line in enumerate(f, 1):
                    if regex.search(line):
                        matches.append(f"{fpath}:{lineno}: {line.rstrip()}")
                        if len(matches) >= max_matches:
                            return
        except (PermissionError, IsADirectoryError, OSError):
            pass

    if base.is_file():
        _search_file(base)
    else:
        for root, _dirs, files in os.walk(base):
            # Skip hidden / common noise directories
            root_path = Path(root)
            parts = root_path.parts
            if any(p.startswith(".") and p not in (".", "..") for p in parts):
                continue
            if any(p in ("node_modules", "__pycache__", ".git") for p in parts):
                continue

            for fname in files:
                if include and not fnmatch.fnmatch(fname, include):
                    continue
                _search_file(root_path / fname)
                if len(matches) >= max_matches:
                    break
            if len(matches) >= max_matches:
                break

    if not matches:
        return "No matches found."
    result = "\n".join(matches)
    if len(matches) >= max_matches:
        result += f"\n... (truncated at {max_matches} matches)"
    return result


def _tool_glob_files(pattern: str, path: str = ".", **_: Any) -> str:
    pol = _sandbox.get_current_policy()
    try:
        base = pol.safe_path(path)
    except _sandbox.SandboxViolation as e:
        return f"Error: {e}"
    if not base.exists():
        return f"Error: Path not found: {path}"

    found = sorted(base.glob(pattern))
    # Filter out hidden dirs
    filtered = [
        str(f) for f in found
        if not any(part.startswith(".") and part not in (".", "..") for part in f.parts)
    ]
    if not filtered:
        return "No files found."
    if len(filtered) > 500:
        return "\n".join(filtered[:500]) + f"\n... ({len(filtered)} total, showing first 500)"
    return "\n".join(filtered)


# ---------------------------------------------------------------------------
# Coordination tools — TeamCreate / SendMessage / TaskUpdate
# ---------------------------------------------------------------------------

def _get_hub():
    """Lazy import to avoid circular dependency."""
    from .hub import get_hub
    return get_hub()


def _tool_team_create(name: str, task: str, role: str = "coder",
                      working_dir: str = "", **_: Any) -> str:
    """Spawn a background worker to run a task in parallel.

    The worker is NOT a first-class Agent — it's a transient background
    task owned by the caller. It inherits the caller's model/provider,
    runs the task to completion, pushes a result entry into the caller's
    task list, then disappears. The UI never sees it as a separate agent.
    """
    import threading as _threading
    import uuid as _uuid
    try:
        hub = _get_hub()
        caller_id = _.get("_caller_agent_id", "") if isinstance(_, dict) else ""
        parent = hub.get_agent(caller_id) if caller_id else None
        if parent is None:
            return ("Error: team_create requires a calling agent context; "
                    "none was found.")

        worker_id = _uuid.uuid4().hex[:8]
        worker_label = f"{role}:{name}" if name else role

        # Record the background job as a task on the PARENT agent so the
        # user can track it from the parent's task list / execution log.
        try:
            t = parent.add_task(
                title=f"[bg:{worker_label}] {task[:80]}",
                description=task,
            )
            task_id = t.id
        except Exception:
            task_id = ""

        def _run_background():
            from .agent import create_agent as _create_agent_fn
            try:
                # Build an ephemeral worker that inherits parent's config.
                # Resolve working directory: explicit > parent's shared_workspace > parent's working_dir
                # This ensures child agents in a project share the same directory.
                _wd = working_dir or parent.shared_workspace or parent.working_dir
                worker = _create_agent_fn(
                    name=f"__bg_{worker_label}_{worker_id}",
                    role=role,
                    model=parent.model,
                    provider=parent.provider,
                    working_dir=_wd,
                    node_id=parent.node_id,
                    parent_id=parent.id,
                )
                # Inherit project context so child knows where to write files
                worker.shared_workspace = parent.shared_workspace
                worker.project_id = parent.project_id
                worker.project_name = parent.project_name
                # Don't register it in the hub — it's transient.
                result_text = ""
                try:
                    result_text = worker.chat(task) or ""
                except Exception as e:
                    result_text = f"Worker error: {e}"
                # Push result back as a completed task entry on the parent
                if task_id:
                    try:
                        parent.update_task(
                            task_id,
                            status="done",
                            result=(result_text or "")[:4000],
                        )
                    except Exception:
                        pass
                # Log to parent event stream for visibility
                try:
                    parent._log("bg_task_complete", {
                        "worker": worker_label,
                        "task_id": task_id,
                        "result_preview": (result_text or "")[:200],
                    })
                except Exception:
                    pass
            except Exception as e:
                if task_id:
                    try:
                        parent.update_task(task_id, status="failed",
                                           result=f"{type(e).__name__}: {e}")
                    except Exception:
                        pass

        th = _threading.Thread(target=_run_background, daemon=True,
                               name=f"bg-{worker_label}-{worker_id}")
        th.start()

        return (
            f"Background worker dispatched.\n"
            f"  Role: {role}\n"
            f"  Worker: {worker_label} (id={worker_id})\n"
            f"  Task ID on parent: {task_id or '(none)'}\n"
            f"  Model inherited: {parent.model or '(default)'} @ "
            f"{parent.provider or '(default)'}\n"
            f"The worker runs in background and will post its result back "
            f"to your task list when done. It is NOT a separately managed agent."
        )
    except Exception as e:
        return f"Error dispatching background worker: {e}"


def _tool_send_message(to_agent: str, content: str,
                       msg_type: str = "task", **_: Any) -> str:
    """Send an inter-agent message."""
    try:
        hub = _get_hub()
        # Resolve agent by name if not an ID
        target = hub.get_agent(to_agent)
        if target is None:
            # Try finding by name
            for a in hub.agents.values():
                if a.name.lower() == to_agent.lower():
                    target = a
                    break
        if target is None:
            available = [f"{a.name} ({a.id})" for a in hub.agents.values()]
            return (
                f"Error: Agent '{to_agent}' not found.\n"
                f"Available agents: {', '.join(available) or 'none'}"
            )
        # Use hub's canonical routing entry point (audited)
        caller_id = _.get("_caller_agent_id", "unknown") if isinstance(_, dict) else "unknown"
        route = getattr(hub, "route_message", None)
        if callable(route):
            route(caller_id, target.id, content, msg_type=msg_type, source="tool_send_message")
        else:
            hub.send_message(caller_id, target.id, content, msg_type=msg_type)
        return (
            f"Message sent to {target.name} ({target.id}).\n"
            f"  Type: {msg_type}\n"
            f"  Content: {content[:200]}"
        )
    except Exception as e:
        return f"Error sending message: {e}"


def _parse_run_at(run_at: str) -> float:
    """Parse run_at spec into a unix timestamp.

    Supported formats:
      '+5m'   → 5 minutes from now
      '+2h'   → 2 hours from now
      '18:30' → today at 18:30 (or tomorrow if already past)
    Returns 0.0 on failure.
    """
    import re as _re
    from datetime import datetime as _dt, timedelta as _td
    run_at = run_at.strip()
    if not run_at:
        return 0.0
    # Relative: +Nm / +Nh / +Ns
    m = _re.match(r'^\+(\d+)\s*([mMhHsS])$', run_at)
    if m:
        val = int(m.group(1))
        unit = m.group(2).lower()
        delta = {'m': _td(minutes=val), 'h': _td(hours=val),
                 's': _td(seconds=val)}[unit]
        return (_dt.now() + delta).timestamp()
    # Absolute: HH:MM
    m = _re.match(r'^(\d{1,2}):(\d{2})$', run_at)
    if m:
        hh, mm = int(m.group(1)), int(m.group(2))
        now = _dt.now()
        target = now.replace(hour=hh, minute=mm, second=0, microsecond=0)
        if target <= now:
            target += _td(days=1)
        return target.timestamp()
    return 0.0


def _tool_task_update(action: str, task_id: str = "", title: str = "",
                      description: str = "", status: str = "",
                      result: str = "",
                      recurrence: str = "once",
                      recurrence_spec: str = "",
                      run_at: str = "", **_: Any) -> str:
    """Manage the shared task list and register with scheduler for execution."""
    try:
        hub = _get_hub()
        # Use the calling agent's task list
        caller_id = _.get("_caller_agent_id", "") if isinstance(_, dict) else ""
        agent = hub.get_agent(caller_id) if caller_id else None

        if action == "list":
            # List all tasks across all agents
            all_tasks = []
            for a in hub.agents.values():
                for t in a.tasks:
                    all_tasks.append(
                        f"  [{t.status.value:>11}] {t.id}: {t.title} "
                        f"(agent: {a.name})"
                    )
            if not all_tasks:
                return "No tasks found."
            return f"Shared task list ({len(all_tasks)} tasks):\n" + "\n".join(all_tasks)

        if action == "create":
            if not title:
                return "Error: 'title' is required for create action."
            rec = (recurrence or "once").lower()
            agent_id = caller_id or (agent.id if agent else "")

            # Route through agent.add_task so recurrence / next_run_at is computed
            if agent:
                new_task = agent.add_task(
                    title=title,
                    description=description,
                    assigned_by=caller_id or "system",
                    source="agent_chat",
                    recurrence=rec,
                    recurrence_spec=recurrence_spec or "",
                )
            else:
                # Fallback: no agent context — create plain task
                from .agent import AgentTask, TaskStatus
                new_task = AgentTask(
                    title=title,
                    description=description,
                    status=TaskStatus(status) if status else TaskStatus.TODO,
                    assigned_by=caller_id or "system",
                    recurrence=rec,
                    recurrence_spec=recurrence_spec or "",
                )

            # ── Register with TaskScheduler for actual execution ──
            # This is the critical bridge: AgentTask → ScheduledJob
            #
            # GUARD: When the agent is running INSIDE a scheduled task,
            # block it from creating new scheduled jobs.  Otherwise
            # "please generate daily report" prompts cause the agent to
            # create duplicate recurring jobs on every execution.
            _in_scheduled = getattr(agent, '_scheduled_context', False) if agent else False
            if _in_scheduled and (rec != "once" or run_at):
                return (
                    f"Task created: {new_task.id} — {title}"
                    f" [NOTE: scheduler registration skipped — "
                    f"you are already running inside a scheduled job]"
                )

            import datetime as _dt
            scheduler_msg = ""
            try:
                from .scheduler import get_scheduler, recurrence_to_cron
                scheduler = get_scheduler()

                if rec != "once":
                    # Recurring task → register as recurring scheduler job
                    cron_expr = recurrence_to_cron(rec, recurrence_spec or "")
                    if cron_expr and scheduler and agent_id:
                        job = scheduler.add_job(
                            agent_id=agent_id,
                            name=title,
                            prompt_template=description or title,
                            job_type="recurring",
                            cron_expr=cron_expr,
                        )
                        nxt = _dt.datetime.fromtimestamp(
                            job.next_run_at).strftime("%Y-%m-%d %H:%M")
                        scheduler_msg = (
                            f" [SCHEDULED: recurring {rec} @ "
                            f"{recurrence_spec or 'default'}, "
                            f"next run: {nxt}, job_id: {job.id}]")

                elif run_at:
                    # One-time delayed task → register as one_time scheduler job
                    run_ts = _parse_run_at(run_at)
                    if run_ts > 0 and scheduler and agent_id:
                        job = scheduler.add_job(
                            agent_id=agent_id,
                            name=title,
                            prompt_template=description or title,
                            job_type="one_time",
                            cron_expr="* * * * *",  # placeholder
                            next_run_at=run_ts,
                        )
                        nxt = _dt.datetime.fromtimestamp(
                            run_ts).strftime("%Y-%m-%d %H:%M")
                        scheduler_msg = (
                            f" [SCHEDULED: one-time at {nxt}, "
                            f"job_id: {job.id}]")
                    elif run_ts <= 0:
                        scheduler_msg = (
                            f" [WARNING: could not parse run_at='{run_at}', "
                            f"task created but NOT scheduled]")

            except Exception as sched_err:
                logger.warning("Failed to register task with scheduler: %s",
                               sched_err)
                scheduler_msg = f" [scheduler registration failed: {sched_err}]"

            return f"Task created: {new_task.id} — {title}{scheduler_msg}"

        if action in ("update", "complete"):
            if not task_id:
                return "Error: 'task_id' is required for update/complete."
            # Find the task across all agents
            from .agent import TaskStatus
            for a in hub.agents.values():
                for t in a.tasks:
                    if t.id == task_id:
                        if action == "complete":
                            t.status = TaskStatus.DONE
                            t.result = result or "Completed"
                        elif status:
                            t.status = TaskStatus(status)
                        if description:
                            t.description = description
                        t.updated_at = time.time()
                        return f"Task {task_id} updated: status={t.status.value}"
            return f"Error: Task '{task_id}' not found."

        return f"Error: Unknown action '{action}'. Use: create | update | complete | list"
    except Exception as e:
        return f"Error managing tasks: {e}"


# ---------------------------------------------------------------------------
# Web tools — DuckDuckGo search & page fetch
# ---------------------------------------------------------------------------

def _tool_web_search(query: str, max_results: int = 8, **_: Any) -> str:
    """Search the internet using DuckDuckGo (API + HTML fallback)."""
    import urllib.request
    import urllib.parse
    import html as html_mod

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        ),
    }

    # ── Strategy 1: DuckDuckGo Instant Answer API (fast, structured) ──
    try:
        api_url = "https://api.duckduckgo.com/?" + urllib.parse.urlencode({
            "q": query, "format": "json", "no_html": "1",
            "skip_disambig": "1", "no_redirect": "1",
        })
        req = urllib.request.Request(api_url, headers=headers)
        import json as _json
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = _json.loads(resp.read().decode("utf-8", errors="replace"))
        results = []
        # Abstract (direct answer)
        if data.get("Abstract"):
            results.append(
                f"0. {data.get('Heading', query)} (Direct Answer)\n"
                f"   URL: {data.get('AbstractURL', '')}\n"
                f"   {data['Abstract']}"
            )
        # Related topics
        for i, topic in enumerate(data.get("RelatedTopics", [])[:max_results]):
            if isinstance(topic, dict) and topic.get("Text"):
                url_ = topic.get("FirstURL", "")
                text = topic["Text"]
                results.append(f"{i+1}. {text[:120]}\n   URL: {url_}")
        if results:
            return f"Search results for: {query}\n\n" + "\n\n".join(results)
    except Exception:
        pass  # Fall through to HTML scraping

    # ── Strategy 2: DuckDuckGo HTML scraping (more results) ──
    url = "https://html.duckduckgo.com/html/?" + urllib.parse.urlencode({"q": query})
    req = urllib.request.Request(url, headers=headers, method="POST",
                                 data=urllib.parse.urlencode({"q": query}).encode())
    try:
        with urllib.request.urlopen(req, timeout=15) as resp:
            body = resp.read().decode("utf-8", errors="replace")
    except Exception as e:
        return f"Error: Web search failed: {e}"

    # Parse results from DuckDuckGo HTML
    results = []
    link_pattern = re.compile(
        r'class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>', re.DOTALL
    )
    snippet_pattern = re.compile(
        r'class="result__snippet"[^>]*>(.*?)</a>', re.DOTALL
    )

    links = link_pattern.findall(body)
    snippets = snippet_pattern.findall(body)

    for i, (raw_url, raw_title) in enumerate(links[:max_results]):
        title = re.sub(r"<[^>]+>", "", raw_title).strip()
        title = html_mod.unescape(title)
        actual_url = raw_url
        m = re.search(r'uddg=([^&]+)', raw_url)
        if m:
            actual_url = urllib.parse.unquote(m.group(1))
        snippet = ""
        if i < len(snippets):
            snippet = re.sub(r"<[^>]+>", "", snippets[i]).strip()
            snippet = html_mod.unescape(snippet)
        results.append(f"{i+1}. {title}\n   URL: {actual_url}\n   {snippet}")

    if not results:
        return "No search results found."
    return f"Search results for: {query}\n\n" + "\n\n".join(results)


def _tool_web_fetch(url: str, max_length: int = 10000, **_: Any) -> str:
    """Fetch the text content of a web page URL."""
    import urllib.request

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
            "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        ),
    }
    req = urllib.request.Request(url, headers=headers)
    try:
        with urllib.request.urlopen(req, timeout=20) as resp:
            content_type = resp.headers.get("Content-Type", "")
            raw = resp.read()
            # Detect encoding
            encoding = "utf-8"
            if "charset=" in content_type:
                encoding = content_type.split("charset=")[-1].split(";")[0].strip()
            body = raw.decode(encoding, errors="replace")
    except Exception as e:
        return f"Error: Failed to fetch URL: {e}"

    # Strip HTML tags to get plain text
    import html as html_mod
    # Remove script and style blocks
    text = re.sub(r"<script[^>]*>.*?</script>", "", body, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
    # Replace block tags with newlines
    text = re.sub(r"<(?:br|p|div|li|tr|h[1-6])[^>]*>", "\n", text, flags=re.IGNORECASE)
    # Remove remaining tags
    text = re.sub(r"<[^>]+>", "", text)
    text = html_mod.unescape(text)
    # Collapse whitespace
    lines = [line.strip() for line in text.splitlines()]
    text = "\n".join(line for line in lines if line)

    if len(text) > max_length:
        text = text[:max_length] + f"\n\n... (truncated at {max_length} characters)"

    return f"[Content from {url}]\n\n{text}"


# ---------------------------------------------------------------------------
# Builtin MCP handlers — audio TTS / STT
# ---------------------------------------------------------------------------

# Audio event queue: Portal UI polls this to play TTS / handle STT
_audio_events: list[dict] = []
_audio_lock = threading.Lock()


def _push_audio_event(event: dict):
    """Push an audio event for Portal UI to consume."""
    with _audio_lock:
        _audio_events.append(event)
        # Keep last 50 events
        if len(_audio_events) > 50:
            _audio_events[:] = _audio_events[-50:]


def get_audio_events(since: int = 0) -> list[dict]:
    """Get audio events (called by portal API)."""
    with _audio_lock:
        return [e for e in _audio_events if e.get("ts", 0) > since]


def _handle_builtin_mcp(target: Any, tool_name: str, arguments: Any,
                        agent: Any) -> str:
    """Handle builtin MCP tools (audio TTS/STT)."""
    import json as _json
    args = arguments if isinstance(arguments, dict) else {}
    if isinstance(arguments, str):
        try:
            args = _json.loads(arguments)
        except Exception:
            args = {}

    mcp_type = (target.command or "").replace("__builtin__", "")

    if mcp_type == "audio_tts":
        if tool_name == "speak":
            text = args.get("text", "")
            if not text:
                return "Error: 'text' argument is required for speak."
            lang = args.get("lang", target.env.get("TTS_LANG", "zh-CN"))
            rate = float(args.get("rate", target.env.get("TTS_RATE", "1.0")))
            voice = args.get("voice", target.env.get("TTS_VOICE", ""))
            _push_audio_event({
                "type": "tts_speak",
                "agent_id": agent.id,
                "agent_name": agent.name,
                "text": text,
                "lang": lang,
                "rate": rate,
                "voice": voice,
                "ts": time.time(),
            })
            return f"Speaking: \"{text[:100]}{'...' if len(text)>100 else ''}\" [lang={lang}, rate={rate}]"

        elif tool_name == "set_voice":
            voice = args.get("voice", "")
            lang = args.get("lang", "")
            return f"Voice preference set: voice={voice}, lang={lang}. Will take effect on next speak()."

        elif tool_name == "list_voices":
            return ("Available voices depend on the user's browser. Common ones:\n"
                    "  - zh-CN: Microsoft Xiaoxiao, Google 普通话\n"
                    "  - en-US: Google US English, Microsoft David\n"
                    "  - ja-JP: Google 日本語\n"
                    "Use speak(text, voice='name') to select a specific voice.")
        else:
            return f"Error: TTS tool '{tool_name}' not found. Available: speak, set_voice, list_voices"

    elif mcp_type == "audio_stt":
        if tool_name == "listen":
            duration = int(args.get("duration", 5))
            lang = args.get("lang", target.env.get("STT_LANG", "zh-CN"))
            _push_audio_event({
                "type": "stt_listen",
                "agent_id": agent.id,
                "agent_name": agent.name,
                "duration": duration,
                "lang": lang,
                "ts": time.time(),
            })
            return (f"Listening request sent to browser (lang={lang}, "
                    f"duration={duration}s). The user's speech will be "
                    f"transcribed and sent as the next user message.")

        elif tool_name == "start_listening":
            lang = args.get("lang", target.env.get("STT_LANG", "zh-CN"))
            _push_audio_event({
                "type": "stt_start",
                "agent_id": agent.id,
                "lang": lang,
                "ts": time.time(),
            })
            return f"Continuous listening started (lang={lang}). Speech will be sent as messages."

        elif tool_name == "stop_listening":
            _push_audio_event({
                "type": "stt_stop",
                "agent_id": agent.id,
                "ts": time.time(),
            })
            return "Listening stopped."
        else:
            return f"Error: STT tool '{tool_name}' not found. Available: listen, start_listening, stop_listening"

    return f"Error: builtin MCP type '{mcp_type}' not recognized."


# ---------------------------------------------------------------------------
# MCP call — thin agent-facing wrapper around the central MCP Call Router.
# ---------------------------------------------------------------------------
#
# Architectural note (READ THIS BEFORE ADDING CODE HERE):
#
# This function used to own ~260 lines of subprocess launch, JSON-RPC
# protocol, and env-variable normalization. All of that has been moved
# into ``app/mcp/dispatcher.py`` (the executor) and ``app/mcp/router.py``
# (the router / auth / classifier). The agent-side API — this function
# — is now a thin adapter whose only job is to turn a tool-call into
# ``client_stub.call(...)``.
#
# If you feel the urge to add subprocess handling, path logic, or env
# injection here again, STOP: those belong in the dispatcher. Keeping
# this function tiny is the architectural invariant that prevents
# path/cwd/env bugs from multiplying across the codebase.
#
# The ``_handle_builtin_mcp`` helper further below is registered with
# the dispatcher at import time so builtin TTS/STT calls flow through
# the same router/dispatcher pipeline as external MCPs.

def _tool_mcp_call(mcp_id: str = "", tool: str = "", arguments: Any = None,
                   list_mcps: bool = False, **_: Any) -> str:
    """Invoke an MCP tool bound to the calling agent.

    Set ``list_mcps=True`` to enumerate the MCPs visible to this
    agent. Otherwise this call is dispatched through the central
    :mod:`app.mcp.client_stub` → :class:`~app.mcp.router.MCPCallRouter`
    → :class:`~app.mcp.dispatcher.NodeMCPDispatcher` pipeline.
    """
    try:
        caller_id = _.get("_caller_agent_id", "") if isinstance(_, dict) else ""
        if not caller_id:
            return "Error: no calling agent context; mcp_call requires an agent."

        from .mcp import client_stub as _stub

        # List mode: delegate to the router's enumeration path.
        if list_mcps or not mcp_id:
            return _stub.list_mcps(caller_id)

        # Normalize arguments — the router/dispatcher wants a dict.
        import json as _json
        args: dict
        if isinstance(arguments, dict):
            args = arguments
        elif isinstance(arguments, str):
            try:
                args = _json.loads(arguments) if arguments.strip() else {}
            except Exception:
                return f"Error: 'arguments' must be a JSON object, got: {arguments!r}"
        elif arguments is None:
            args = {}
        else:
            return f"Error: 'arguments' must be a JSON object, got: {type(arguments).__name__}"

        return _stub.call(
            caller_id=caller_id,
            mcp_id=mcp_id,
            tool=tool,
            arguments=args,
        )
    except Exception as e:
        return f"Error in mcp_call: {e}"


# Register the in-process builtin handler with the dispatcher so
# builtin MCPs (audio TTS/STT) flow through the same router pipeline
# as external stdio MCPs. This replaces the old inline branch in
# _tool_mcp_call. Registration happens at module import time.
try:
    from .mcp.dispatcher import register_builtin_handler as _register_builtin
    _register_builtin("__builtin__audio", _handle_builtin_mcp)
    _register_builtin("builtin", _handle_builtin_mcp)
except Exception:
    # Not fatal — builtins simply won't be dispatchable until
    # registration is retried. Never block module import for this.
    pass


# ---------------------------------------------------------------------------
# New daily-work tools — screenshot, http, datetime, json, text
# ---------------------------------------------------------------------------

def _tool_web_screenshot(url: str, output_path: str = "", full_page: bool = False,
                         width: int = 1280, height: int = 720, **_: Any) -> str:
    """Take a screenshot of a web page using Playwright or Selenium."""
    import importlib
    import base64

    # Determine output path
    if not output_path:
        import hashlib
        url_hash = hashlib.md5(url.encode()).hexdigest()[:8]
        ts = int(time.time())
        output_path = f"/tmp/screenshot_{url_hash}_{ts}.png"

    # Strategy 1: Try Playwright (preferred)
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": width, "height": height})
            page.goto(url, wait_until="networkidle", timeout=30000)
            page.screenshot(path=output_path, full_page=full_page)
            browser.close()
        size = os.path.getsize(output_path)
        return f"Screenshot saved to {output_path} ({size} bytes, {width}x{height})"
    except ImportError:
        logger.debug("Playwright not installed, trying fallback methods")
    except Exception as e:
        # Playwright installed but failed; try fallback
        logger.debug("Playwright screenshot failed: %s, trying fallback methods", e)

    # Strategy 2: Try subprocess with playwright CLI
    try:
        cmd = f'python3 -c "from playwright.sync_api import sync_playwright; p=sync_playwright().start(); b=p.chromium.launch(headless=True); pg=b.new_page(viewport={{\'width\':{width},\'height\':{height}}}); pg.goto(\'{url}\',wait_until=\'networkidle\',timeout=30000); pg.screenshot(path=\'{output_path}\',full_page={full_page}); b.close(); p.stop(); print(\'ok\')"'
        result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=45)
        if result.returncode == 0 and os.path.exists(output_path):
            size = os.path.getsize(output_path)
            return f"Screenshot saved to {output_path} ({size} bytes, {width}x{height})"
    except Exception as e:
        logger.debug("Playwright subprocess failed: %s, trying other methods", e)

    # Strategy 3: Use cutycapt or wkhtmltoimage if available
    for cmd_name, cmd_tpl in [
        ("wkhtmltoimage", f"wkhtmltoimage --width {width} --height {height} '{url}' '{output_path}'"),
        ("cutycapt", f"cutycapt --url='{url}' --out='{output_path}' --min-width={width} --min-height={height}"),
    ]:
        try:
            result = subprocess.run(f"which {cmd_name}", shell=True, capture_output=True, timeout=5)
            if result.returncode == 0:
                result = subprocess.run(cmd_tpl, shell=True, capture_output=True, text=True, timeout=30)
                if os.path.exists(output_path):
                    size = os.path.getsize(output_path)
                    return f"Screenshot saved to {output_path} ({size} bytes, {width}x{height})"
        except Exception:
            continue

    return (
        "Error: Screenshot tools not available. Please install one of:\n"
        "  pip install playwright && playwright install chromium\n"
        "  apt install wkhtmltopdf\n"
        "You can also use the 'browser' MCP (Puppeteer) for screenshots."
    )


def _tool_http_request(url: str, method: str = "GET", headers: dict = None,
                       body: str = "", json_body: dict = None,
                       timeout: int = 30, **_: Any) -> str:
    """Make an HTTP request to any URL."""
    import urllib.request
    import urllib.parse
    import json as _json

    method = method.upper()
    req_headers = {
        "User-Agent": "TudouClaw-Agent/1.0",
    }
    if headers:
        req_headers.update(headers)

    data = None
    if json_body is not None:
        data = _json.dumps(json_body).encode("utf-8")
        req_headers.setdefault("Content-Type", "application/json")
    elif body:
        data = body.encode("utf-8")

    req = urllib.request.Request(url, data=data, headers=req_headers, method=method)
    try:
        timeout = max(1, min(int(timeout), 120))
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            status = resp.status
            resp_headers = dict(resp.headers)
            resp_body = resp.read().decode("utf-8", errors="replace")
    except urllib.error.HTTPError as e:
        status = e.code
        resp_headers = dict(e.headers) if hasattr(e, 'headers') else {}
        try:
            resp_body = e.read().decode("utf-8", errors="replace")
        except Exception:
            resp_body = str(e)
    except Exception as e:
        return f"Error: HTTP request failed: {e}"

    # Format response
    result = f"HTTP {status} {method} {url}\n"
    result += "--- Headers ---\n"
    for k, v in list(resp_headers.items())[:20]:
        result += f"  {k}: {v}\n"
    result += "--- Body ---\n"
    if len(resp_body) > MAX_HTTP_RESPONSE_CHARS:
        resp_body = resp_body[:MAX_HTTP_RESPONSE_CHARS] + f"\n... (truncated at {MAX_HTTP_RESPONSE_CHARS} chars, total: {len(resp_body)})"
    result += resp_body
    return result


def _tool_datetime_calc(action: str, date: str = "", date2: str = "",
                        days: int = 0, hours: int = 0, minutes: int = 0,
                        timezone: str = "", format: str = "",
                        **_: Any) -> str:
    """Perform date/time calculations."""
    from datetime import datetime, timedelta
    import zoneinfo

    def _parse_date(s: str) -> datetime:
        """Try multiple date formats."""
        for fmt in [
            "%Y-%m-%dT%H:%M:%S", "%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M",
            "%Y-%m-%d", "%Y/%m/%d %H:%M:%S", "%Y/%m/%d",
            "%m/%d/%Y", "%d-%m-%Y",
        ]:
            try:
                return datetime.strptime(s, fmt)
            except ValueError:
                continue
        raise ValueError(f"Cannot parse date: {s}")

    def _get_tz(tz_name: str):
        if not tz_name:
            return None
        try:
            return zoneinfo.ZoneInfo(tz_name)
        except Exception:
            return None

    try:
        if action == "now":
            tz = _get_tz(timezone)
            now = datetime.now(tz)
            fmt = format or "%Y-%m-%d %H:%M:%S %Z"
            return f"Current time: {now.strftime(fmt)}\nTimezone: {timezone or 'local'}\nISO: {now.isoformat()}"

        elif action == "diff":
            d1 = _parse_date(date)
            d2 = _parse_date(date2)
            delta = d2 - d1
            total_secs = int(delta.total_seconds())
            abs_secs = abs(total_secs)
            d = abs_secs // 86400
            h = (abs_secs % 86400) // 3600
            m = (abs_secs % 3600) // 60
            sign = "-" if total_secs < 0 else ""
            return (f"Date 1: {d1}\nDate 2: {d2}\n"
                    f"Difference: {sign}{d} days, {h} hours, {m} minutes\n"
                    f"Total days: {delta.days}\n"
                    f"Total seconds: {total_secs}")

        elif action == "add":
            d = _parse_date(date)
            delta = timedelta(days=days, hours=hours, minutes=minutes)
            result = d + delta
            fmt = format or "%Y-%m-%d %H:%M:%S"
            return (f"Original: {d.strftime(fmt)}\n"
                    f"Added: {days}d {hours}h {minutes}m\n"
                    f"Result: {result.strftime(fmt)}\n"
                    f"ISO: {result.isoformat()}")

        elif action == "format":
            d = _parse_date(date)
            fmt = format or "%Y-%m-%d %H:%M:%S"
            return f"Formatted: {d.strftime(fmt)}\nISO: {d.isoformat()}"

        elif action == "convert":
            d = _parse_date(date)
            tz = _get_tz(timezone)
            if tz is None:
                return f"Error: Unknown timezone: {timezone}"
            if d.tzinfo is None:
                import zoneinfo as zi
                d = d.replace(tzinfo=zi.ZoneInfo("UTC"))
            converted = d.astimezone(tz)
            fmt = format or "%Y-%m-%d %H:%M:%S %Z"
            return (f"Original: {d.strftime(fmt)}\n"
                    f"Converted to {timezone}: {converted.strftime(fmt)}\n"
                    f"ISO: {converted.isoformat()}")
        else:
            return f"Error: Unknown action '{action}'. Use: now, diff, add, format, convert"
    except Exception as e:
        return f"Error: {e}"


def _tool_json_process(action: str, data: str, path: str = "",
                       data2: str = "", **_: Any) -> str:
    """Process JSON data."""
    import json as _json

    def _load(s: str):
        """Load JSON from string or file."""
        s = s.strip()
        if s.startswith("/") or s.startswith("./"):
            try:
                with open(s, "r", encoding="utf-8") as f:
                    return _json.load(f)
            except Exception as e:
                raise ValueError(f"Failed to read file {s}: {e}")
        return _json.loads(s)

    def _extract(obj, path_str: str):
        """Simple JSONPath-like extraction: 'a.b[0].c'"""
        parts = re.split(r'\.|\[(\d+)\]', path_str)
        parts = [p for p in parts if p is not None and p != '']
        for p in parts:
            if isinstance(obj, dict):
                obj = obj[p]
            elif isinstance(obj, list):
                obj = obj[int(p)]
            else:
                raise KeyError(f"Cannot navigate '{p}' in {type(obj).__name__}")
        return obj

    try:
        if action == "parse":
            obj = _load(data)
            formatted = _json.dumps(obj, indent=2, ensure_ascii=False)
            return f"Valid JSON ({type(obj).__name__}):\n{formatted[:10000]}"

        elif action == "extract":
            obj = _load(data)
            result = _extract(obj, path)
            if isinstance(result, (dict, list)):
                return _json.dumps(result, indent=2, ensure_ascii=False)[:MAX_JSON_RESULT_CHARS]
            return str(result)

        elif action == "keys":
            obj = _load(data)
            if isinstance(obj, dict):
                return f"Keys ({len(obj)}): " + ", ".join(str(k) for k in obj.keys())
            elif isinstance(obj, list):
                return f"Array with {len(obj)} items"
            return f"Type: {type(obj).__name__}, Value: {str(obj)[:200]}"

        elif action == "flatten":
            obj = _load(data)
            flat = {}
            def _flatten(o, prefix=""):
                if isinstance(o, dict):
                    for k, v in o.items():
                        _flatten(v, f"{prefix}{k}.")
                elif isinstance(o, list):
                    for i, v in enumerate(o):
                        _flatten(v, f"{prefix}{i}.")
                else:
                    flat[prefix.rstrip(".")] = o
            _flatten(obj)
            return _json.dumps(flat, indent=2, ensure_ascii=False)[:10000]

        elif action == "to_csv":
            obj = _load(data)
            if not isinstance(obj, list) or not obj:
                return "Error: Input must be a non-empty JSON array of objects"
            headers = list(obj[0].keys()) if isinstance(obj[0], dict) else []
            if not headers:
                return "Error: Array items must be objects"
            lines = [",".join(headers)]
            for item in obj:
                vals = [str(item.get(h, "")).replace(",", ";").replace("\n", " ") for h in headers]
                lines.append(",".join(vals))
            return "\n".join(lines)

        elif action == "from_csv":
            lines = data.strip().splitlines()
            if len(lines) < 2:
                return "Error: CSV must have at least header + 1 data row"
            headers = [h.strip() for h in lines[0].split(",")]
            result = []
            for line in lines[1:]:
                vals = [v.strip() for v in line.split(",")]
                result.append(dict(zip(headers, vals)))
            return _json.dumps(result, indent=2, ensure_ascii=False)[:10000]

        elif action == "merge":
            obj1 = _load(data)
            obj2 = _load(data2)
            if isinstance(obj1, dict) and isinstance(obj2, dict):
                merged = {**obj1, **obj2}
            elif isinstance(obj1, list) and isinstance(obj2, list):
                merged = obj1 + obj2
            else:
                return "Error: Both inputs must be same type (both objects or both arrays)"
            return _json.dumps(merged, indent=2, ensure_ascii=False)[:10000]

        elif action == "count":
            obj = _load(data)
            if isinstance(obj, list):
                return f"Array: {len(obj)} items"
            elif isinstance(obj, dict):
                return f"Object: {len(obj)} keys"
            return f"Type: {type(obj).__name__}"

        else:
            return f"Error: Unknown action '{action}'. Use: parse, extract, keys, flatten, to_csv, from_csv, merge, count"
    except Exception as e:
        return f"Error: {e}"


def _tool_text_process(action: str, text: str, pattern: str = "",
                       replacement: str = "", n: int = 10,
                       algorithm: str = "sha256", delimiter: str = "\n",
                       **_: Any) -> str:
    """Process and transform text."""
    try:
        if action == "count":
            lines = text.splitlines()
            words = text.split()
            chars = len(text)
            return f"Lines: {len(lines)}\nWords: {len(words)}\nCharacters: {chars}"

        elif action == "replace":
            if not pattern:
                return "Error: 'pattern' required for replace"
            result = re.sub(pattern, replacement, text)
            count = len(re.findall(pattern, text))
            return f"Replaced {count} occurrences.\n\n{result[:10000]}"

        elif action == "extract":
            if not pattern:
                return "Error: 'pattern' required for extract"
            matches = re.findall(pattern, text)
            if not matches:
                return "No matches found."
            return f"Found {len(matches)} matches:\n" + "\n".join(str(m) for m in matches[:200])

        elif action == "sort":
            lines = text.splitlines()
            sorted_lines = sorted(lines)
            return "\n".join(sorted_lines)

        elif action == "dedup":
            lines = text.splitlines()
            seen = set()
            unique = []
            for line in lines:
                if line not in seen:
                    seen.add(line)
                    unique.append(line)
            removed = len(lines) - len(unique)
            return f"Removed {removed} duplicates ({len(unique)} unique lines):\n\n" + "\n".join(unique)

        elif action == "base64_encode":
            import base64
            encoded = base64.b64encode(text.encode("utf-8")).decode("ascii")
            return encoded

        elif action == "base64_decode":
            import base64
            decoded = base64.b64decode(text).decode("utf-8", errors="replace")
            return decoded

        elif action == "url_encode":
            import urllib.parse
            return urllib.parse.quote(text, safe="")

        elif action == "url_decode":
            import urllib.parse
            return urllib.parse.unquote(text)

        elif action == "hash":
            import hashlib
            algo = algorithm.lower()
            if algo == "md5":
                h = hashlib.md5(text.encode("utf-8")).hexdigest()
            elif algo == "sha1":
                h = hashlib.sha1(text.encode("utf-8")).hexdigest()
            elif algo == "sha256":
                h = hashlib.sha256(text.encode("utf-8")).hexdigest()
            else:
                return f"Error: Unknown algorithm '{algo}'. Use: md5, sha1, sha256"
            return f"{algo}: {h}"

        elif action == "head":
            lines = text.splitlines()
            n = max(1, min(n, len(lines)))
            return "\n".join(lines[:n])

        elif action == "tail":
            lines = text.splitlines()
            n = max(1, min(n, len(lines)))
            return "\n".join(lines[-n:])

        elif action == "split":
            parts = text.split(delimiter)
            return f"Split into {len(parts)} parts:\n" + "\n---\n".join(parts[:100])

        else:
            return (f"Error: Unknown action '{action}'. Use: count, replace, extract, sort, dedup, "
                    "base64_encode, base64_decode, url_encode, url_decode, hash, head, tail, split")
    except Exception as e:
        return f"Error: {e}"


def _tool_save_experience(
    scene: str,
    core_knowledge: str,
    action_rules: list[str] | None = None,
    taboo_rules: list[str] | None = None,
    priority: str = "medium",
    tags: list[str] | None = None,
    exp_type: str = "retrospective",
    source: str = "",
    role: str = "",
    **ctx: Any,
) -> str:
    """Persist an experience entry into the role-based experience library.

    Skills (installable capability packages) are managed by the Skill Registry and
    are NOT written through this tool. Experiences are short, scene-anchored
    lessons that get auto-injected into the relevant role's system prompt.
    """
    try:
        if not scene or not core_knowledge:
            return "Error: 'scene' and 'core_knowledge' are required"

        # Resolve role: explicit arg > calling agent's role > 'default'
        resolved_role = (role or "").strip()
        if not resolved_role:
            try:
                caller_id = ctx.get("_caller_agent_id", "") if isinstance(ctx, dict) else ""
                if caller_id:
                    hub = _get_hub()
                    agent = hub.get_agent(caller_id) if hub else None
                    if agent is not None:
                        resolved_role = (getattr(agent, "role", "") or "").strip()
            except Exception:
                pass
        if not resolved_role:
            resolved_role = "default"

        # Validate priority
        pri = (priority or "medium").strip().lower()
        if pri not in ("high", "medium", "low"):
            pri = "medium"

        # Validate exp_type
        etype = (exp_type or "retrospective").strip().lower()
        if etype not in ("retrospective", "active_learning"):
            etype = "retrospective"

        from .experience_library import get_experience_library, Experience

        lib = get_experience_library()
        exp = Experience(
            exp_type=etype,
            source=source or "agent.save_experience",
            scene=scene.strip(),
            core_knowledge=core_knowledge.strip(),
            action_rules=[str(r).strip() for r in (action_rules or []) if str(r).strip()],
            taboo_rules=[str(r).strip() for r in (taboo_rules or []) if str(r).strip()],
            priority=pri,
            tags=[str(t).strip() for t in (tags or []) if str(t).strip()],
        )
        saved = lib.add_experience(resolved_role, exp)
        return (
            f"✓ Experience saved: id={saved.id} role={resolved_role} "
            f"priority={saved.priority} scene={saved.scene[:60]}"
        )
    except Exception as e:
        return f"Error saving experience: {e}"


def _tool_knowledge_lookup(query: str = "", entry_id: str = "",
                           agent_id: str = "", **kw: Any) -> str:
    """Look up entries in the knowledge base.

    Routing is determined by the agent's rag_mode:
      - "shared"  → query global shared knowledge (default)
      - "private" → query agent's private collection
      - "both"    → query private first, then shared
      - "none"    → return empty

    If entry_id is provided, returns that entry's full content from shared KB.
    Otherwise searches by query using the agent's configured RAG routing.
    """
    try:
        # If entry_id is provided, fetch that entry directly from shared KB
        if entry_id:
            entry = _knowledge.get_entry(entry_id)
            if entry:
                return json.dumps({
                    "status": "success",
                    "entry": entry
                }, ensure_ascii=False, indent=2)
            else:
                return json.dumps({
                    "status": "error",
                    "message": f"Entry '{entry_id}' not found"
                })

        if not query:
            return json.dumps({
                "status": "error",
                "message": "Either 'query' or 'entry_id' must be provided"
            })

        # --- RAG-routed search ---
        # Try to get agent profile for rag_mode routing
        agent_profile = kw.get("_agent_profile")
        rag_mode = "shared"  # default
        if agent_profile:
            rag_mode = getattr(agent_profile, "rag_mode", "shared") or "shared"

        if rag_mode == "none":
            return json.dumps({
                "status": "not_found",
                "message": "RAG is disabled for this agent (rag_mode=none)"
            })

        results_combined = []

        # Private / Both: search via RAG provider registry
        if rag_mode in ("private", "both") and agent_id:
            try:
                from .rag_provider import search_for_agent
                rag_results = search_for_agent(agent_profile, query,
                                               agent_id=agent_id, top_k=5)
                for r in rag_results:
                    results_combined.append({
                        "id": r.get("id", ""),
                        "title": r.get("title", ""),
                        "content": r.get("content", ""),
                        "tags": r.get("metadata", {}).get("tags", "").split(",")
                               if r.get("metadata", {}).get("tags") else [],
                        "source": "private_rag",
                    })
            except Exception as e:
                logger.warning("RAG provider search failed: %s", e)

        # Shared / Both: also search the classic shared knowledge base
        if rag_mode in ("shared", "both"):
            shared_results = _knowledge.search(query)
            for e in shared_results:
                # Avoid duplicates (same ID already from RAG)
                if not any(r["id"] == e["id"] for r in results_combined):
                    results_combined.append(e)

        if not results_combined:
            return json.dumps({
                "status": "not_found",
                "message": f"No knowledge entries found matching '{query}'"
            }, ensure_ascii=False)

        # Check for exact title match
        query_lower = query.lower().strip()
        for entry in results_combined:
            if entry.get("title", "").lower().strip() == query_lower:
                return json.dumps({
                    "status": "success",
                    "entry": entry
                }, ensure_ascii=False, indent=2)

        # No exact match — return list for refinement
        matches = [
            {"id": e.get("id", ""), "title": e.get("title", ""),
             "tags": e.get("tags", []),
             "source": e.get("source", "shared")}
            for e in results_combined[:20]
        ]

        return json.dumps({
            "status": "partial",
            "message": f"Found {len(matches)} matching entries. Use entry_id to read full content.",
            "matches": matches
        }, ensure_ascii=False, indent=2)

    except Exception as e:
        return json.dumps({
            "status": "error",
            "message": f"Error querying knowledge base: {str(e)}"
        })


def _tool_share_knowledge(title: str, content: str,
                          tags: list[str] | None = None, **ctx: Any) -> str:
    """Share knowledge with all agents via the shared Knowledge Base."""
    try:
        if not title or not title.strip():
            return "Error: 'title' is required"
        if not content or not content.strip():
            return "Error: 'content' is required"

        # Resolve caller agent info for attribution
        caller_name = ""
        caller_role = ""
        try:
            caller_id = ctx.get("_caller_agent_id", "") if isinstance(ctx, dict) else ""
            if caller_id:
                hub = _get_hub()
                agent = hub.get_agent(caller_id) if hub else None
                if agent is not None:
                    caller_name = getattr(agent, "name", "") or ""
                    caller_role = getattr(agent, "role", "") or ""
        except Exception:
            pass

        # Add source attribution
        source_info = ""
        if caller_name or caller_role:
            source_info = f"\n\n---\nShared by: {caller_name} (role: {caller_role})"

        resolved_tags = [str(t).strip() for t in (tags or []) if str(t).strip()]
        if caller_role:
            resolved_tags += ["shared-by-agent", caller_role]

        entry = _knowledge.add_entry(
            title=title.strip(),
            content=content.strip() + source_info,
            tags=resolved_tags,
        )
        return (
            f"Knowledge shared successfully: '{title.strip()}' (id: {entry['id']}). "
            f"All agents can now access this via knowledge_lookup."
        )
    except Exception as e:
        return f"Failed to share knowledge: {e}"


def _tool_learn_from_peers(source_role: str, topic: str = "",
                           limit: int = 5, **ctx: Any) -> str:
    """Learn from other agents' experiences by importing from another role."""
    try:
        if not source_role or not source_role.strip():
            return "Error: 'source_role' is required"

        source_role = source_role.strip()
        limit = max(1, min(int(limit), 20))

        # Resolve caller's role
        caller_role = ""
        try:
            caller_id = ctx.get("_caller_agent_id", "") if isinstance(ctx, dict) else ""
            if caller_id:
                hub = _get_hub()
                agent = hub.get_agent(caller_id) if hub else None
                if agent is not None:
                    caller_role = getattr(agent, "role", "") or ""
        except Exception:
            pass

        from .experience_library import _get_global_library
        library = _get_global_library()
        experiences = library.import_cross_role(
            source_role=source_role,
            target_role=caller_role or "default",
            topic=topic.strip() if topic else "",
            limit=limit,
        )
        if not experiences:
            msg = f"No experiences found for role '{source_role}'"
            if topic:
                msg += f" on topic '{topic}'"
            return msg

        lines = [f"Imported {len(experiences)} experiences from role '{source_role}':"]
        for i, exp in enumerate(experiences, 1):
            lines.append(f"\n{i}. [{exp.priority}] {exp.scene}")
            lines.append(f"   Knowledge: {exp.core_knowledge[:200]}")
            if exp.action_rules:
                lines.append(f"   Rules: {'; '.join(exp.action_rules[:3])}")
            lines.append(f"   Success rate: {exp.success_rate:.0%}")
        return "\n".join(lines)
    except Exception as e:
        return f"Failed to learn from peers: {e}"


def _tool_pip_install(packages: str, upgrade: bool = False, **_: Any) -> str:
    """Install or upgrade Python packages using pip."""
    import sys

    if not packages or not packages.strip():
        return "Error: packages parameter is required"

    try:
        pkg_list = packages.split()
        cmd = [sys.executable, "-m", "pip", "install"]
        if upgrade:
            cmd.append("--upgrade")
        cmd.extend(pkg_list)
        cmd.append("--break-system-packages")

        result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)

        if result.returncode == 0:
            return f"✓ Successfully installed: {', '.join(pkg_list)}"
        else:
            return f"Error installing packages: {result.stderr}"
    except Exception as e:
        return f"Error: {e}"


def _tool_create_pptx(output_path: str, slides: list, title: str = "", **_: Any) -> str:
    """Create a PowerPoint presentation file."""
    import sys

    try:
        # Auto-install python-pptx if not available
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "install", "python-pptx", "--break-system-packages"],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode != 0:
                return f"Error installing python-pptx: {result.stderr}"
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN

        # Validate path
        pol = _sandbox.get_current_policy()
        output_file = pol.safe_path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Create presentation
        prs = Presentation()

        # Set title if provided
        if title:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title

        # Add slides
        for slide_data in slides:
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", "")
            layout_type = slide_data.get("layout", "title_content").lower()

            # Select layout
            if layout_type == "title":
                layout = prs.slide_layouts[0]
            elif layout_type == "content":
                layout = prs.slide_layouts[5]  # Blank with title
            elif layout_type == "blank":
                layout = prs.slide_layouts[6]  # Blank
            else:  # title_content
                layout = prs.slide_layouts[1]  # Title and content

            slide = prs.slides.add_slide(layout)

            # Add title
            if slide_title and len(slide.shapes) > 0:
                title_shape = slide.shapes.title
                if title_shape:
                    title_shape.text = slide_title

            # Add content
            if content and len(slide.shapes) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()

                # Split content by lines and add as bullet points
                for line in content.split('\n'):
                    if line.strip():
                        p = tf.add_paragraph()
                        p.text = line.strip()
                        p.level = 0

            # Add images onto the slide
            for img_spec in (slide_data.get("images") or []):
                if not isinstance(img_spec, dict):
                    continue
                img_path_raw = img_spec.get("path", "")
                if not img_path_raw:
                    continue
                img_file = pol.safe_path(img_path_raw)
                left_v = Inches(float(img_spec.get("left", 1)))
                top_v = Inches(float(img_spec.get("top", 2)))
                kw: dict = {}
                if img_spec.get("width"):
                    kw["width"] = Inches(float(img_spec["width"]))
                if img_spec.get("height"):
                    kw["height"] = Inches(float(img_spec["height"]))
                slide.shapes.add_picture(str(img_file), left_v, top_v, **kw)

        # Save presentation
        prs.save(str(output_file))
        return f"✓ Created presentation: {output_file}"
    except Exception as e:
        return f"Error creating presentation: {e}"


def _tool_create_pptx_advanced(
    output_path: str,
    slides: list,
    theme: dict | None = None,
    **_: Any,
) -> str:
    """Create an advanced PowerPoint with shapes, charts, tables, and infographics."""
    import sys

    try:
        # Auto-install python-pptx if needed
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
        except ImportError:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "install", "python-pptx", "--break-system-packages"],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode != 0:
                return f"Error installing python-pptx: {result.stderr}"
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor

        # Validate path
        pol = _sandbox.get_current_policy()
        output_file = pol.safe_path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Theme defaults
        th = theme or {}
        T_PRIMARY = th.get("primary", "E8590C")
        T_SECONDARY = th.get("secondary", "2B2B2B")
        T_ACCENT = th.get("accent", "F4A261")
        T_BG = th.get("background", "FFFFFF")
        T_TITLE_FONT = th.get("title_font", "Microsoft YaHei")
        T_BODY_FONT = th.get("body_font", "Microsoft YaHei")

        def _rgb(hex_str: str) -> RGBColor:
            """Convert hex string to RGBColor."""
            h = hex_str.lstrip("#")
            if len(h) == 6:
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            return RGBColor(0, 0, 0)

        def _resolve_color(val: str) -> str:
            """Replace theme placeholders like 'primary' with actual hex."""
            m = {"primary": T_PRIMARY, "secondary": T_SECONDARY,
                 "accent": T_ACCENT, "background": T_BG}
            return m.get(val, val)

        # Shape type mapping
        SHAPE_MAP = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
            "chevron": MSO_SHAPE.CHEVRON,
            "diamond": MSO_SHAPE.DIAMOND,
            "pentagon": MSO_SHAPE.REGULAR_PENTAGON,
            "hexagon": MSO_SHAPE.HEXAGON,
            "star": MSO_SHAPE.STAR_5_POINT,
        }

        # Chart type mapping
        CHART_MAP = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "radar": XL_CHART_TYPE.RADAR,
            "area": XL_CHART_TYPE.AREA,
        }

        # Alignment mapping
        ALIGN_MAP = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }
        VALIGN_MAP = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }

        # Create presentation (16:9 widescreen)
        prs = Presentation()
        SLIDE_W = 10.0   # inches
        SLIDE_H = 5.625  # inches
        MARGIN = 0.15     # minimum margin from edge
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)

        def _clamp_bounds(el: dict) -> dict:
            """Clamp element x/y/w/h so it stays within slide boundaries.

            Fixes the most common LLM layout mistake: placing 3+ items
            in a row where the last one overflows the right/bottom edge.
            """
            x = float(el.get("x", 0))
            y = float(el.get("y", 0))
            w = float(el.get("w", 1))
            h = float(el.get("h", 1))

            # Clamp negative positions
            if x < 0:
                x = 0
            if y < 0:
                y = 0

            max_w = SLIDE_W - MARGIN
            max_h = SLIDE_H - MARGIN

            # If right edge overflows, shrink width first; if still bad, shift left
            if x + w > max_w:
                # Try shrinking width (keep at least 30% of original)
                new_w = max_w - x
                if new_w >= w * 0.3 and new_w > 0.3:
                    w = new_w
                else:
                    # Shift left to fit
                    x = max(0, max_w - w)
                    if x + w > max_w:
                        w = max_w - x

            # Same for bottom edge
            if y + h > max_h:
                new_h = max_h - y
                if new_h >= h * 0.3 and new_h > 0.3:
                    h = new_h
                else:
                    y = max(0, max_h - h)
                    if y + h > max_h:
                        h = max_h - y

            el["x"] = round(x, 3)
            el["y"] = round(y, 3)
            el["w"] = round(w, 3)
            el["h"] = round(h, 3)
            return el

        def _add_text_element(slide, el):
            """Add a text box to the slide."""
            left = Inches(el.get("x", 0))
            top = Inches(el.get("y", 0))
            width = Inches(el.get("w", 8))
            height = Inches(el.get("h", 1))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            # Background color
            bg = el.get("bg_color", "")
            if bg:
                fill = txBox.fill
                fill.solid()
                fill.fore_color.rgb = _rgb(_resolve_color(bg))

            # Vertical alignment
            va = el.get("valign", "")
            if va in VALIGN_MAP:
                tf.paragraphs[0].alignment  # ensure exists
                txBox.text_frame._txBody.attrib  # access
                try:
                    tf._txBody[0].attrib  # bodyPr
                except Exception:
                    pass
                # Set via the text frame directly
                from pptx.oxml.ns import qn
                bodyPr = tf._txBody.find(qn("a:bodyPr"))
                if bodyPr is not None:
                    anchor_val = {"top": "t", "middle": "ctr", "bottom": "b"}.get(va, "t")
                    bodyPr.set("anchor", anchor_val)

            content = el.get("content", "")
            lines = content.split("\\n") if "\\n" in content else content.split("\n")

            font_size = el.get("font_size", 14)
            font_name = el.get("font_name", T_BODY_FONT)
            bold = el.get("bold", False)
            italic = el.get("italic", False)
            color = _resolve_color(el.get("color", T_SECONDARY))
            align = el.get("align", "left")
            line_spacing = el.get("line_spacing", 0)

            for i, line_text in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line_text
                p.font.size = Pt(font_size)
                p.font.name = font_name
                p.font.bold = bold
                p.font.italic = italic
                if color:
                    p.font.color.rgb = _rgb(color)
                if align in ALIGN_MAP:
                    p.alignment = ALIGN_MAP[align]
                if line_spacing and line_spacing > 0:
                    p.line_spacing = Pt(font_size * line_spacing)

        def _add_shape_element(slide, el):
            """Add a shape to the slide."""
            shape_type_name = el.get("shape_type", "rectangle")
            mso_shape = SHAPE_MAP.get(shape_type_name, MSO_SHAPE.RECTANGLE)
            left = Inches(el.get("x", 0))
            top = Inches(el.get("y", 0))
            width = Inches(el.get("w", 1))
            height = Inches(el.get("h", 1))

            shape = slide.shapes.add_shape(mso_shape, left, top, width, height)

            fill_color = el.get("fill_color", "")
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = _rgb(_resolve_color(fill_color))
            else:
                shape.fill.background()  # transparent

            line_color = el.get("line_color", "")
            line_width = el.get("line_width", 0)
            if line_color:
                shape.line.color.rgb = _rgb(_resolve_color(line_color))
                shape.line.width = Pt(line_width or 1)
            else:
                shape.line.fill.background()  # no border

            rotation = el.get("rotation", 0)
            if rotation:
                shape.rotation = rotation

        def _add_line_element(slide, el):
            """Add a line connector."""
            x = Inches(el.get("x", 0))
            y = Inches(el.get("y", 0))
            w = Inches(el.get("w", 1))
            h = Inches(el.get("h", 0))

            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                x, y, x + w, y + h,
            )
            lc = el.get("line_color", "CCCCCC")
            if lc:
                connector.line.color.rgb = _rgb(_resolve_color(lc))
            lw = el.get("line_width", 1)
            connector.line.width = Pt(lw)

        def _add_chart_element(slide, el):
            """Add a chart to the slide."""
            from pptx.chart.data import CategoryChartData

            chart_type_name = el.get("chart_type", "column")
            xl_chart = CHART_MAP.get(chart_type_name, XL_CHART_TYPE.COLUMN_CLUSTERED)

            x = Inches(el.get("x", 0.5))
            y = Inches(el.get("y", 1.5))
            w = Inches(el.get("w", 5))
            h = Inches(el.get("h", 3.5))

            chart_data = CategoryChartData()
            categories = el.get("categories", [])
            chart_data.categories = categories

            for s in (el.get("series") or []):
                chart_data.add_series(
                    s.get("name", "Series"),
                    s.get("values", []),
                )

            chart_frame = slide.shapes.add_chart(xl_chart, x, y, w, h, chart_data)
            chart = chart_frame.chart

            # Apply custom colors
            colors_list = el.get("colors", [])
            if colors_list:
                try:
                    plot = chart.plots[0]
                    if chart_type_name in ("pie", "doughnut"):
                        # For pie/doughnut, color each point
                        if plot.series and len(plot.series) > 0:
                            series_obj = plot.series[0]
                            for idx, c in enumerate(colors_list):
                                if idx < len(categories):
                                    point = series_obj.points[idx]
                                    point.format.fill.solid()
                                    point.format.fill.fore_color.rgb = _rgb(_resolve_color(c))
                    else:
                        # For bar/column/line, color each series
                        for idx, c in enumerate(colors_list):
                            if idx < len(plot.series):
                                s = plot.series[idx]
                                s.format.fill.solid()
                                s.format.fill.fore_color.rgb = _rgb(_resolve_color(c))
                except Exception:
                    pass  # color fail is non-critical

            # Labels
            if el.get("show_labels") or el.get("show_percent"):
                try:
                    plot = chart.plots[0]
                    plot.has_data_labels = True
                    data_labels = plot.data_labels
                    if el.get("show_percent") and chart_type_name in ("pie", "doughnut"):
                        data_labels.show_percentage = True
                        data_labels.show_value = False
                    else:
                        data_labels.show_value = True
                    data_labels.font.size = Pt(10)
                except Exception:
                    pass

            # Legend
            if el.get("show_legend") is False:
                chart.has_legend = False
            elif el.get("show_legend"):
                chart.has_legend = True

        def _add_table_element(slide, el):
            """Add a table to the slide."""
            headers = el.get("headers", [])
            rows_data = el.get("rows", [])
            total_rows = len(rows_data) + (1 if headers else 0)
            total_cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 1)

            x = Inches(el.get("x", 0.5))
            y = Inches(el.get("y", 1.5))
            w = Inches(el.get("w", 9))
            h = Inches(el.get("h", 3))

            table_shape = slide.shapes.add_table(total_rows, total_cols, x, y, w, h)
            table = table_shape.table

            # Set column widths evenly
            col_w = Emu(int(Inches(el.get("w", 9)) / total_cols))
            for ci in range(total_cols):
                table.columns[ci].width = col_w

            header_bg = el.get("header_color", T_PRIMARY)
            header_fc = el.get("header_font_color", "FFFFFF")
            stripe_bg = el.get("stripe_color", "")

            row_offset = 0
            if headers:
                for ci, htext in enumerate(headers):
                    cell = table.cell(0, ci)
                    cell.text = str(htext)
                    # Header style
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(_resolve_color(header_bg))
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = _rgb(_resolve_color(header_fc))
                        paragraph.font.name = T_BODY_FONT
                        paragraph.alignment = PP_ALIGN.CENTER
                row_offset = 1

            for ri, row in enumerate(rows_data):
                for ci, cval in enumerate(row):
                    if ci >= total_cols:
                        break
                    cell = table.cell(ri + row_offset, ci)
                    cell.text = str(cval)
                    # Stripe
                    if stripe_bg and ri % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(_resolve_color(stripe_bg))
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(11)
                        paragraph.font.name = T_BODY_FONT
                        paragraph.alignment = PP_ALIGN.CENTER

        def _add_icon_circle(slide, el):
            """Add a circle with text inside (for numbering, icons, etc.)."""
            x = Inches(el.get("x", 0))
            y = Inches(el.get("y", 0))
            w = Inches(el.get("w", 0.8))
            h = Inches(el.get("h", 0.8))

            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
            fill_color = el.get("fill_color", T_PRIMARY)
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(_resolve_color(fill_color))
            shape.line.fill.background()  # no border

            # Add text inside
            tf = shape.text_frame
            tf.word_wrap = False
            from pptx.oxml.ns import qn
            bodyPr = tf._txBody.find(qn("a:bodyPr"))
            if bodyPr is not None:
                bodyPr.set("anchor", "ctr")

            p = tf.paragraphs[0]
            p.text = el.get("text", "")
            p.font.size = Pt(el.get("font_size", 16))
            p.font.bold = True
            p.font.color.rgb = _rgb(_resolve_color(el.get("font_color", "FFFFFF")))
            p.font.name = T_TITLE_FONT
            p.alignment = PP_ALIGN.CENTER

        def _add_image_element(slide, el):
            """Add an image to the slide."""
            img_path = el.get("path", "")
            if not img_path:
                return
            img_file = pol.safe_path(img_path)
            if not img_file.exists():
                return
            x = Inches(el.get("x", 0))
            y = Inches(el.get("y", 0))
            kw = {}
            if el.get("w"):
                kw["width"] = Inches(el["w"])
            if el.get("h"):
                kw["height"] = Inches(el["h"])
            slide.shapes.add_picture(str(img_file), x, y, **kw)

        # Element dispatcher
        ELEMENT_HANDLERS = {
            "text": _add_text_element,
            "shape": _add_shape_element,
            "line": _add_line_element,
            "chart": _add_chart_element,
            "table": _add_table_element,
            "icon_circle": _add_icon_circle,
            "image": _add_image_element,
        }

        # Layout engine integration
        try:
            from .utils.pptx_layouts import generate_layout
        except ImportError:
            generate_layout = None

        # Build slides
        blank_layout = prs.slide_layouts[6]  # Blank
        slide_count = 0
        theme_dict = {
            "primary": T_PRIMARY, "secondary": T_SECONDARY,
            "accent": T_ACCENT, "background": T_BG,
            "title_font": T_TITLE_FONT, "body_font": T_BODY_FONT,
        }

        for slide_data in slides:
            slide = prs.slides.add_slide(blank_layout)
            slide_count += 1

            # Slide background
            bg_color = slide_data.get("background", T_BG)
            if bg_color:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = _rgb(_resolve_color(bg_color))

            # If layout spec is present, auto-generate elements from it
            all_elements = []
            layout_spec = slide_data.get("layout")
            if layout_spec and generate_layout:
                try:
                    auto_els = generate_layout(layout_spec, theme_dict)
                    all_elements.extend(auto_els)
                except Exception as _le:
                    import sys as _sys
                    print(f"[pptx_advanced] layout error: {_le}", file=_sys.stderr)

            # Append any manually-specified elements (can supplement layout)
            all_elements.extend(slide_data.get("elements") or [])

            # Add elements in order (z-order: first = bottom)
            for el in all_elements:
                el_type = el.get("type", "")
                handler = ELEMENT_HANDLERS.get(el_type)
                if handler:
                    try:
                        # Clamp bounds to prevent overflow
                        if any(k in el for k in ("x", "y", "w", "h")):
                            _clamp_bounds(el)
                        handler(slide, el)
                    except Exception as e:
                        # Non-critical: log but continue
                        import sys as _sys
                        print(f"[pptx_advanced] element error ({el_type}): {e}",
                              file=_sys.stderr)

        prs.save(str(output_file))
        return f"✓ Created advanced presentation ({slide_count} slides): {output_file}"

    except Exception as e:
        return f"Error creating presentation: {e}"


def _tool_desktop_screenshot(output_path: str = "", region: dict | None = None, **_: Any) -> str:
    """Take a screenshot of the desktop."""
    try:
        from datetime import datetime
        import os

        # Default output path
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"screenshot_{timestamp}.png"

        # Validate output path
        pol = _sandbox.get_current_policy()
        output_file = pol.safe_path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Try mss first (cross-platform)
        try:
            import mss
            with mss.mss() as sct:
                monitor = sct.monitors[1]  # Primary monitor

                if region:
                    screenshot = sct.grab({
                        'left': region.get('x', 0),
                        'top': region.get('y', 0),
                        'width': region.get('w', monitor['width']),
                        'height': region.get('h', monitor['height'])
                    })
                else:
                    screenshot = sct.grab(monitor)

                import mss.tools
                mss.tools.to_png(screenshot.rgb, screenshot.size, output=str(output_file))
                return f"✓ Screenshot saved: {output_path}"
        except ImportError:
            pass

        # Try PIL/Pillow ImageGrab
        try:
            from PIL import ImageGrab

            if region:
                bbox = (region.get('x', 0), region.get('y', 0),
                        region.get('x', 0) + region.get('w', 1920),
                        region.get('y', 0) + region.get('h', 1080))
                img = ImageGrab.grab(bbox=bbox)
            else:
                img = ImageGrab.grab()

            img.save(str(output_file), 'PNG')
            return f"✓ Screenshot saved: {output_path}"
        except ImportError:
            pass

        # Fallback to platform-specific commands
        if os.name == 'posix':
            # Unix/Linux - try scrot
            cmd = ["scrot", str(output_file)]
            result = subprocess.run(cmd, capture_output=True, timeout=10)
            if result.returncode == 0:
                return f"✓ Screenshot saved: {output_path}"

            # macOS - try screencapture
            cmd = ["screencapture", "-x", str(output_file)]
            result = subprocess.run(cmd, capture_output=True, timeout=10)
            if result.returncode == 0:
                return f"✓ Screenshot saved: {output_path}"

        return "Error: Could not take screenshot (mss, PIL, scrot, or screencapture required)"
    except Exception as e:
        return f"Error taking screenshot: {e}"


def _tool_create_video(output_path: str, frames: list, fps: int = 24, audio_path: str = "", **_: Any) -> str:
    """Create a video from image frames."""
    import sys

    try:
        # Validate output path
        pol = _sandbox.get_current_policy()
        output_file = pol.safe_path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)

        # Validate frame paths
        frame_list = []
        for frame in frames:
            img_path = frame.get("image_path", "")
            if not img_path:
                return "Error: Each frame must have image_path"
            img_file = pol.safe_path(img_path)
            if not img_file.exists():
                return f"Error: Image file not found: {img_path}"
            duration = frame.get("duration", 3)
            frame_list.append((str(img_file), duration))

        # Try moviepy first
        try:
            import moviepy.editor as mpy
        except ImportError:
            result = subprocess.run(
                [sys.executable, "-m", "pip", "install", "moviepy", "--break-system-packages"],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode != 0:
                return f"Error installing moviepy: {result.stderr}"
            import moviepy.editor as mpy

        # Create video from frames
        clips = []
        for img_path, duration in frame_list:
            clip = mpy.ImageClip(img_path).set_duration(duration)
            clips.append(clip)

        video = mpy.concatenate_videoclips(clips)

        # Add audio if provided
        if audio_path:
            audio_file = pol.safe_path(audio_path)
            if audio_file.exists():
                audio = mpy.AudioFileClip(str(audio_file))
                video = video.set_audio(audio)

        # Write video file
        video.write_videofile(str(output_file), fps=fps, verbose=False, logger=None)
        video.close()

        return f"✓ Video created: {output_path}"
    except Exception as e:
        return f"Error creating video: {e}"


# ---------------------------------------------------------------------------
# Skill guide on-demand loader
# ---------------------------------------------------------------------------

def _tool_get_skill_guide(arguments: dict) -> str:
    """Load the full SKILL.md guide + ancillary file list for a granted skill.

    Returns the complete instructions (with frontmatter stripped) so the
    agent can follow the step-by-step guide and run scripts from skill_dir.
    """
    name = (arguments.get("name") or "").strip()
    if not name:
        return "Error: name is required"

    try:
        import sys as _sys
        reg = None
        # Try hub first (hot path in production)
        _llm_mod = _sys.modules.get("app.llm")
        hub = getattr(_llm_mod, "_active_hub", None) if _llm_mod else None
        reg = getattr(hub, "skill_registry", None) if hub else None
        if reg is None:
            # Fallback: module-level singleton
            from .skills.engine import get_registry
            reg = get_registry()
        if reg is None:
            # Last resort: check if skill_store has a registry
            try:
                from . import skill_store as _ss
                store = _ss.get_store()
                if store and store._registry:
                    reg = store._registry
            except Exception:
                pass
        if reg is None:
            return "Error: skill registry not available (hub not started)"

        # Find the skill by name (fuzzy: accept name, id, or name@version)
        found = None
        for inst in reg.list_all():
            if inst.manifest.name == name or inst.id == name:
                found = inst
                break
            if name in inst.id:
                found = inst
        if found is None:
            available = [i.manifest.name for i in reg.list_all()]
            return f"Error: skill '{name}' not found. Available: {', '.join(available)}"

        # Determine effective skill_dir: prefer agent-local workspace copy
        # The caller may pass agent_id so we can look up the agent's workspace.
        install_dir = found.install_dir
        agent_id = (arguments.get("agent_id") or "").strip()
        effective_dir = install_dir

        if agent_id:
            try:
                _llm_mod2 = _sys.modules.get("app.llm")
                hub2 = getattr(_llm_mod2, "_active_hub", None) if _llm_mod2 else None
                agent_obj = hub2.get_agent(agent_id) if hub2 and hasattr(hub2, "get_agent") else None
                if agent_obj and hasattr(agent_obj, "get_skill_workspace_dir"):
                    local_dir = agent_obj.get_skill_workspace_dir(found.manifest.name)
                    if local_dir:
                        effective_dir = str(local_dir)
            except Exception:
                pass

        entry_file = found.manifest.entry or "SKILL.md"

        # Read the SKILL.md body (strip frontmatter) — prefer agent-local copy
        import re as _re
        from pathlib import Path as _Path
        md_path = _Path(effective_dir) / entry_file
        if not md_path.exists():
            # Fallback to global install_dir if agent-local copy missing
            md_path = _Path(install_dir) / entry_file
        body = ""
        if md_path.exists():
            text = md_path.read_text(encoding="utf-8")
            fm = _re.match(r"^---\s*\n.*?\n---\s*\n?", text, _re.DOTALL)
            body = text[fm.end():] if fm else text

        # List ancillary files (scripts, references, etc.)
        files = []
        base = _Path(effective_dir)
        if not base.is_dir():
            base = _Path(install_dir)
        for fp in sorted(base.rglob("*")):
            if fp.is_file() and fp.name.lower() != "skill.md":
                try:
                    rel = str(fp.relative_to(base))
                except ValueError:
                    rel = fp.name
                files.append(rel)

        # Also list reference .md files whose content may be needed
        ref_mds = []
        for fp in base.glob("*.md"):
            if fp.name.lower() != "skill.md":
                ref_mds.append(fp.name)

        result_parts = [
            f"## Skill: {found.manifest.name}",
            f"**skill_dir**: `{effective_dir}`",
            f"**runtime**: {found.manifest.runtime}",
            "",
            "运行脚本时先 cd 到 skill_dir:",
            f"```bash",
            f"cd {effective_dir}",
            f"```",
            "",
        ]
        if files:
            result_parts.append("**附属文件**: " + ", ".join(files))
            result_parts.append("")
        if ref_mds:
            result_parts.append("**参考文档** (需要时用 read_file 读取): " +
                                ", ".join(f"`{effective_dir}/{m}`" for m in ref_mds))
            result_parts.append("")
        result_parts.append("---")
        result_parts.append("")
        result_parts.append(body)
        return "\n".join(result_parts)

    except Exception as e:
        return f"Error loading skill guide: {e}"


def _tool_propose_skill(role: str = "", topic: str = "", **ctx: Any) -> str:
    """Scan experience library and propose skill drafts via SkillForge.

    Returns a summary of generated drafts (pending admin approval).
    """
    try:
        from .skills._skill_forge import get_skill_forge

        # Resolve role from caller agent if not specified
        if not role:
            try:
                caller_id = ctx.get("_caller_agent_id", "") if isinstance(ctx, dict) else ""
                if caller_id:
                    hub = _get_hub()
                    agent = hub.get_agent(caller_id) if hub else None
                    if agent is not None:
                        role = (getattr(agent, "role", "") or "").strip()
            except Exception:
                pass

        forge = get_skill_forge()
        candidates = forge.scan_for_candidates(role=role or "")

        if not candidates:
            return (
                "未发现可以生成技能的经验模式。需要至少 3 个相似的高成功率经验。"
                "请继续积累经验（通过 save_experience 工具），之后再试。"
            )

        # Export packages for all candidates
        results = []
        for draft in candidates:
            try:
                export_dir = forge.export_package(draft)
                results.append(
                    f"✓ 技能草稿: {draft.name} (ID: {draft.id})\n"
                    f"  描述: {draft.description}\n"
                    f"  置信度: {draft.confidence:.0%}\n"
                    f"  来源经验: {len(draft.source_experiences)} 条\n"
                    f"  导出目录: {export_dir}\n"
                    f"  状态: 等待管理员审批"
                )
            except Exception as e:
                results.append(f"✗ 草稿 {draft.name} 导出失败: {e}")

        summary = (
            f"已生成 {len(candidates)} 个技能草稿，等待管理员在 Portal 审批：\n\n"
            + "\n\n".join(results)
        )
        return summary

    except Exception as e:
        return f"Error proposing skill: {e}"


def _tool_submit_skill(dir_name: str, **ctx: Any) -> str:
    """Submit a skill package from the agent's workspace for admin approval.

    Reads manifest.yaml, SKILL.md, and *.py from {workspace}/{dir_name},
    validates required manifest fields, creates a SkillDraft, and saves it
    to the SkillForge review queue.
    """
    import os
    import time as _time

    try:
        import yaml as _yaml
    except ImportError:
        return "Error: PyYAML not installed. Run pip install pyyaml."

    from .skills._skill_forge import get_skill_forge, SkillDraft

    # Resolve workspace directory from sandbox policy (agent's working dir)
    pol = _sandbox.get_current_policy()
    workspace = str(pol.root) if getattr(pol, "root", None) else None

    # Fallback: try to find workspace from caller agent
    if not workspace:
        try:
            caller_id = ctx.get("_caller_agent_id", "") if isinstance(ctx, dict) else ""
            if caller_id:
                hub = _get_hub()
                agent = hub.get_agent(caller_id) if hub else None
                if agent and hasattr(agent, "working_dir"):
                    workspace = str(agent.working_dir)
        except Exception:
            pass

    if not workspace:
        return "Error: Cannot determine workspace directory."

    skill_dir = os.path.join(workspace, dir_name)
    if not os.path.isdir(skill_dir):
        return f"Error: Directory not found: {skill_dir}"

    # Read manifest.yaml
    manifest_path = os.path.join(skill_dir, "manifest.yaml")
    if not os.path.isfile(manifest_path):
        return (
            "Error: manifest.yaml not found in skill directory. "
            "Please create manifest.yaml with required fields: "
            "name, version, description, runtime, author, entry"
        )

    manifest_yaml = open(manifest_path, "r", encoding="utf-8").read()
    try:
        m = _yaml.safe_load(manifest_yaml) or {}
    except Exception as e:
        return f"Error: Invalid YAML in manifest.yaml: {e}"

    # Validate required fields
    required = ["name", "version", "description", "runtime", "author", "entry"]
    missing = [f for f in required if not m.get(f)]
    if missing:
        return (
            f"Error: manifest.yaml missing required fields: {', '.join(missing)}. "
            "All of these are required: name, version, description, runtime, author, entry"
        )

    rt = m.get("runtime", "")
    if rt not in ("python", "shell", "markdown"):
        return f"Error: runtime must be 'python', 'shell', or 'markdown', got '{rt}'"

    # Read SKILL.md
    skill_md = ""
    skill_md_path = os.path.join(skill_dir, "SKILL.md")
    if os.path.isfile(skill_md_path):
        skill_md = open(skill_md_path, "r", encoding="utf-8").read()
    else:
        return (
            "Error: SKILL.md not found in skill directory. "
            "Please create SKILL.md documenting what the skill does and how to use it."
        )

    # Collect code files (*.py)
    code_files: dict[str, str] = {}
    for fn in os.listdir(skill_dir):
        fp = os.path.join(skill_dir, fn)
        if os.path.isfile(fp) and fn.endswith(".py"):
            try:
                code_files[fn] = open(fp, "r", encoding="utf-8").read()
            except Exception:
                pass

    # If runtime is python, entry file must exist
    entry = m.get("entry", "")
    if rt == "python" and entry.endswith(".py") and entry not in code_files:
        return f"Error: Entry file '{entry}' not found in skill directory."

    # Build description string
    desc = m.get("description", "")
    if isinstance(desc, dict):
        desc = desc.get("zh-CN") or desc.get("en") or str(desc)
    triggers = m.get("triggers", [])

    # Check for duplicate: same name + same version = reject
    forge = get_skill_forge()
    skill_name = m["name"]
    skill_version = m.get("version", "")
    for existing in forge._drafts.values():
        if existing.name == skill_name and existing.status in ("draft", "exported", "approved"):
            # Parse existing version from manifest
            existing_version = ""
            if existing.manifest_yaml:
                try:
                    em = _yaml.safe_load(existing.manifest_yaml) or {}
                    existing_version = em.get("version", "")
                except Exception:
                    pass
            if existing_version == skill_version:
                return (
                    f"Error: 技能 '{skill_name}' v{skill_version} 已存在"
                    f"（ID: {existing.id}, 状态: {existing.status}）。\n"
                    f"请修改 manifest.yaml 中的 version 字段后重新提交。"
                )

    draft_id = f"SF-{_time.strftime('%Y%m%d')}-SUB-{os.urandom(3).hex()}"
    draft = SkillDraft(
        id=draft_id,
        name=m["name"],
        description=str(desc),
        source_experiences=[],
        role=ctx.get("_caller_role", "") if isinstance(ctx, dict) else "",
        scene_pattern="",
        triggers=triggers if isinstance(triggers, list) else [triggers],
        manifest_yaml=manifest_yaml,
        skill_md=skill_md,
        confidence=0.95,
        created_at=_time.time(),
        status="exported",
        runtime=rt,
        code_files=code_files,
    )

    forge._drafts[draft_id] = draft
    forge._save_drafts()

    return (
        f"✓ 技能已提交审批！\n"
        f"  草稿 ID: {draft_id}\n"
        f"  名称: {m['name']}\n"
        f"  运行时: {rt}\n"
        f"  代码文件: {', '.join(code_files.keys()) or '(无)'}\n"
        f"  状态: 等待管理员在 Portal → 技能锻造 中审批\n\n"
        f"管理员审批通过后，技能将自动出现在 Skill Store 中。"
    )


# ---------------------------------------------------------------------------
# Tool dispatcher
# ---------------------------------------------------------------------------

_TOOL_FUNCS: dict[str, callable] = {
    "read_file": _tool_read_file,
    "write_file": _tool_write_file,
    "edit_file": _tool_edit_file,
    "bash": _tool_bash,
    "search_files": _tool_search_files,
    "glob_files": _tool_glob_files,
    "web_search": _tool_web_search,
    "web_fetch": _tool_web_fetch,
    # New daily-work tools
    "web_screenshot": _tool_web_screenshot,
    "http_request": _tool_http_request,
    "datetime_calc": _tool_datetime_calc,
    "json_process": _tool_json_process,
    "text_process": _tool_text_process,
    # Coordination tools
    "team_create": _tool_team_create,
    "send_message": _tool_send_message,
    "task_update": _tool_task_update,
    "mcp_call": _tool_mcp_call,
    # Experience persistence + skill generation
    "save_experience": _tool_save_experience,
    "propose_skill": _tool_propose_skill,
    "submit_skill": _tool_submit_skill,
    # Knowledge management tools
    "knowledge_lookup": _tool_knowledge_lookup,
    "share_knowledge": _tool_share_knowledge,
    "learn_from_peers": _tool_learn_from_peers,
    # Human-in-the-loop tools (handled specially by agent, not dispatched here)
    "request_web_login": lambda **kw: "ERROR: request_web_login must be handled by agent directly",
    # System and productivity tools
    "pip_install": _tool_pip_install,
    "create_pptx": _tool_create_pptx,
    "create_pptx_advanced": _tool_create_pptx_advanced,
    "desktop_screenshot": _tool_desktop_screenshot,
    "create_video": _tool_create_video,
    "get_skill_guide": _tool_get_skill_guide,
}


# Tool name aliases (LLMs sometimes call with different names)
_TOOL_ALIASES: dict[str, str] = {
    "exec": "bash",
    "execute": "bash",
    "shell": "bash",
    "run_command": "bash",
    "cmd": "bash",
    "run_bash": "bash",
    "read": "read_file",
    "write": "write_file",
    "edit": "edit_file",
    "search": "search_files",
    "grep": "search_files",
    "glob": "glob_files",
    "find": "glob_files",
    "fetch": "web_fetch",
    "fetch_url": "web_fetch",
    "screenshot": "web_screenshot",
    "capture": "web_screenshot",
    "http": "http_request",
    "request": "http_request",
    "api_call": "http_request",
    "curl": "http_request",
    "datetime": "datetime_calc",
    "date": "datetime_calc",
    "time": "datetime_calc",
    "json": "json_process",
    "parse_json": "json_process",
    "text": "text_process",
    "string": "text_process",
    "knowledge": "knowledge_lookup",
    "look_up_knowledge": "knowledge_lookup",
    "search_knowledge": "knowledge_lookup",
    "share": "share_knowledge",
    "publish_knowledge": "share_knowledge",
    "learn_peers": "learn_from_peers",
    "cross_role_learn": "learn_from_peers",
    "pip": "pip_install",
    "install": "pip_install",
    "pptx": "create_pptx",
    "pptx_advanced": "create_pptx_advanced",
    "advanced_pptx": "create_pptx_advanced",
    "powerpoint": "create_pptx",
    "presentation": "create_pptx",
    "screenshot": "desktop_screenshot",
    "snap": "desktop_screenshot",
    "screen_capture": "desktop_screenshot",
    "video": "create_video",
    "make_video": "create_video",
    "stitch_frames": "create_video",
    "skill_guide": "get_skill_guide",
    "load_skill": "get_skill_guide",
    "read_skill": "get_skill_guide",
    "generate_skill": "propose_skill",
    "create_skill": "propose_skill",
    "forge_skill": "propose_skill",
    "submit_skill_package": "submit_skill",
    "publish_skill": "submit_skill",
}


# ---------------------------------------------------------------------------
# ToolRegistry initialization
# ---------------------------------------------------------------------------

def _init_registry() -> ToolRegistry:
    """
    Initialize the module-level tool registry from existing TOOL_DEFINITIONS
    and _TOOL_FUNCS. This is called once to populate the singleton.
    """
    registry = ToolRegistry()

    # Map tool names to their toolset categories
    toolset_map = {
        # Core file operations
        "read_file": "core",
        "write_file": "core",
        "edit_file": "core",
        "bash": "core",
        "search_files": "core",
        "glob_files": "core",

        # Web tools
        "web_search": "web",
        "web_fetch": "web",
        "web_screenshot": "web",
        "http_request": "web",

        # Data processing
        "json_process": "data",
        "text_process": "data",
        "datetime_calc": "data",

        # Coordination / messaging
        "team_create": "coordination",
        "send_message": "coordination",
        "task_update": "coordination",
        "mcp_call": "coordination",

        # Skill management
        "save_experience": "coordination",
        "propose_skill": "skill",
        "submit_skill": "skill",

        # Knowledge management
        "knowledge_lookup": "coordination",
        "share_knowledge": "coordination",
        "learn_from_peers": "coordination",

        # Human-in-the-loop
        "request_web_login": "coordination",
        # System and productivity tools
        "pip_install": "system",
        "create_pptx": "productivity",
        "create_pptx_advanced": "productivity",
        "desktop_screenshot": "system",
        "create_video": "productivity",
        "get_skill_guide": "skill",
    }

    # Find tool schema definitions by name
    schema_map = {}
    for tool_def in TOOL_DEFINITIONS:
        if tool_def.get("type") == "function":
            tool_name = tool_def["function"].get("name")
            if tool_name:
                schema_map[tool_name] = tool_def["function"]

    # Register each tool from _TOOL_FUNCS
    for tool_name, handler in _TOOL_FUNCS.items():
        toolset = toolset_map.get(tool_name, "other")
        schema = schema_map.get(tool_name, {})
        description = schema.get("description", "")

        # Determine risk level
        if tool_name in ("bash", "write_file", "edit_file"):
            risk = "dangerous"
        elif tool_name in ("web_fetch", "web_search", "http_request", "pip_install"):
            risk = "moderate"
        else:
            risk = "safe"

        registry.register(
            name=tool_name,
            toolset=toolset,
            schema=schema,
            handler=handler,
            description=description,
            risk_level=risk,
        )

    # Register aliases
    for alias, canonical in _TOOL_ALIASES.items():
        try:
            registry.add_alias(alias, canonical)
        except ValueError:
            logger.warning(f"Failed to register alias '{alias}' → '{canonical}'")

    return registry


# Module-level singleton instance
tool_registry = _init_registry()


def execute_tool(name: str, arguments: dict) -> str:
    """
    Execute a tool by name with the given arguments.
    Delegates to tool_registry.dispatch() but maintains backward compatibility
    with existing code that calls execute_tool() directly.
    Returns the result string.
    """
    return tool_registry.dispatch(name, arguments)


def get_tool_definitions() -> list[dict]:
    """
    Return tool definitions in function-calling JSON schema format.
    Delegates to tool_registry.get_definitions() for available tools.
    """
    return tool_registry.get_definitions()
