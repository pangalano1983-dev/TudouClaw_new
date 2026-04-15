"""Extract text content from uploaded files for LLM consumption."""

import logging

logger = logging.getLogger(__name__)

MAX_CHARS = 50000

_TEXT_EXTS = frozenset({
    'txt', 'md', 'csv', 'tsv', 'json', 'yaml', 'yml', 'py', 'js',
    'ts', 'html', 'css', 'xml', 'log', 'ini', 'cfg', 'sh', 'bat',
})


def extract_file_text(file_path: str, mime: str = "") -> str:
    """Extract text content from a file for LLM consumption.

    Supports: PDF, DOCX, PPTX, XLSX/XLS, TXT, MD, CSV, JSON, YAML, HTML, code files.
    Returns extracted text or empty string on failure.
    """
    ext = file_path.rsplit('.', 1)[-1].lower() if '.' in file_path else ''

    # Text-based files: read directly
    if ext in _TEXT_EXTS:
        for enc in ('utf-8', 'gbk', 'gb2312', 'latin-1'):
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()[:MAX_CHARS]
            except (UnicodeDecodeError, Exception):
                continue
        return ""

    # PDF
    if ext == 'pdf':
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages[:100]:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
            return "\n\n".join(text_parts)[:MAX_CHARS]
        except Exception:
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(file_path)
                text_parts = [page.get_text() for page in doc[:100]]
                doc.close()
                return "\n\n".join(text_parts)[:MAX_CHARS]
            except Exception:
                return ""

    # DOCX
    if ext == 'docx':
        try:
            from docx import Document
            doc = Document(file_path)
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    text_parts.append(" | ".join(cells))
            return "\n".join(text_parts)[:MAX_CHARS]
        except Exception:
            return ""

    # PPTX — extract text + structure + styles so the LLM can learn
    # how to CREATE similar presentations, not just read text
    if ext == 'pptx':
        try:
            return _extract_pptx_rich(file_path)
        except Exception as e:
            logger.debug("PPTX rich extraction failed: %s", e)
            return ""

    # XLSX / XLS
    if ext in ('xlsx', 'xls'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for ws in wb.worksheets[:10]:
                rows = []
                for row in ws.iter_rows(max_row=500, values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    rows.append(" | ".join(cells))
                if rows:
                    text_parts.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
            wb.close()
            return "\n\n".join(text_parts)[:MAX_CHARS]
        except Exception:
            return ""

    return ""


# ---------------------------------------------------------------------------
# Rich PPTX extraction — captures structure + styles for template learning
# ---------------------------------------------------------------------------

def _emu_to_cm(emu):
    """Convert EMU (English Metric Units) to cm, rounded."""
    if emu is None:
        return None
    return round(emu / 914400 * 2.54, 1)


def _rgb_str(color):
    """Extract RGB hex string from a pptx color object."""
    try:
        if color and color.rgb:
            return f"#{color.rgb}"
    except Exception:
        pass
    return None


def _font_desc(font):
    """Describe a font object concisely."""
    parts = []
    if font.name:
        parts.append(font.name)
    if font.size:
        parts.append(f"{round(font.size.pt)}pt")
    if font.bold:
        parts.append("bold")
    if font.italic:
        parts.append("italic")
    c = _rgb_str(font.color)
    if c:
        parts.append(c)
    return ", ".join(parts) if parts else None


def _extract_pptx_rich(file_path: str) -> str:
    """Extract PPTX with full structure and style information.

    Output is structured text that lets an LLM understand the template's
    design system: layouts, fonts, colors, positions, so it can recreate
    similar presentations.
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []

    # ── Global info ──
    w = _emu_to_cm(prs.slide_width)
    h = _emu_to_cm(prs.slide_height)
    parts.append(f"[PPTX Template] Slide size: {w} x {h} cm, Slides: {len(prs.slides)}")

    # ── Slide layouts available ──
    try:
        layouts = []
        for layout in prs.slide_layouts:
            layouts.append(layout.name)
        if layouts:
            parts.append(f"[Slide Layouts] {', '.join(layouts)}")
    except Exception:
        pass

    # ── Per-slide details ──
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        # Layout name
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass
        slide_parts.append(f"[Slide {i+1}] layout: {layout_name or 'unknown'}")

        # Background
        try:
            bg = slide.background
            if bg and bg.fill and bg.fill.type is not None:
                fill_type = str(bg.fill.type)
                fill_info = f"background: {fill_type}"
                try:
                    fc = bg.fill.fore_color
                    if fc and fc.rgb:
                        fill_info += f" #{fc.rgb}"
                except Exception:
                    pass
                slide_parts.append(f"  {fill_info}")
        except Exception:
            pass

        # Shapes
        for shape in slide.shapes:
            shape_info = _describe_shape(shape)
            if shape_info:
                slide_parts.append(shape_info)

        parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)[:MAX_CHARS]


def _describe_shape(shape, indent="  "):
    """Describe a single shape with position, size, style, and content."""
    lines = []

    # Basic info
    shape_type = shape.shape_type
    type_name = str(shape_type).split('(')[0].strip() if shape_type else "shape"
    name = shape.name or ""

    # Position & size
    pos = ""
    try:
        l, t = _emu_to_cm(shape.left), _emu_to_cm(shape.top)
        w, h = _emu_to_cm(shape.width), _emu_to_cm(shape.height)
        pos = f"pos:({l},{t})cm size:({w}x{h})cm"
    except Exception:
        pass

    lines.append(f"{indent}[{type_name}] {name}  {pos}")

    # Image
    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        try:
            img = shape.image
            lines.append(f"{indent}  image: {img.content_type}, {len(img.blob)//1024}KB")
        except Exception:
            lines.append(f"{indent}  image: (embedded)")

    # Table
    if shape.has_table:
        try:
            table = shape.table
            rows_info = f"{len(table.rows)}rows x {len(table.columns)}cols"
            lines.append(f"{indent}  table: {rows_info}")
            if table.rows:
                header = [cell.text.strip() for cell in table.rows[0].cells]
                lines.append(f"{indent}  header: {' | '.join(header)}")
        except Exception:
            pass

    # Text content with style
    if shape.has_text_frame:
        tf = shape.text_frame
        for para in tf.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Paragraph-level style
            style_parts = []
            try:
                if para.alignment is not None:
                    style_parts.append(f"align:{para.alignment}")
            except Exception:
                pass
            # Run-level font (use first run as representative)
            if para.runs:
                fd = _font_desc(para.runs[0].font)
                if fd:
                    style_parts.append(fd)
            style_str = f"  ({', '.join(style_parts)})" if style_parts else ""
            lines.append(f"{indent}  \"{text}\"{style_str}")

    return "\n".join(lines) if lines else ""
