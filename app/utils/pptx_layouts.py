"""
PPTX Smart Layout Engine — auto-calculate element positions from high-level specs.

Agent says *what* content goes on a slide; this module figures out *where*.
"""
from __future__ import annotations

import math
from typing import Any

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
SLIDE_W = 10.0
SLIDE_H = 5.625
SAFE_LEFT = 0.5
SAFE_RIGHT = 9.5
SAFE_TOP = 0.3
SAFE_BOTTOM = 5.35
CONTENT_W = SAFE_RIGHT - SAFE_LEFT  # 9.0
CONTENT_H = SAFE_BOTTOM - SAFE_TOP  # 5.05

# Header bar occupies top 1 inch on most content slides
HEADER_H = 1.0
BODY_TOP = HEADER_H + 0.3  # content starts after header + gap


def _header_elements(title: str, page_num: int | str, theme: dict) -> list[dict]:
    """Standard page header: colored bar + title text."""
    num_str = str(page_num).zfill(3) if isinstance(page_num, int) else str(page_num)
    return [
        {"type": "shape", "shape_type": "rectangle",
         "x": 0, "y": 0, "w": SLIDE_W, "h": HEADER_H,
         "fill_color": theme.get("primary", "1E40AF")},
        {"type": "text",
         "content": f"{num_str} | {title}",
         "x": SAFE_LEFT, "y": 0.15, "w": CONTENT_W, "h": 0.7,
         "font_size": 24, "bold": True, "color": "FFFFFF",
         "font_name": theme.get("title_font", "Microsoft YaHei")},
    ]


def _footer_line(theme: dict) -> dict:
    """Subtle bottom decoration line."""
    return {"type": "line",
            "x": SAFE_LEFT, "y": SAFE_BOTTOM, "w": CONTENT_W, "h": 0,
            "line_color": theme.get("accent", "3B82F6"), "line_width": 0.5}


# ═══════════════════════════════════════════════════════════════════════════
# Layout: cover
# ═══════════════════════════════════════════════════════════════════════════
def layout_cover(spec: dict, theme: dict) -> list[dict]:
    """Full-bleed cover page with left accent bar."""
    title = spec.get("title", "标题")
    subtitle = spec.get("subtitle", "")
    date = spec.get("date", "")
    author = spec.get("author", "")

    primary = theme.get("primary", "1E40AF")
    bg = theme.get("background", "FFFFFF")
    accent = theme.get("accent", "3B82F6")

    els: list[dict] = [
        # Full background
        {"type": "shape", "shape_type": "rectangle",
         "x": 0, "y": 0, "w": SLIDE_W, "h": SLIDE_H,
         "fill_color": bg},
        # Left accent bar (30%)
        {"type": "shape", "shape_type": "rectangle",
         "x": 0, "y": 0, "w": 3.0, "h": SLIDE_H,
         "fill_color": primary},
        # Decorative circle on accent bar
        {"type": "icon_circle",
         "x": 1.0, "y": 3.8, "w": 1.0, "h": 1.0,
         "fill_color": accent, "text": "★", "font_size": 20, "font_color": "FFFFFF"},
        # Title
        {"type": "text", "content": title,
         "x": 3.5, "y": 1.5, "w": 5.5, "h": 1.4,
         "font_size": 36, "bold": True, "color": primary,
         "font_name": theme.get("title_font", "Microsoft YaHei")},
        # Divider line
        {"type": "line",
         "x": 3.5, "y": 3.1, "w": 4.0, "h": 0,
         "line_color": accent, "line_width": 2},
    ]

    if subtitle:
        els.append({
            "type": "text", "content": subtitle,
            "x": 3.5, "y": 3.3, "w": 5.5, "h": 0.8,
            "font_size": 18, "color": theme.get("secondary", "2B2B2B"),
            "font_name": theme.get("body_font", "Microsoft YaHei")})

    bottom_text = " | ".join(filter(None, [author, date]))
    if bottom_text:
        els.append({
            "type": "text", "content": bottom_text,
            "x": 3.5, "y": 4.5, "w": 5.5, "h": 0.5,
            "font_size": 12, "color": "999999",
            "font_name": theme.get("body_font", "Microsoft YaHei")})

    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: toc (目录)
# ═══════════════════════════════════════════════════════════════════════════
def layout_toc(spec: dict, theme: dict) -> list[dict]:
    """Table of contents — items auto-arranged in 2-column grid."""
    title = spec.get("title", "目录")
    items = spec.get("items", [])  # [{"num": "01", "text": "概述"}, ...]

    els = _header_elements(title, spec.get("page_num", 1), theme)
    primary = theme.get("primary", "1E40AF")
    accent = theme.get("accent", "3B82F6")

    n = len(items)
    cols = 2
    rows = math.ceil(n / cols)
    card_w = 4.2
    card_h = 0.6
    gap_x = 0.6
    gap_y = 0.2
    start_x = SAFE_LEFT
    start_y = BODY_TOP + 0.2

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        num = item.get("num", str(i + 1).zfill(2))
        text = item.get("text", "")

        # Card background
        els.append({"type": "shape", "shape_type": "rounded_rect",
                     "x": x, "y": y, "w": card_w, "h": card_h,
                     "fill_color": "F0F4FF", "line_color": "E0E7FF", "line_width": 1})
        # Number circle
        els.append({"type": "icon_circle",
                     "x": x + 0.1, "y": y + (card_h - 0.5) / 2,
                     "w": 0.5, "h": 0.5,
                     "fill_color": primary, "text": num,
                     "font_size": 12, "font_color": "FFFFFF"})
        # Text
        els.append({"type": "text", "content": text,
                     "x": x + 0.75, "y": y, "w": card_w - 0.9, "h": card_h,
                     "font_size": 14, "color": theme.get("secondary", "2B2B2B"),
                     "valign": "middle"})

    # Vertical center divider
    if n > 1:
        mid_x = start_x + card_w + gap_x / 2
        els.append({"type": "line",
                     "x": mid_x, "y": BODY_TOP,
                     "w": 0, "h": rows * (card_h + gap_y),
                     "line_color": "E0E7FF", "line_width": 0.5})

    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: section (章节分隔)
# ═══════════════════════════════════════════════════════════════════════════
def layout_section(spec: dict, theme: dict) -> list[dict]:
    """Section divider — big number + title, centered."""
    title = spec.get("title", "")
    num = spec.get("num", "01")
    subtitle = spec.get("subtitle", "")
    primary = theme.get("primary", "1E40AF")

    return [
        {"type": "shape", "shape_type": "rectangle",
         "x": 0, "y": 0, "w": SLIDE_W, "h": SLIDE_H,
         "fill_color": primary},
        {"type": "icon_circle",
         "x": 4.2, "y": 1.2, "w": 1.6, "h": 1.6,
         "fill_color": theme.get("accent", "3B82F6"),
         "text": str(num), "font_size": 36, "font_color": "FFFFFF"},
        {"type": "text", "content": title,
         "x": 1.5, "y": 3.2, "w": 7.0, "h": 1.0,
         "font_size": 32, "bold": True, "color": "FFFFFF", "align": "center",
         "font_name": theme.get("title_font", "Microsoft YaHei")},
        {"type": "text", "content": subtitle,
         "x": 2.0, "y": 4.2, "w": 6.0, "h": 0.6,
         "font_size": 16, "color": "FFFFFFCC", "align": "center"},
    ]


# ═══════════════════════════════════════════════════════════════════════════
# Layout: cards (N 个卡片自动排列)
# ═══════════════════════════════════════════════════════════════════════════
def layout_cards(spec: dict, theme: dict) -> list[dict]:
    """Auto-arrange N cards in optimal grid.

    items: [{"title": "...", "detail": "...", "icon": "01"}, ...]
    """
    title = spec.get("title", "")
    items = spec.get("items", [])
    n = len(items)
    if n == 0:
        return _header_elements(title, spec.get("page_num", ""), theme)

    primary = theme.get("primary", "1E40AF")
    accent = theme.get("accent", "3B82F6")

    # Decide grid: 1→1x1, 2→1x2, 3→1x3, 4→2x2, 5-6→2x3, 7-9→3x3
    if n <= 2:
        cols, rows = n, 1
    elif n == 3:
        cols, rows = 3, 1
    elif n == 4:
        cols, rows = 2, 2
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols = 3
        rows = math.ceil(n / cols)

    gap = 0.2
    total_gap_x = gap * (cols - 1)
    total_gap_y = gap * (rows - 1)
    card_w = (CONTENT_W - total_gap_x) / cols
    avail_h = SAFE_BOTTOM - BODY_TOP - 0.1
    card_h = min((avail_h - total_gap_y) / rows, 2.2)

    els = _header_elements(title, spec.get("page_num", ""), theme)

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        x = SAFE_LEFT + col * (card_w + gap)
        y = BODY_TOP + row * (card_h + gap)

        icon_text = item.get("icon", str(i + 1).zfill(2))
        card_title = item.get("title", "")
        detail = item.get("detail", "")

        # Card background
        els.append({"type": "shape", "shape_type": "rounded_rect",
                     "x": x, "y": y, "w": card_w, "h": card_h,
                     "fill_color": "F0F4FF", "line_color": "E0E7FF", "line_width": 1})
        # Icon circle
        els.append({"type": "icon_circle",
                     "x": x + 0.15, "y": y + 0.15,
                     "w": 0.6, "h": 0.6,
                     "fill_color": primary, "text": icon_text,
                     "font_size": 14, "font_color": "FFFFFF"})
        # Card title
        els.append({"type": "text", "content": card_title,
                     "x": x + 0.9, "y": y + 0.15,
                     "w": card_w - 1.1, "h": 0.5,
                     "font_size": 16, "bold": True,
                     "color": theme.get("secondary", "2B2B2B")})
        # Detail text
        if detail:
            els.append({"type": "text", "content": detail,
                         "x": x + 0.2, "y": y + 0.8,
                         "w": card_w - 0.4, "h": card_h - 1.0,
                         "font_size": 12, "color": "666666",
                         "line_spacing": 1.4})

    els.append(_footer_line(theme))
    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: process (流程)
# ═══════════════════════════════════════════════════════════════════════════
def layout_process(spec: dict, theme: dict) -> list[dict]:
    """Horizontal process flow with N steps + arrows between them."""
    title = spec.get("title", "")
    steps = spec.get("items", [])  # [{"title": "...", "detail": "..."}, ...]
    n = len(steps)
    if n == 0:
        return _header_elements(title, spec.get("page_num", ""), theme)

    primary = theme.get("primary", "1E40AF")
    accent = theme.get("accent", "3B82F6")
    els = _header_elements(title, spec.get("page_num", ""), theme)

    arrow_w = 0.3
    total_arrows = max(n - 1, 0)
    total_arrow_space = total_arrows * (arrow_w + 0.1)  # arrows + small gaps
    card_w = (CONTENT_W - total_arrow_space - 0.1 * (n - 1)) / n
    card_w = min(card_w, 3.5)  # cap individual card width
    card_h = 2.2
    y = BODY_TOP + 0.2

    # Recalculate total width and center
    total_w = n * card_w + total_arrows * (arrow_w + 0.2)
    start_x = SAFE_LEFT + (CONTENT_W - total_w) / 2

    cursor_x = start_x
    for i, step in enumerate(steps):
        step_title = step.get("title", f"步骤 {i + 1}")
        detail = step.get("detail", "")
        icon = step.get("icon", str(i + 1).zfill(2))

        # Card background
        els.append({"type": "shape", "shape_type": "rounded_rect",
                     "x": cursor_x, "y": y, "w": card_w, "h": card_h,
                     "fill_color": "F0F4FF", "line_color": "E0E7FF", "line_width": 1})
        # Icon circle at top-center of card
        els.append({"type": "icon_circle",
                     "x": cursor_x + (card_w - 0.6) / 2, "y": y + 0.15,
                     "w": 0.6, "h": 0.6,
                     "fill_color": primary if i == 0 else accent,
                     "text": icon, "font_size": 14, "font_color": "FFFFFF"})
        # Step title
        els.append({"type": "text", "content": step_title,
                     "x": cursor_x + 0.1, "y": y + 0.85,
                     "w": card_w - 0.2, "h": 0.4,
                     "font_size": 14, "bold": True, "align": "center",
                     "color": theme.get("secondary", "2B2B2B")})
        # Detail
        if detail:
            els.append({"type": "text", "content": detail,
                         "x": cursor_x + 0.1, "y": y + 1.3,
                         "w": card_w - 0.2, "h": card_h - 1.5,
                         "font_size": 11, "align": "center", "color": "666666",
                         "line_spacing": 1.3})

        cursor_x += card_w

        # Arrow between steps
        if i < n - 1:
            arrow_x = cursor_x + 0.05
            arrow_y = y + card_h / 2 - 0.15
            els.append({"type": "shape", "shape_type": "arrow_right",
                         "x": arrow_x, "y": arrow_y, "w": arrow_w, "h": 0.3,
                         "fill_color": accent})
            cursor_x += arrow_w + 0.2

    # Bottom summary text area
    summary = spec.get("summary", "")
    if summary:
        els.append({"type": "text", "content": summary,
                     "x": SAFE_LEFT, "y": y + card_h + 0.3,
                     "w": CONTENT_W, "h": 0.5,
                     "font_size": 12, "color": "888888", "align": "center"})

    els.append(_footer_line(theme))
    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: kpi (数据看板)
# ═══════════════════════════════════════════════════════════════════════════
def layout_kpi(spec: dict, theme: dict) -> list[dict]:
    """KPI dashboard — big numbers with labels."""
    title = spec.get("title", "")
    kpis = spec.get("items", [])  # [{"value": "99%", "label": "准确率", "icon": "★"}, ...]
    n = len(kpis)
    if n == 0:
        return _header_elements(title, spec.get("page_num", ""), theme)

    primary = theme.get("primary", "1E40AF")
    accent = theme.get("accent", "3B82F6")
    els = _header_elements(title, spec.get("page_num", ""), theme)

    cols = min(n, 4)
    gap = 0.3
    card_w = (CONTENT_W - gap * (cols - 1)) / cols
    card_h = 2.5
    y = BODY_TOP + 0.5

    # Center if fewer items
    total_w = cols * card_w + (cols - 1) * gap
    start_x = SAFE_LEFT + (CONTENT_W - total_w) / 2

    for i, kpi in enumerate(kpis):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap)
        cy = y + row * (card_h + 0.3)

        value = kpi.get("value", "0")
        label = kpi.get("label", "")
        icon = kpi.get("icon", "◆")

        # Card background
        els.append({"type": "shape", "shape_type": "rounded_rect",
                     "x": x, "y": cy, "w": card_w, "h": card_h,
                     "fill_color": "F8FAFC", "line_color": "E2E8F0", "line_width": 1})
        # Icon circle
        els.append({"type": "icon_circle",
                     "x": x + (card_w - 0.5) / 2, "y": cy + 0.2,
                     "w": 0.5, "h": 0.5,
                     "fill_color": accent, "text": icon,
                     "font_size": 12, "font_color": "FFFFFF"})
        # Big number
        els.append({"type": "text", "content": str(value),
                     "x": x, "y": cy + 0.8, "w": card_w, "h": 0.9,
                     "font_size": 42, "bold": True, "align": "center",
                     "color": primary})
        # Label
        els.append({"type": "text", "content": label,
                     "x": x, "y": cy + 1.7, "w": card_w, "h": 0.5,
                     "font_size": 14, "align": "center", "color": "888888"})

    els.append(_footer_line(theme))
    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: comparison (左右对比)
# ═══════════════════════════════════════════════════════════════════════════
def layout_comparison(spec: dict, theme: dict) -> list[dict]:
    """Side-by-side comparison with left/right content."""
    title = spec.get("title", "")
    left = spec.get("left", {})   # {"title": "...", "items": ["...", "..."]}
    right = spec.get("right", {})
    primary = theme.get("primary", "1E40AF")
    accent = theme.get("accent", "3B82F6")

    els = _header_elements(title, spec.get("page_num", ""), theme)

    col_w = 4.2
    left_x = SAFE_LEFT
    right_x = SAFE_LEFT + col_w + 0.6
    y = BODY_TOP

    for side, x, color in [(left, left_x, primary), (right, right_x, accent)]:
        side_title = side.get("title", "")
        items = side.get("items", [])

        # Column header
        els.append({"type": "shape", "shape_type": "rounded_rect",
                     "x": x, "y": y, "w": col_w, "h": 0.6,
                     "fill_color": color})
        els.append({"type": "text", "content": side_title,
                     "x": x + 0.2, "y": y + 0.05, "w": col_w - 0.4, "h": 0.5,
                     "font_size": 16, "bold": True, "color": "FFFFFF",
                     "align": "center", "valign": "middle"})

        # Items
        for j, item_text in enumerate(items):
            iy = y + 0.8 + j * 0.45
            els.append({"type": "text",
                         "content": f"• {item_text}",
                         "x": x + 0.2, "y": iy, "w": col_w - 0.4, "h": 0.4,
                         "font_size": 13, "color": theme.get("secondary", "2B2B2B")})

    # VS divider
    mid_x = left_x + col_w + 0.05
    els.append({"type": "line",
                 "x": mid_x, "y": BODY_TOP, "w": 0, "h": 3.5,
                 "line_color": "DDDDDD", "line_width": 1})
    els.append({"type": "icon_circle",
                 "x": mid_x - 0.2, "y": BODY_TOP + 1.5,
                 "w": 0.5, "h": 0.5,
                 "fill_color": "FFFFFF", "text": "VS",
                 "font_size": 10, "font_color": "999999"})

    els.append(_footer_line(theme))
    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: timeline (时间轴)
# ═══════════════════════════════════════════════════════════════════════════
def layout_timeline(spec: dict, theme: dict) -> list[dict]:
    """Horizontal timeline with milestones."""
    title = spec.get("title", "")
    items = spec.get("items", [])  # [{"date": "2024-Q1", "text": "..."}, ...]
    n = len(items)
    if n == 0:
        return _header_elements(title, spec.get("page_num", ""), theme)

    primary = theme.get("primary", "1E40AF")
    accent = theme.get("accent", "3B82F6")
    els = _header_elements(title, spec.get("page_num", ""), theme)

    # Horizontal axis
    axis_y = BODY_TOP + 1.5
    els.append({"type": "line",
                 "x": SAFE_LEFT, "y": axis_y, "w": CONTENT_W, "h": 0,
                 "line_color": primary, "line_width": 2})

    gap = CONTENT_W / max(n, 1)
    for i, item in enumerate(items):
        cx = SAFE_LEFT + gap * i + gap / 2
        date = item.get("date", "")
        text = item.get("text", "")

        # Node circle on axis
        els.append({"type": "icon_circle",
                     "x": cx - 0.25, "y": axis_y - 0.25,
                     "w": 0.5, "h": 0.5,
                     "fill_color": primary if i % 2 == 0 else accent,
                     "text": str(i + 1), "font_size": 12, "font_color": "FFFFFF"})

        # Alternate above/below the axis
        if i % 2 == 0:
            # Above
            els.append({"type": "text", "content": date,
                         "x": cx - 0.8, "y": axis_y - 1.0,
                         "w": 1.6, "h": 0.35,
                         "font_size": 11, "bold": True, "align": "center",
                         "color": primary})
            els.append({"type": "text", "content": text,
                         "x": cx - 0.8, "y": axis_y - 0.65,
                         "w": 1.6, "h": 0.35,
                         "font_size": 10, "align": "center", "color": "666666"})
        else:
            # Below
            els.append({"type": "text", "content": date,
                         "x": cx - 0.8, "y": axis_y + 0.35,
                         "w": 1.6, "h": 0.35,
                         "font_size": 11, "bold": True, "align": "center",
                         "color": primary})
            els.append({"type": "text", "content": text,
                         "x": cx - 0.8, "y": axis_y + 0.7,
                         "w": 1.6, "h": 0.35,
                         "font_size": 10, "align": "center", "color": "666666"})

    els.append(_footer_line(theme))
    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: chart_page (图表页 — 标题+描述框架，chart 由 agent 在 elements 补充)
# ═══════════════════════════════════════════════════════════════════════════
def layout_chart_page(spec: dict, theme: dict) -> list[dict]:
    """Chart page scaffold — header + description. Agent adds chart element manually."""
    title = spec.get("title", "")
    description = spec.get("description", "")

    els = _header_elements(title, spec.get("page_num", ""), theme)

    if description:
        els.append({"type": "text", "content": description,
                     "x": SAFE_LEFT, "y": SAFE_BOTTOM - 0.5,
                     "w": CONTENT_W, "h": 0.4,
                     "font_size": 12, "color": "888888", "align": "center"})

    els.append(_footer_line(theme))
    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: closing (结束页)
# ═══════════════════════════════════════════════════════════════════════════
def layout_closing(spec: dict, theme: dict) -> list[dict]:
    """Closing/thank-you page."""
    title = spec.get("title", "Thank You")
    subtitle = spec.get("subtitle", "")
    contact = spec.get("contact", "")
    primary = theme.get("primary", "1E40AF")

    els: list[dict] = [
        {"type": "shape", "shape_type": "rectangle",
         "x": 0, "y": 0, "w": SLIDE_W, "h": SLIDE_H,
         "fill_color": theme.get("background", "FFFFFF")},
        {"type": "shape", "shape_type": "rectangle",
         "x": 0, "y": 0, "w": 3.0, "h": SLIDE_H,
         "fill_color": primary},
        {"type": "icon_circle",
         "x": 1.0, "y": 2.0, "w": 1.0, "h": 1.0,
         "fill_color": theme.get("accent", "3B82F6"),
         "text": "✓", "font_size": 24, "font_color": "FFFFFF"},
        {"type": "text", "content": title,
         "x": 3.5, "y": 1.8, "w": 5.5, "h": 1.2,
         "font_size": 32, "bold": True, "color": primary,
         "font_name": theme.get("title_font", "Microsoft YaHei")},
        {"type": "line",
         "x": 3.5, "y": 3.2, "w": 4.0, "h": 0,
         "line_color": theme.get("accent", "3B82F6"), "line_width": 2},
    ]

    if subtitle:
        els.append({"type": "text", "content": subtitle,
                     "x": 3.5, "y": 3.4, "w": 5.5, "h": 0.6,
                     "font_size": 16, "color": theme.get("secondary", "2B2B2B")})

    if contact:
        els.append({"type": "text", "content": contact,
                     "x": 3.5, "y": 4.2, "w": 5.5, "h": 0.6,
                     "font_size": 12, "color": "999999"})

    return els


# ═══════════════════════════════════════════════════════════════════════════
# Layout: table_page (表格页)
# ═══════════════════════════════════════════════════════════════════════════
def layout_table_page(spec: dict, theme: dict) -> list[dict]:
    """Table page with header bar. The table element itself is generated here."""
    title = spec.get("title", "")
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    summary = spec.get("summary", "")

    els = _header_elements(title, spec.get("page_num", ""), theme)

    if headers or rows:
        row_count = len(rows) + (1 if headers else 0)
        table_h = min(0.4 * row_count + 0.2, SAFE_BOTTOM - BODY_TOP - 0.8)
        els.append({
            "type": "table",
            "x": SAFE_LEFT, "y": BODY_TOP, "w": CONTENT_W, "h": table_h,
            "headers": headers,
            "rows": rows,
            "header_color": theme.get("primary", "1E40AF"),
            "header_font_color": "FFFFFF",
            "stripe_color": "F8FAFC",
        })

    if summary:
        els.append({"type": "text", "content": summary,
                     "x": SAFE_LEFT, "y": SAFE_BOTTOM - 0.5,
                     "w": CONTENT_W, "h": 0.4,
                     "font_size": 12, "color": "888888"})

    els.append(_footer_line(theme))
    return els


# ═══════════════════════════════════════════════════════════════════════════
# Registry
# ═══════════════════════════════════════════════════════════════════════════
LAYOUT_REGISTRY: dict[str, Any] = {
    "cover": layout_cover,
    "toc": layout_toc,
    "section": layout_section,
    "cards": layout_cards,
    "grid": layout_cards,        # alias
    "grid_2x2": layout_cards,    # alias
    "grid_2x3": layout_cards,    # alias
    "two_column": layout_cards,  # 2 items → auto 1x2
    "three_column": layout_cards,  # 3 items → auto 1x3
    "process": layout_process,
    "kpi": layout_kpi,
    "comparison": layout_comparison,
    "timeline": layout_timeline,
    "chart": layout_chart_page,
    "chart_page": layout_chart_page,
    "table": layout_table_page,
    "table_page": layout_table_page,
    "closing": layout_closing,
}


def generate_layout(layout_spec: dict, theme: dict) -> list[dict]:
    """Public entry point. Returns a list of element dicts.

    layout_spec:
      {"type": "process", "title": "实施流程", "page_num": 3,
       "items": [{"title": "调研", "detail": "..."}, ...]}
    """
    layout_type = layout_spec.get("type", "")
    fn = LAYOUT_REGISTRY.get(layout_type)
    if fn is None:
        return []
    return fn(layout_spec, theme)


# ═══════════════════════════════════════════════════════════════════════════
# Part 2: Template extraction from user-uploaded PPTX
# ═══════════════════════════════════════════════════════════════════════════

def extract_pptx_template(file_path: str) -> dict:
    """Analyze an uploaded PPTX and extract reusable style info.

    Returns a template dict:
    {
        "slide_size": {"w": 10.0, "h": 5.625},
        "colors": {"primary": "1E40AF", ...},
        "fonts": {"title": "...", "body": "..."},
        "slides": [
            {
                "index": 0,
                "layout_name": "...",
                "elements": [
                    {"type": "shape|text|...", "x": ..., "y": ..., "w": ..., "h": ...,
                     "font_size": ..., "color": ..., "text_preview": "..."}
                ]
            }
        ]
    }
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(file_path)
    EMU_TO_IN = 1 / 914400

    # Collect colors & fonts across all slides
    all_colors: dict[str, int] = {}  # hex → frequency
    all_fonts: dict[str, int] = {}   # font name → frequency

    slides_info = []

    for si, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
        elements = []

        for shape in slide.shapes:
            el: dict[str, Any] = {
                "name": shape.name,
                "x": round(shape.left * EMU_TO_IN, 2) if shape.left else 0,
                "y": round(shape.top * EMU_TO_IN, 2) if shape.top else 0,
                "w": round(shape.width * EMU_TO_IN, 2) if shape.width else 0,
                "h": round(shape.height * EMU_TO_IN, 2) if shape.height else 0,
            }

            # Detect type
            st = str(shape.shape_type) if shape.shape_type else ""
            if "TABLE" in st:
                el["type"] = "table"
                try:
                    tbl = shape.table
                    el["rows"] = tbl.rows.__len__()
                    el["cols"] = tbl.columns.__len__()
                except Exception:
                    pass
            elif "CHART" in st:
                el["type"] = "chart"
            elif "PICTURE" in st or "MEDIA" in st:
                el["type"] = "image"
            elif "TEXT_BOX" in st:
                el["type"] = "text"
            else:
                el["type"] = "shape"

            # Extract text info
            if shape.has_text_frame:
                text = shape.text_frame.text[:100]
                el["text_preview"] = text
                for para in shape.text_frame.paragraphs:
                    if para.font and para.font.name:
                        fname = para.font.name
                        all_fonts[fname] = all_fonts.get(fname, 0) + 1
                    if para.font and para.font.size:
                        el["font_size"] = round(para.font.size / 12700, 1)  # EMU → pt
                    try:
                        if para.font and para.font.color and para.font.color.type is not None:
                            c = str(para.font.color.rgb)
                            el["color"] = c
                            all_colors[c] = all_colors.get(c, 0) + 1
                    except (AttributeError, TypeError):
                        pass

            # Extract fill color
            try:
                if hasattr(shape, "fill") and shape.fill.type is not None:
                    fc = shape.fill.fore_color
                    if fc and fc.rgb:
                        c = str(fc.rgb)
                        el["fill_color"] = c
                        all_colors[c] = all_colors.get(c, 0) + 3  # fill colors are more important
            except Exception:
                pass

            elements.append(el)

        slides_info.append({
            "index": si,
            "layout_name": layout_name,
            "element_count": len(elements),
            "elements": elements,
        })

    # Derive theme from frequency analysis
    sorted_colors = sorted(all_colors.items(), key=lambda x: -x[1])
    # Filter out white/black/near-white/near-black
    significant_colors = [
        c for c, _ in sorted_colors
        if c not in ("FFFFFF", "000000", "ffffff", "000000")
        and not c.startswith("F") and not c.startswith("0")
    ]
    if not significant_colors and sorted_colors:
        significant_colors = [c for c, _ in sorted_colors[:3]]

    sorted_fonts = sorted(all_fonts.items(), key=lambda x: -x[1])

    template = {
        "slide_size": {
            "w": round(prs.slide_width * EMU_TO_IN, 2),
            "h": round(prs.slide_height * EMU_TO_IN, 2),
        },
        "colors": {
            "detected": [c for c, _ in sorted_colors[:10]],
            "primary": significant_colors[0] if len(significant_colors) > 0 else "1E40AF",
            "secondary": significant_colors[1] if len(significant_colors) > 1 else "2B2B2B",
            "accent": significant_colors[2] if len(significant_colors) > 2 else "3B82F6",
        },
        "fonts": {
            "detected": [f for f, _ in sorted_fonts[:5]],
            "title": sorted_fonts[0][0] if sorted_fonts else "Microsoft YaHei",
            "body": sorted_fonts[1][0] if len(sorted_fonts) > 1 else sorted_fonts[0][0] if sorted_fonts else "Microsoft YaHei",
        },
        "slide_count": len(slides_info),
        "slides": slides_info,
    }

    return template
